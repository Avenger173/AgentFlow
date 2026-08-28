"""把已核验的文档草稿交付为可编辑项目方案 PPTX。

这里刻意不调用模型：草稿生成与来源核验已经在 Document Agent Runtime 完成，演示文稿层只把
同一份受控快照转换成可确认的页面计划和确定性文件。这样“导出 PPT”不会悄悄引入新事实、
重新消耗模型预算，或因模型波动让同一草稿每次导出不同内容。
"""

from __future__ import annotations

import hashlib
import json
import re
from datetime import UTC, datetime
from pathlib import Path
from uuid import uuid4

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from app.core.config import settings
from app.database.task_repository import append_workflow_artifact
from app.schemas.document_agent import DocumentAgentRunResponse, DocumentSourceRef
from app.schemas.presentation import (
    PresentationExportRequest,
    PresentationExportResponse,
    PresentationPreflight,
    PresentationPreflightFinding,
    PresentationPreviewResponse,
    PresentationSlidePlan,
    PresentationVerification,
)
from app.schemas.workflow import WorkflowArtifact
from app.services.document_agent import get_document_agent_result
from app.services.project_review import (
    ProjectReviewServiceError,
    evaluate_project_document_material,
)
from app.services.workspace_documents import WorkspaceDocumentError, read_workspace_document_chunks
from app.workflow.dry_run import clear_dry_run_memory_cache


_MAX_CONTENT_SLIDES = 8
_PRESENTATION_TYPE = "project_proposal"
_PPTX_MIME_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_DOCUMENT_AGENT_ID = "document_agent"

# 固定模板的色彩与现有客户端的浅色工作台保持一致。这里不做自由主题系统，避免 V1 在没有
# 模板设计器、字体预检和客户品牌资产的情况下承诺“任意风格 PPT”。
_COLOR_NAVY = RGBColor(20, 34, 74)
_COLOR_BLUE = RGBColor(35, 100, 232)
_COLOR_SKY = RGBColor(236, 245, 255)
_COLOR_TEXT = RGBColor(31, 48, 84)
_COLOR_MUTED = RGBColor(93, 111, 140)
_COLOR_WHITE = RGBColor(255, 255, 255)
_COLOR_PALE = RGBColor(247, 250, 255)


class PresentationDeliveryError(RuntimeError):
    """演示文稿工作流可预期的业务错误。"""


class PresentationNotFoundError(PresentationDeliveryError):
    """关联不到一个可交付的文档草稿快照。"""


class PresentationConflictError(PresentationDeliveryError):
    """计划过期、命名冲突或文件验证失败。"""


class PresentationConfirmationError(PresentationDeliveryError):
    """客户端未明确确认本次写入。"""


def build_project_proposal_preview(*, task_id: str) -> PresentationPreviewResponse:
    """根据已完成、已核验的文档草稿构建只读幻灯片计划。"""

    result = _load_exportable_draft(task_id)
    context = result.document_context
    warnings: list[str] = []
    sections = context.draft_sections
    if len(sections) > _MAX_CONTENT_SLIDES:
        # 当前 DocumentContext 的正式协议最多八章，这层显式校验是未来扩展时的安全护栏。
        raise PresentationDeliveryError(
            f"当前草稿包含 {len(sections)} 个章节，超过项目方案 PPT V1 的 {_MAX_CONTENT_SLIDES} 章上限；"
            "请先在草稿中合并章节后再生成演示文稿。"
        )
    # 自动预检与草稿来源范围绑定：它只重读实际支撑本次草稿的材料，不把 workspace 里不相关
    # 的文件混入项目判断。预检只给计划附加质量事实，文件写入仍由用户一次确认控制。
    preflight = _build_project_delivery_preflight(context)

    source_version_id = context.draft_version.version_id if context.draft_version is not None else task_id
    title = _normalize_text(context.draft_title, limit=120)
    source_refs = _unique_source_refs(
        source for section in sections for source in section.source_refs
    )
    cover_sources = source_refs[:4]
    agenda_items = [section.heading for section in sections]
    slides: list[PresentationSlidePlan] = [
        PresentationSlidePlan(
            slide_id="cover",
            role="cover",
            title=title,
            bullets=["基于已核验的文档草稿生成", "导出前请确认每页内容与来源范围"],
            source_refs=cover_sources,
        ),
        PresentationSlidePlan(
            slide_id="agenda",
            role="agenda",
            title="方案目录",
            bullets=agenda_items,
        ),
    ]
    for index, section in enumerate(sections, start=1):
        bullets = _extract_section_bullets(section.body)
        if not bullets:
            # DocumentDraftSection.body 受 schema 限制必非空；保留显式警告可避免渲染出“标题有、内容空”的
            # 误导幻灯片。这里用原章节短文本，而不是让导出层推断或补写内容。
            bullets = [_normalize_text(section.body, limit=150)]
            warnings.append(f"“{section.heading}”未识别到分段要点，已按原文短段落排版。")
        slides.append(
            PresentationSlidePlan(
                slide_id=f"content_{index}",
                role="content",
                title=_normalize_text(section.heading, limit=120),
                bullets=bullets,
                source_refs=section.source_refs,
            )
        )

    summary_bullets = _build_summary_bullets(sections)
    slides.append(
        PresentationSlidePlan(
            slide_id="summary",
            role="summary",
            title="交付前核对",
            bullets=summary_bullets,
            source_refs=source_refs[:4],
        )
    )
    source_labels = [_source_label(source) for source in source_refs]
    if len(source_labels) > 12:
        source_labels = source_labels[:12]
        warnings.append("来源清单超过 12 条，演示文稿仅展示前 12 条；完整来源仍保留在原文档任务中。")
    slides.append(
        PresentationSlidePlan(
            slide_id="sources",
            role="sources",
            title="来源与追溯",
            bullets=source_labels or ["原任务未提供可展示来源，当前草稿不应进入导出。"],
            source_refs=source_refs[:4],
        )
    )

    plan_id = _plan_id(
        task_id=task_id,
        source_version_id=source_version_id,
        title=title,
        slides=slides,
        preflight=preflight,
    )
    return PresentationPreviewResponse(
        source_task_id=task_id,
        source_version_id=source_version_id,
        presentation_type=_PRESENTATION_TYPE,
        plan_id=plan_id,
        title=title,
        slides=slides,
        preflight=preflight,
        warnings=warnings,
    )


def export_project_proposal_presentation(
    *,
    task_id: str,
    request: PresentationExportRequest,
) -> PresentationExportResponse:
    """确认后渲染 PPTX、回读验证，并把交付物追加到原任务的审计链。"""

    if not request.confirmed:
        raise PresentationConfirmationError("导出项目方案 PPT 前需要用户确认。")
    preview = build_project_proposal_preview(task_id=task_id)
    if request.presentation_type != preview.presentation_type:
        raise PresentationDeliveryError("当前仅支持“项目方案”演示文稿类型。")
    if request.plan_id != preview.plan_id:
        # 计划 hash 绑定任务、版本、章节正文和来源。草稿变化后必须重新打开预览，不能把旧确认
        # 静默套到新内容上。
        raise PresentationConflictError("演示文稿计划已过期，请重新查看预览并确认导出。")

    filename = _safe_presentation_filename(
        request.filename,
        fallback_title=preview.title,
        task_id=task_id,
    )
    output_root = settings.document_presentation_output_dir
    output_root.mkdir(parents=True, exist_ok=True)
    target = (output_root / filename).resolve()
    try:
        target.relative_to(output_root)
    except ValueError as exc:
        raise PresentationDeliveryError("演示文稿文件名未通过受控输出目录校验。") from exc

    try:
        # 先用 x 模式原子占位，再把同一个文件句柄交给渲染器。这样任何同名文件都不会被覆盖。
        with target.open("xb") as target_file:
            _render_project_proposal_presentation(target_file, preview)
    except FileExistsError as exc:
        raise PresentationConflictError(
            f"output/document_presentations 中已存在同名文件“{filename}”，请改名后再次确认导出。"
        ) from exc
    except Exception as exc:
        target.unlink(missing_ok=True)
        raise PresentationDeliveryError(f"无法生成项目方案 PPT：{exc}") from exc

    try:
        verification = _verify_presentation(target, preview)
        if not verification.passed:
            raise PresentationConflictError("项目方案 PPT 回读验证失败。")
    except Exception as exc:
        target.unlink(missing_ok=True)
        if isinstance(exc, PresentationDeliveryError):
            raise
        raise PresentationDeliveryError(f"项目方案 PPT 回读验证失败：{exc}") from exc

    artifact_id = f"{task_id}:project_proposal_ppt:{uuid4().hex[:10]}"
    relative_path = f"output/document_presentations/{filename}"
    artifact = WorkflowArtifact(
        artifact_id=artifact_id,
        task_id=task_id,
        step_id="presentation_export",
        agent_id=_DOCUMENT_AGENT_ID,
        kind="file",
        name=filename,
        summary=(
            f"用户确认导出的项目方案演示文稿，包含 {verification.slide_count} 页，"
            "内容来自已核验文档草稿。"
        ),
        uri=f"agentflow-output://document_presentations/{filename}",
        mime_type=_PPTX_MIME_TYPE,
        metadata={
            "runtime": True,
            "output_scope": "document_presentations",
            "output_path": str(target),
            "relative_output_path": relative_path,
            "confirmed_by": "local_user",
            "source_task_id": task_id,
            "source_version_id": preview.source_version_id,
            "presentation_type": preview.presentation_type,
            "plan_id": preview.plan_id,
            "slide_count": verification.slide_count,
            "source_slide_count": verification.source_slide_count,
            "verification_passed": verification.passed,
        },
        created_at=datetime.now(UTC).isoformat(timespec="seconds"),
    )
    try:
        append_workflow_artifact(
            artifact=artifact,
            event_name="artifact_saved",
            message=f"用户已确认导出项目方案演示文稿：{relative_path}",
        )
    except Exception as exc:
        # 审计失败时不能留下一个历史页不可见的“幽灵交付物”。文件是本调用 x 模式新建的，
        # 因此只能撤回它自身，不触碰用户此前的任何文件。
        target.unlink(missing_ok=True)
        raise PresentationDeliveryError("演示文稿导出审计失败，已撤回本次新建文件。") from exc

    clear_dry_run_memory_cache()
    return PresentationExportResponse(
        task_id=task_id,
        artifact_id=artifact_id,
        filename=filename,
        relative_path=relative_path,
        artifact_uri=artifact.uri,
        slide_count=verification.slide_count,
        verification=verification,
        message="项目方案 PPT 已导出，可在任务历史中查看交付记录。",
    )


def _load_exportable_draft(task_id: str) -> DocumentAgentRunResponse:
    """恢复并验证一个可用于交付的草稿快照。"""

    result = get_document_agent_result(task_id)
    if result is None:
        raise PresentationNotFoundError("未找到对应的文档草稿任务。")
    if result.status != "completed":
        raise PresentationDeliveryError("只有已完成的文档草稿可以生成项目方案 PPT。")
    context = result.document_context
    if not context.draft_title or not context.draft_sections:
        raise PresentationDeliveryError("当前任务不是可交付的文档草稿预览。")
    if context.draft_verification_state != "verified":
        raise PresentationDeliveryError("当前草稿尚未完成来源核验，请先处理待确认事实后再导出。")
    if not any(section.source_refs for section in context.draft_sections):
        raise PresentationDeliveryError("当前草稿缺少可追溯来源，不能生成项目方案 PPT。")
    return result


def _plan_id(
    *,
    task_id: str,
    source_version_id: str,
    title: str,
    slides: list[PresentationSlidePlan],
    preflight: PresentationPreflight,
) -> str:
    """生成绑定草稿快照的稳定计划身份，供确认导出时做过期保护。"""

    payload = {
        "template": _PRESENTATION_TYPE,
        "task_id": task_id,
        "source_version_id": source_version_id,
        "title": title,
        "slides": [slide.model_dump(mode="json") for slide in slides],
        # 预检读取的是当前受控材料，纳入 hash 后，材料更新或规则升级都会要求用户重新查看计划。
        "preflight": preflight.model_dump(mode="json"),
    }
    encoded = json.dumps(payload, ensure_ascii=False, sort_keys=True, separators=(",", ":"))
    return hashlib.sha256(encoded.encode("utf-8")).hexdigest()[:48]


def _build_project_delivery_preflight(context: object) -> PresentationPreflight:
    """自动运行项目材料规则检查，并返回供 PPT 预览展示的轻量结果。

    这里不调用 ``run_project_document_review``，因为那个入口的职责是“客户主动要一份审查报告”
    并会写入任务历史。PPT 预检只复用相同的纯规则函数，避免每次打开预览都制造额外历史任务。
    """

    source_paths: list[str] = []
    unsourced_sections: list[str] = []
    for section in getattr(context, "draft_sections", []):
        section_sources = list(getattr(section, "source_refs", []) or [])
        if not section_sources:
            unsourced_sections.append(str(getattr(section, "heading", "未命名章节")))
            continue
        for source in section_sources:
            relative_path = str(getattr(source, "relative_path", "")).strip()
            if relative_path and relative_path not in source_paths:
                source_paths.append(relative_path)

    if unsourced_sections:
        labels = "、".join(f"“{item}”" for item in unsourced_sections[:3])
        raise PresentationDeliveryError(f"{labels} 缺少可追溯来源，不能进入项目方案 PPT 导出。")
    if not source_paths:
        raise PresentationDeliveryError("当前草稿缺少可用于自动预检的项目材料来源。")
    if len(source_paths) > 4:
        # Document Agent 的草稿第一版不应跨太多材料交付。这个上限既控制预检耗时，也避免把
        # 互不相关的项目资料混作一份 PPT；用户可先整理为更聚焦的草稿。
        raise PresentationDeliveryError("当前草稿引用超过 4 份材料，请先收敛来源范围后再制作项目方案 PPT。")

    reports = []
    for relative_path in source_paths:
        try:
            chunks = read_workspace_document_chunks(relative_path=relative_path)
            reports.append(
                evaluate_project_document_material(
                    requested_document_type="auto",
                    chunks=chunks,
                )
            )
        except (WorkspaceDocumentError, ProjectReviewServiceError) as exc:
            raise PresentationDeliveryError(f"自动交付预检无法读取“{relative_path}”：{exc}") from exc

    check_total = sum(len(report.checks) for report in reports)
    passed_check_total = sum(
        1 for report in reports for check in report.checks if check.status == "passed"
    )
    attention_check_total = check_total - passed_check_total
    findings = sorted(
        (finding for report in reports for finding in report.findings),
        key=lambda item: ({"high": 0, "medium": 1, "low": 2}[item.severity], item.rule_id),
    )[:12]
    high_attention_total = sum(item.severity == "high" for item in findings)
    preflight_findings = [
        PresentationPreflightFinding(
            severity=finding.severity,
            category=finding.category,
            title=finding.title,
            suggestion=finding.suggestion,
            source_refs=finding.source_refs,
        )
        for finding in findings
    ]
    if attention_check_total:
        summary = (
            f"系统已自动核验 {len(reports)} 份项目材料，共 {check_total} 项规则："
            f"{passed_check_total} 项找到明确表述，{attention_check_total} 项建议在交付前补充。"
        )
    else:
        summary = (
            f"系统已自动核验 {len(reports)} 份项目材料，共 {check_total} 项规则均找到明确表述；"
            "仍建议由项目负责人确认业务事实。"
        )
    return PresentationPreflight(
        status="attention" if attention_check_total else "passed",
        summary=summary,
        checked_documents=source_paths,
        check_total=check_total,
        passed_check_total=passed_check_total,
        attention_check_total=attention_check_total,
        high_attention_total=high_attention_total,
        findings=preflight_findings,
        warnings=[
            "自动预检检查材料表述与来源范围，不替代项目负责人对实际执行、合规和行业要求的确认。"
        ],
    )


def _unique_source_refs(sources: object) -> list[DocumentSourceRef]:
    """按文件与定位去重来源，保留第一次出现的稳定顺序。"""

    unique: list[DocumentSourceRef] = []
    seen: set[tuple[str, int, int, str]] = set()
    for source in sources:
        if not isinstance(source, DocumentSourceRef):
            continue
        key = (source.relative_path, source.start_line, source.end_line, source.source_locator)
        if key in seen:
            continue
        seen.add(key)
        unique.append(source)
    return unique


def _build_summary_bullets(sections: list[object]) -> list[str]:
    """从每章第一个已有要点组成核对页，不新增行动项或结论。"""

    bullets: list[str] = []
    for section in sections:
        body = getattr(section, "body", "")
        heading = getattr(section, "heading", "")
        points = _extract_section_bullets(str(body))
        if points:
            bullets.append(_normalize_text(f"{heading}：{points[0]}", limit=150))
    return bullets[:6] or ["请在导出前回看各章节内容与来源。"]


def _extract_section_bullets(body: str) -> list[str]:
    """把已核验 Markdown 草稿压成适合单页阅读的短要点。"""

    candidates: list[str] = []
    for raw_line in body.splitlines():
        line = re.sub(r"^\s{0,3}(?:[-*+]\s+|\d+[.)]\s+|#{1,6}\s*)", "", raw_line).strip()
        line = re.sub(r"\*{1,3}([^*]+)\*{1,3}", r"\1", line)
        if not line or line.startswith(">"):
            continue
        # 一行过长时仅按可见句号/分号切分，绝不做语义重写。
        candidates.extend(part.strip() for part in re.split(r"(?<=[。！？!?；;])\s*", line) if part.strip())

    bullets: list[str] = []
    seen: set[str] = set()
    for candidate in candidates:
        clean = _normalize_text(candidate, limit=150)
        key = clean.casefold()
        if len(clean) < 4 or key in seen:
            continue
        seen.add(key)
        bullets.append(clean)
        if len(bullets) == 5:
            break
    return bullets


def _normalize_text(value: str, *, limit: int) -> str:
    """收紧空白并截断到模板可承载范围，避免写入不可读的溢出页面。"""

    normalized = re.sub(r"\s+", " ", value).strip()
    if len(normalized) <= limit:
        return normalized
    return f"{normalized[: max(1, limit - 1)].rstrip()}…"


def _source_label(source: DocumentSourceRef) -> str:
    """生成不含绝对路径的来源标签，供 PPT 页脚和来源页复核。"""

    locator = source.source_locator.replace("\n", " ").strip()
    if not locator:
        if source.source_kind == "page":
            locator = f"第 {source.start_line} 页"
        elif source.start_line == source.end_line:
            locator = f"第 {source.start_line} 行"
        else:
            locator = f"第 {source.start_line}-{source.end_line} 行"
    return _normalize_text(f"{source.relative_path} · {locator}", limit=130)


def _safe_presentation_filename(raw_filename: str, *, fallback_title: str, task_id: str) -> str:
    """将客户命名限制为单个 .pptx 文件名，禁止目录与覆盖。"""

    candidate = raw_filename.strip()
    if not candidate:
        candidate = f"{fallback_title or 'AgentFlow 项目方案'}-{task_id[-6:]}.pptx"
    if "/" in candidate or "\\" in candidate or Path(candidate).name != candidate:
        raise PresentationDeliveryError("演示文稿名称只能是文件名，不能包含目录或路径分隔符。")
    if not candidate.lower().endswith(".pptx"):
        raise PresentationDeliveryError("演示文稿名称必须以 .pptx 结尾。")
    stem = candidate[:-5].strip()
    sanitized_stem = re.sub(r'[<>:"|?*\x00-\x1f]+', "-", stem)
    sanitized_stem = re.sub(r"\s+", " ", sanitized_stem).strip(" .-")
    if not sanitized_stem:
        sanitized_stem = f"AgentFlow 项目方案-{task_id[-6:]}"
    return f"{sanitized_stem[:96]}.pptx"


def _render_project_proposal_presentation(target_file: object, preview: PresentationPreviewResponse) -> None:
    """按固定浅色项目方案模板渲染全部幻灯片。"""

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    for index, slide_plan in enumerate(preview.slides, start=1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        _paint_background(slide)
        if slide_plan.role == "cover":
            _render_cover_slide(slide, slide_plan, preview.title)
        else:
            _render_standard_slide(slide, slide_plan, index=index, total=len(preview.slides))
    presentation.save(target_file)


def _paint_background(slide: object) -> None:
    """创建固定边框和浅色画布，避免依赖客户端样式或本机主题。"""

    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = _COLOR_PALE
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.16))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = _COLOR_BLUE
    top_bar.line.fill.background()


def _render_cover_slide(slide: object, plan: PresentationSlidePlan, title: str) -> None:
    """渲染封面；只表现导出身份和来源状态，不包装不存在的客户品牌信息。"""

    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.15), Inches(11.9), Inches(4.9))
    panel.fill.solid()
    panel.fill.fore_color.rgb = _COLOR_SKY
    panel.line.color.rgb = RGBColor(199, 220, 252)
    _add_text_box(slide, Inches(1.15), Inches(1.7), Inches(10.9), Inches(1.35), title, font_size=30, color=_COLOR_NAVY, bold=True)
    _add_text_box(
        slide,
        Inches(1.18),
        Inches(3.1),
        Inches(9.8),
        Inches(0.7),
        "项目方案演示文稿 · 基于已核验文档草稿",
        font_size=16,
        color=_COLOR_BLUE,
    )
    _add_text_box(slide, Inches(1.18), Inches(4.25), Inches(9.8), Inches(0.5), "导出前请确认内容范围与来源。", font_size=13, color=_COLOR_MUTED)
    _add_source_footer(slide, plan.source_refs, page_label="AgentFlow 文档助手")


def _render_standard_slide(slide: object, plan: PresentationSlidePlan, *, index: int, total: int) -> None:
    """渲染目录、正文、核对和来源页的统一结构。"""

    _add_text_box(slide, Inches(0.78), Inches(0.62), Inches(11.1), Inches(0.62), plan.title, font_size=24, color=_COLOR_NAVY, bold=True)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.78), Inches(1.35), Inches(0.75), Inches(0.07))
    accent.fill.solid()
    accent.fill.fore_color.rgb = _COLOR_BLUE
    accent.line.fill.background()
    font_size = 16 if plan.role != "sources" else 13
    _add_bullet_box(
        slide,
        Inches(1.0),
        Inches(1.72),
        Inches(11.15),
        Inches(4.65),
        plan.bullets,
        font_size=font_size,
        compact=plan.role == "sources",
    )
    _add_source_footer(slide, plan.source_refs, page_label=f"{index} / {total}")


def _add_text_box(
    slide: object,
    left: object,
    top: object,
    width: object,
    height: object,
    text: str,
    *,
    font_size: int,
    color: RGBColor,
    bold: bool = False,
) -> None:
    """集中处理字体、边距和自动换行，保证模板中的文字不会因默认样式漂移。"""

    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    paragraph = text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.LEFT
    paragraph.font.name = "Microsoft YaHei UI"
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color


def _add_bullet_box(
    slide: object,
    left: object,
    top: object,
    width: object,
    height: object,
    bullets: list[str],
    *,
    font_size: int,
    compact: bool,
) -> None:
    """按固定条数、字号和行距渲染要点，避免把长文本硬塞入单页。"""

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = _COLOR_WHITE
    box.line.color.rgb = RGBColor(215, 228, 249)
    text_frame = box.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.34)
    text_frame.margin_right = Inches(0.28)
    text_frame.margin_top = Inches(0.24)
    text_frame.margin_bottom = Inches(0.2)
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    for index, bullet in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        # python-pptx 不提供稳定的跨模板项目符号开关；把项目符号作为真实文本写入，可避免
        # 不同 Office 主题下丢失 bullet 样式，也便于回读验证与无 Office 环境验收。
        paragraph.text = f"• {_normalize_text(bullet, limit=170 if compact else 150)}"
        paragraph.level = 0
        paragraph.font.name = "Microsoft YaHei UI"
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = _COLOR_TEXT
        paragraph.space_after = Pt(5 if compact else 10)
        paragraph.line_spacing = 1.15


def _add_source_footer(slide: object, sources: list[DocumentSourceRef], *, page_label: str) -> None:
    """在每页底部留下紧凑、可阅读的来源提示，不显示本机绝对路径。"""

    source_text = "；".join(_source_label(source) for source in sources[:2])
    if not source_text:
        source_text = "来源：已核验草稿快照"
    _add_text_box(slide, Inches(0.8), Inches(6.86), Inches(10.8), Inches(0.32), source_text, font_size=9, color=_COLOR_MUTED)
    _add_text_box(slide, Inches(11.75), Inches(6.84), Inches(0.8), Inches(0.32), page_label, font_size=9, color=_COLOR_MUTED)


def _verify_presentation(path: Path, preview: PresentationPreviewResponse) -> PresentationVerification:
    """重新打开新建文件并验证页数、标题、正文与来源页，防止“生成成功”只是空壳。"""

    if not path.exists() or path.stat().st_size <= 0:
        raise PresentationConflictError("项目方案 PPT 未生成有效文件。")
    opened = Presentation(path)
    if len(opened.slides) != len(preview.slides):
        raise PresentationConflictError("项目方案 PPT 页数与已确认计划不一致。")

    for index, (slide, planned) in enumerate(zip(opened.slides, preview.slides), start=1):
        text = "\n".join(shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False))
        if planned.title not in text:
            raise PresentationConflictError(f"第 {index} 页缺少已确认标题，已停止交付。")
        if planned.role == "content" and not any(bullet in text for bullet in planned.bullets):
            raise PresentationConflictError(f"第 {index} 页缺少已确认正文要点，已停止交付。")

    source_slide_count = sum(1 for slide in preview.slides if slide.role == "sources")
    if source_slide_count != 1:
        raise PresentationConflictError("项目方案 PPT 必须包含唯一的来源追溯页。")
    return PresentationVerification(
        passed=True,
        slide_count=len(opened.slides),
        source_slide_count=source_slide_count,
    )
