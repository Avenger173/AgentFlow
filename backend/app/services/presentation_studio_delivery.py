"""把 PPT 制作 V2 的已确认创作计划渲染为可编辑演示文稿。

创作计划和文件交付刻意拆开：前者允许模型判断叙事与视觉方向，后者只读取已持久化的计划、
使用内置视觉 token 渲染、回读验证并记录 artifact。这样确认后的文件不会再次调用模型，
也不会静默联网下载素材。
"""

from __future__ import annotations

import re
from collections import Counter
from collections.abc import Callable
from dataclasses import dataclass
from datetime import UTC, datetime
from io import BytesIO
from pathlib import Path
from urllib.parse import urlparse
from uuid import uuid4

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from app.core.config import settings
from app.database.task_repository import append_workflow_artifact
from app.schemas.presentation import PresentationExportResponse, PresentationVerification
from app.schemas.presentation_studio import (
    PresentationStudioExportRequest,
    PresentationStudioAssetSlot,
    PresentationStudioDataPlan,
    PresentationStudioPlanResponse,
    PresentationStudioSlidePlan,
)
from app.schemas.workflow import WorkflowArtifact
from app.services.pexels_assets import PexelsAssetResolution, PexelsImageAsset, fetch_pexels_images
from app.services.presentation_research_gateway import (
    ResearchGatewayChartData,
    ResearchGatewayDataPoint,
    ResearchGatewayResolution,
    complete_research_resolution_with_ai_draft,
    fetch_ai_knowledge_draft_chart_data,
    fetch_research_gateway_chart_data,
)
from app.services.presentation_native_animations import (
    NativePresentationMotionSummary,
    apply_native_presentation_motion,
    inspect_native_presentation_motion,
)
from app.services.presentation_studio import get_presentation_studio_result
from app.services.seedream_assets import SeedreamAssetResolution, SeedreamImageAsset, generate_seedream_images
from app.services.wikimedia_research import (
    WikimediaResearchResolution,
    WikimediaResearchSource,
    fetch_wikimedia_references,
)
from app.services.world_bank_data import (
    WorldBankChartData,
    WorldBankDataResolution,
    WorldBankDataPoint,
    fetch_world_bank_chart_data,
)
from app.workflow.dry_run import clear_dry_run_memory_cache


_DOCUMENT_AGENT_ID = "document_agent"
_PPTX_MIME_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_StudioVisualAsset = PexelsImageAsset | SeedreamImageAsset
_StudioVisualAssetResolution = PexelsAssetResolution | SeedreamAssetResolution
_StructuredDataChart = WorldBankChartData | ResearchGatewayChartData
_StructuredDataResolution = WorldBankDataResolution | ResearchGatewayResolution
_ExportProgressCallback = Callable[[str, str, str], None]


class PresentationStudioDeliveryError(RuntimeError):
    """PPT 制作 V2 写入阶段的可预期错误。"""


class PresentationStudioPlanNotFoundError(PresentationStudioDeliveryError):
    """任务不存在，或不是可交付的 PPT 创作计划。"""


class PresentationStudioPlanConflictError(PresentationStudioDeliveryError):
    """计划身份不一致、同名文件或回读验证失败。"""


class PresentationStudioConfirmationError(PresentationStudioDeliveryError):
    """客户端没有完成明确确认，不能创建文件。"""


@dataclass(frozen=True)
class _StudioThemePalette:
    # 设计系统决定构图规则，不能把主题退化成仅替换几个 RGB 值。
    design_system: str
    background: RGBColor
    surface: RGBColor
    primary: RGBColor
    accent: RGBColor
    heading: RGBColor
    text: RGBColor
    muted: RGBColor
    border: RGBColor


_THEME_PALETTES = {
    "executive_blue": _StudioThemePalette(
        design_system="executive",
        background=RGBColor(245, 248, 253),
        surface=RGBColor(255, 255, 255),
        primary=RGBColor(35, 100, 232),
        accent=RGBColor(119, 183, 255),
        heading=RGBColor(20, 34, 74),
        text=RGBColor(31, 48, 84),
        muted=RGBColor(93, 111, 140),
        border=RGBColor(215, 228, 249),
    ),
    "technology_emerald": _StudioThemePalette(
        design_system="technology",
        background=RGBColor(243, 251, 249),
        surface=RGBColor(255, 255, 255),
        primary=RGBColor(15, 137, 116),
        accent=RGBColor(112, 205, 178),
        heading=RGBColor(19, 67, 61),
        text=RGBColor(34, 78, 72),
        muted=RGBColor(85, 118, 111),
        border=RGBColor(204, 232, 224),
    ),
    "narrative_warm": _StudioThemePalette(
        design_system="editorial",
        background=RGBColor(254, 248, 243),
        surface=RGBColor(255, 255, 255),
        primary=RGBColor(197, 96, 61),
        accent=RGBColor(239, 176, 118),
        heading=RGBColor(86, 43, 29),
        text=RGBColor(95, 59, 43),
        muted=RGBColor(132, 104, 89),
        border=RGBColor(243, 219, 200),
    ),
    "impact_contrast": _StudioThemePalette(
        design_system="impact",
        background=RGBColor(244, 247, 251),
        surface=RGBColor(255, 255, 255),
        primary=RGBColor(22, 43, 85),
        accent=RGBColor(240, 105, 75),
        heading=RGBColor(20, 31, 54),
        text=RGBColor(50, 62, 84),
        muted=RGBColor(105, 117, 137),
        border=RGBColor(216, 224, 236),
    ),
}


def export_presentation_studio_plan(
    *,
    task_id: str,
    request: PresentationStudioExportRequest,
    progress_callback: _ExportProgressCallback | None = None,
) -> PresentationExportResponse:
    """确认后写入新 PPTX、回读验证并将 artifact 追加到创作计划任务。"""

    if not request.confirmed:
        raise PresentationStudioConfirmationError("导出 PPT 前需要用户确认当前创作计划。")
    _emit_progress(
        progress_callback,
        "presentation_export_started",
        "已确认本次导出，正在校验创作计划与受控输出位置。",
    )
    plan = _load_plan(task_id)
    if request.plan_id != plan.plan_id:
        raise PresentationStudioPlanConflictError("PPT 创作计划已失效，请重新查看计划后再确认导出。")

    filename = _safe_filename(request.filename, fallback_title=plan.brief.title, task_id=task_id)
    output_root = settings.document_presentation_output_dir
    output_root.mkdir(parents=True, exist_ok=True)
    target = (output_root / filename).resolve()
    try:
        target.relative_to(output_root)
    except ValueError as exc:
        raise PresentationStudioDeliveryError("演示文稿名称未通过受控输出目录校验。") from exc
    if target.exists():
        raise PresentationStudioPlanConflictError(
            f"output/document_presentations 中已存在同名文件“{filename}”，请改名后再次确认导出。"
        )

    # 先完成路径与覆盖检查再调用外部视觉 Provider，避免重复导出白白消耗图库或生成额度。
    asset_slots = _effective_asset_slots(plan)
    if request.fetch_external_assets or request.fetch_licensed_assets:
        _emit_progress(progress_callback, "presentation_visual_started", "正在准备本次确认的外部视觉素材。")
    assets = _resolve_visual_assets(plan=plan, request=request, asset_slots=asset_slots)
    assets_by_slide_id = _assign_assets_to_slides(asset_slots=asset_slots, assets=assets.images)
    if request.fetch_public_research:
        _emit_progress(progress_callback, "presentation_public_reference_started", "正在读取本次确认的公开资料参考。")
    research = _resolve_public_research(plan=plan, request=request)
    structured_data = _resolve_structured_data(
        plan=plan,
        request=request,
        progress_callback=progress_callback,
    )
    structured_data_charts = _structured_data_charts(structured_data)
    data_contract_gap = _structured_data_contract_gap(plan.data_plan, structured_data_charts)

    _emit_progress(progress_callback, "presentation_render_started", "正在写入可编辑 PPTX，并准备回读验证。")

    try:
        # 使用 x 模式拒绝覆盖。渲染器只拥有本次新建句柄，失败时也只撤回它创建的文件。
        with target.open("xb") as target_file:
            motion = _render_studio_presentation(
                target_file,
                plan,
                assets=assets.images,
                assets_by_slide_id=assets_by_slide_id,
                research_sources=research.sources,
                structured_data=structured_data_charts,
            )
    except FileExistsError as exc:
        raise PresentationStudioPlanConflictError(
            f"output/document_presentations 中已存在同名文件“{filename}”，请改名后再次确认导出。"
        ) from exc
    except Exception as exc:
        target.unlink(missing_ok=True)
        raise PresentationStudioDeliveryError(f"无法生成 PPT 创作文件：{exc}") from exc

    try:
        verification = _verify_studio_presentation(
            target,
            plan,
            assets=assets,
            research=research,
            structured_data=structured_data,
            motion=motion,
        )
    except Exception as exc:
        target.unlink(missing_ok=True)
        if isinstance(exc, PresentationStudioDeliveryError):
            raise
        raise PresentationStudioDeliveryError(f"PPT 创作文件回读验证失败：{exc}") from exc

    _emit_progress(progress_callback, "presentation_render_verified", "PPTX 已通过回读验证，正在写入任务交付记录。")

    artifact_id = f"{task_id}:presentation_studio_ppt:{uuid4().hex[:10]}"
    relative_path = f"output/document_presentations/{filename}"
    artifact = WorkflowArtifact(
        artifact_id=artifact_id,
        task_id=task_id,
        step_id="presentation_studio_export",
        agent_id=_DOCUMENT_AGENT_ID,
        kind="file",
        name=filename,
        summary=(
            f"用户确认导出的 PPT 创作文件，包含 {verification.slide_count} 页多版式页面和 "
            f"{len(assets.images)} 张外部视觉素材、{len(research.sources)} 条公开资料参考"
            f"、{len(structured_data_charts)} 个可编辑数据视图和 "
            f"{motion.entrance_effect_count} 个原生入场动效。"
        ),
        uri=f"agentflow-output://document_presentations/{filename}",
        mime_type=_PPTX_MIME_TYPE,
        metadata={
            "runtime": True,
            "output_scope": "document_presentations",
            "output_path": str(target),
            "relative_output_path": relative_path,
            "confirmed_by": "local_user",
            "presentation_mode": "studio_v2",
            "plan_id": plan.plan_id,
            "theme": plan.brief.theme,
            "slide_count": verification.slide_count,
            "source_slide_count": verification.source_slide_count,
            "verification_passed": verification.passed,
            "native_presentation_motion": motion.audit_metadata(),
            "external_assets_fetched": bool(assets.images),
            "asset_plan_state": plan.asset_plan.state,
            "asset_provider": assets.provider if assets.images else "",
            "asset_count": len(assets.images),
            "asset_slot_count": len(asset_slots),
            "asset_assignments": [
                {
                    "slide_id": slot.slide_id,
                    "slide_title": slot.slide_title,
                    "query": slot.query,
                    "embedded": slot.slide_id in assets_by_slide_id,
                }
                for slot in asset_slots
            ],
            "asset_sources": [image.audit_metadata() for image in assets.images],
            "asset_warnings": list(assets.warnings),
            "public_research_requested": request.fetch_public_research,
            "public_research_provider": research.provider if research.sources else "",
            "public_research_count": len(research.sources),
            "public_research_sources": [source.audit_metadata() for source in research.sources],
            "public_research_warnings": list(research.warnings),
            "structured_data_requested": request.fetch_structured_data,
            "structured_data_provider": _structured_data_provider(structured_data),
            # 单数键保留给旧历史详情；新数据章节完整记录每个实际落盘的视图。
            "structured_data_chart": structured_data_charts[0].audit_metadata() if structured_data_charts else None,
            "structured_data_charts": [chart.audit_metadata() for chart in structured_data_charts],
            "structured_data_warnings": list(structured_data.warnings),
            "structured_data_contract_complete": not bool(data_contract_gap),
            "structured_data_contract_gap": data_contract_gap,
        },
        created_at=datetime.now(UTC).isoformat(timespec="seconds"),
    )
    try:
        append_workflow_artifact(
            artifact=artifact,
            event_name="artifact_saved",
            message=f"用户已确认导出 PPT 创作文件：{relative_path}",
        )
    except Exception as exc:
        target.unlink(missing_ok=True)
        raise PresentationStudioDeliveryError("PPT 创作文件审计失败，已撤回本次新建文件。") from exc

    clear_dry_run_memory_cache()
    if data_contract_gap:
        delivery_message = (
            "PPT 创作文件已导出并通过文件回读，但客户明确的数据视图合同仅部分完成："
            f"{data_contract_gap}。{_motion_delivery_suffix(motion)}"
        )
    elif assets.images:
        delivery_message = (
            f"PPT 创作文件已导出并通过回读验证，已嵌入 {len(assets.images)} 张{assets.label}，"
            f"并记录 {len(research.sources)} 条公开资料参考和 "
            f"{len(structured_data_charts)} 个可编辑数据视图。{_motion_delivery_suffix(motion)}"
        )
    elif research.sources or structured_data_charts:
        delivery_message = (
            "PPT 创作文件已按多版式内置视觉 token 导出并通过回读验证，已记录 "
            f"{len(research.sources)} 条公开资料参考和 {len(structured_data_charts)} 个可编辑数据视图。"
            f"{_motion_delivery_suffix(motion)}"
        )
    else:
        delivery_message = "PPT 创作文件已按多版式内置视觉 token 导出并通过回读验证。" + _motion_delivery_suffix(motion)
    return PresentationExportResponse(
        task_id=task_id,
        artifact_id=artifact_id,
        filename=filename,
        relative_path=relative_path,
        artifact_uri=artifact.uri,
        slide_count=verification.slide_count,
        verification=verification,
        message=delivery_message,
    )


def _emit_progress(
    callback: _ExportProgressCallback | None,
    event: str,
    message: str,
    level: str = "info",
) -> None:
    """同步交付层只报告已经发生的阶段，不伪造百分比。"""

    if callback is not None:
        callback(event, message, level)


def _motion_delivery_suffix(motion: NativePresentationMotionSummary) -> str:
    """把真实动效状态体现在交付回执中，避免用户只能自行打开 PowerPoint 才发现降级。"""

    if motion.enabled:
        return f"已写入 {motion.transition_slide_count} 页淡入转场和 {motion.entrance_effect_count} 个点击入场动效。"
    return "本次未写入原生动效，已自动保留为无动画的可编辑版本。"


def _load_plan(task_id: str) -> PresentationStudioPlanResponse:
    plan = get_presentation_studio_result(task_id)
    if plan is None:
        raise PresentationStudioPlanNotFoundError("未找到可导出的 PPT 创作计划。")
    if not plan.slides or plan.slides[0].role != "cover" or plan.slides[-1].role != "sources":
        raise PresentationStudioPlanConflictError("PPT 创作计划缺少完整的封面或事实核验页，不能导出。")
    return plan


def _safe_filename(raw_filename: str, *, fallback_title: str, task_id: str) -> str:
    candidate = raw_filename.strip() or f"{fallback_title or 'AgentFlow 演示方案'}-{task_id[-6:]}.pptx"
    if "/" in candidate or "\\" in candidate or Path(candidate).name != candidate:
        raise PresentationStudioDeliveryError("演示文稿名称只能是文件名，不能包含目录或路径分隔符。")
    if not candidate.lower().endswith(".pptx"):
        raise PresentationStudioDeliveryError("演示文稿名称必须以 .pptx 结尾。")
    stem = re.sub(r'[<>:"|?*\x00-\x1f]+', "-", candidate[:-5])
    stem = re.sub(r"\s+", " ", stem).strip(" .-")
    return f"{(stem or f'AgentFlow 演示方案-{task_id[-6:]}')[:96]}.pptx"


def _resolve_visual_assets(
    *,
    plan: PresentationStudioPlanResponse,
    request: PresentationStudioExportRequest,
    asset_slots: tuple[PresentationStudioAssetSlot, ...],
) -> _StudioVisualAssetResolution:
    """只在二次确认后的导出阶段调用已选视觉 Provider，并允许无图降级完成交付。"""

    fetch_external_assets = request.fetch_external_assets or request.fetch_licensed_assets
    if not fetch_external_assets:
        notice = (
            "本次未选择获取外部视觉素材，已使用无图设计版式。"
            if plan.asset_plan.state == "planned"
            else "本次创作计划未请求外部视觉素材，已使用无图设计版式。"
        )
        return PexelsAssetResolution(images=(), warnings=(notice,))
    if not request.network_confirmed:
        raise PresentationStudioConfirmationError("使用外部视觉素材前需要明确确认本次调用。")
    if plan.asset_plan.state != "planned" or not plan.asset_plan.provider:
        return PexelsAssetResolution(
            images=(),
            warnings=("当前计划没有可执行的外部视觉素材建议，已使用无图设计版式。",),
        )
    if plan.asset_plan.provider == "pexels":
        return fetch_pexels_images(plan.asset_plan.queries, limit=len(asset_slots))
    if plan.asset_plan.provider == "seedream":
        return generate_seedream_images(plan.asset_plan.queries, limit=min(4, len(asset_slots)))
    return PexelsAssetResolution(
        images=(),
        warnings=("当前计划引用了不受支持的视觉 Provider，已使用无图设计版式。",),
    )


def _resolve_public_research(
    *,
    plan: PresentationStudioPlanResponse,
    request: PresentationStudioExportRequest,
) -> WikimediaResearchResolution:
    """仅在计划和本次确认都同意时读取固定公开资料接口。"""

    if not request.fetch_public_research:
        return WikimediaResearchResolution(sources=(), warnings=())
    if not request.network_confirmed:
        raise PresentationStudioConfirmationError("补充公开资料前需要明确确认本次联网调用。")
    if plan.research_plan.state != "planned" or plan.research_plan.provider != "wikimedia":
        return WikimediaResearchResolution(
            sources=(),
            warnings=("当前计划未预留公开资料参考，本次不会扩大联网范围。",),
        )
    queries = _research_queries(plan, limit=plan.research_plan.max_sources)
    return fetch_wikimedia_references(queries, limit=plan.research_plan.max_sources)


def _resolve_structured_data(
    *,
    plan: PresentationStudioPlanResponse,
    request: PresentationStudioExportRequest,
    progress_callback: _ExportProgressCallback | None = None,
) -> _StructuredDataResolution:
    """执行已固化的数据计划；普通创作默认直出模型数据，联网研究不是前置条件。"""

    if not request.fetch_structured_data:
        return WorldBankDataResolution(chart=None, warnings=())
    if plan.data_plan.state == "research_planned" and plan.data_plan.provider == "research_gateway":
        # 客户已经确认写入当前 PPTX；普通数据型演示只需调用已配置的模型。此前先联网取证
        # 再补 AI 草稿，会让“我只要数据和图表”的请求被无关的网页失败卡住。
        if plan.data_plan.evidence_mode == "ai_direct":
            return fetch_ai_knowledge_draft_chart_data(
                plan.data_plan,
                progress_callback=progress_callback,
            )
        if not request.network_confirmed:
            raise PresentationStudioConfirmationError("联网核验数据图表前需要明确确认本次联网调用。")
        verified = fetch_research_gateway_chart_data(
            plan.data_plan,
            progress_callback=progress_callback,
        )
        return complete_research_resolution_with_ai_draft(
            plan.data_plan,
            verified,
            progress_callback=progress_callback,
        )
    if (
        plan.data_plan.state not in {"planned", "provider_planned"}
        or plan.data_plan.provider != "world_bank"
    ):
        return WorldBankDataResolution(
            chart=None,
            warnings=("当前计划未预留可验证数据图表，本次不会扩大数据接口范围。",),
        )
    # 固定 World Bank Provider 仍是真实联网读取；只有普通 ``ai_direct`` 数据创作移除了
    # 联网确认前置条件，不能把这项安全边界一并放开。
    if not request.network_confirmed:
        raise PresentationStudioConfirmationError("读取 World Bank 数据前需要明确确认本次联网调用。")
    return fetch_world_bank_chart_data(plan.data_plan)


def _structured_data_provider(resolution: _StructuredDataResolution) -> str:
    """artifact 只声明本次实际成功的 Provider，避免计划状态被误报成已读取数据。"""

    if resolution.chart is None:
        return ""
    if isinstance(resolution.chart, WorldBankChartData):
        return "world_bank"
    charts = _structured_data_charts(resolution)
    evidence_levels = {
        chart.evidence_level for chart in charts if isinstance(chart, ResearchGatewayChartData)
    }
    if evidence_levels == {"ai_knowledge_draft"}:
        return "ai_knowledge_draft"
    if "ai_knowledge_draft" in evidence_levels:
        return "research_gateway+ai_knowledge_draft"
    return "research_gateway"


def _structured_data_charts(
    resolution: _StructuredDataResolution,
) -> tuple[_StructuredDataChart, ...]:
    """统一旧 Provider 的单图返回和 ResearchGateway 的数据章节返回。"""

    if isinstance(resolution, ResearchGatewayResolution) and resolution.charts:
        return resolution.charts
    return (resolution.chart,) if resolution.chart is not None else ()


def _structured_data_contract_gap(
    contract: PresentationStudioDataPlan,
    charts: tuple[_StructuredDataChart, ...],
) -> str:
    """逐类统计真实交付物；只有用户明确点名数量时才把缺口提升为部分完成。"""

    if not contract.visual_contract_explicit:
        return ""
    table_total = sum(chart.chart_type in {"comparison_table", "trend_table"} for chart in charts)
    bar_total = sum(chart.chart_type in {"comparison_bar", "grouped_bar", "horizontal_bar"} for chart in charts)
    line_total = sum(chart.chart_type in {"trend_line", "trend_area"} for chart in charts)
    gaps: list[str] = []
    if table_total < contract.required_table_count:
        gaps.append(f"表格 {table_total}/{contract.required_table_count}")
    if bar_total < contract.required_bar_chart_count:
        gaps.append(f"柱状图 {bar_total}/{contract.required_bar_chart_count}")
    if line_total < contract.required_line_chart_count:
        gaps.append(f"折线图 {line_total}/{contract.required_line_chart_count}")
    if len(charts) < contract.required_visual_count:
        gaps.append(f"数据视图 {len(charts)}/{contract.required_visual_count}")
    return "、".join(gaps)


def _research_queries(plan: PresentationStudioPlanResponse, *, limit: int) -> tuple[str, ...]:
    """从已确认计划派生短查询，不允许模型在导出阶段重新生成联网指令。"""

    candidates = [plan.brief.title]
    candidates.extend(slide.title for slide in plan.slides if slide.role == "content")
    queries: list[str] = []
    seen: set[str] = set()
    for candidate in candidates:
        query = " ".join(candidate.split()).strip()
        key = query.casefold()
        if len(query) < 2 or key in seen:
            continue
        seen.add(key)
        queries.append(query[:140])
        if len(queries) >= max(1, min(limit, 3)):
            break
    return tuple(queries)


def _effective_asset_slots(plan: PresentationStudioPlanResponse) -> tuple[PresentationStudioAssetSlot, ...]:
    """为 V2.1 历史计划补出稳定的页面槽位，不让未导出的旧计划在升级后失效。

    新计划必须持久化 ``slots``，以便每张图片由自己的语义查询绑定页面。历史 V2.1 快照只有
    有序查询词；这里按当时既有的“封面优先、随后正文页”规则在内存恢复同一映射，只供这一次
    导出和 artifact 审计使用，不修改原 SQLite 快照。
    """

    if plan.asset_plan.slots:
        return tuple(plan.asset_plan.slots)

    eligible_slides = [
        slide for slide in plan.slides if slide.role == "cover" or slide.role == "content"
    ]
    return tuple(
        PresentationStudioAssetSlot(
            slide_id=slide.slide_id,
            slide_title=slide.title,
            query=query,
            purpose="兼容历史 PPT 计划的有序外部图片映射。",
        )
        for slide, query in zip(eligible_slides, plan.asset_plan.queries, strict=False)
    )


def _assign_assets_to_slides(
    *,
    asset_slots: tuple[PresentationStudioAssetSlot, ...],
    assets: tuple[_StudioVisualAsset, ...],
) -> dict[str, _StudioVisualAsset]:
    """只按已持久化的槽位映射素材；任一查询失败都不会让后续图片滑到错误页面。"""

    images_by_query = {image.query.casefold(): image for image in assets}
    return {
        slot.slide_id: image
        for slot in asset_slots
        if (image := images_by_query.get(slot.query.casefold())) is not None
    }


def _render_studio_presentation(
    target_file: object,
    plan: PresentationStudioPlanResponse,
    *,
    assets: tuple[_StudioVisualAsset, ...],
    assets_by_slide_id: dict[str, _StudioVisualAsset],
    research_sources: tuple[WikimediaResearchSource, ...],
    structured_data: tuple[_StructuredDataChart, ...],
) -> NativePresentationMotionSummary:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    palette = _THEME_PALETTES[plan.brief.theme]
    for index, slide_plan in enumerate(plan.slides, start=1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        _paint_background(slide, palette)
        if slide_plan.role == "cover":
            _render_cover(slide, plan, palette, image=assets_by_slide_id.get(slide_plan.slide_id))
        else:
            image = _content_image_for_slide(slide_plan, assets_by_slide_id=assets_by_slide_id)
            _render_standard(
                slide,
                slide_plan,
                index=index,
                total=len(plan.slides),
                palette=palette,
                image=image,
                asset_sources=assets,
                research_sources=research_sources,
                structured_data=structured_data,
            )
    # 先完成所有可编辑对象，再注入动画时间线。这样任何 animation fallback 都不会影响图表、
    # 表格和文本形状本身，且动效仍引用最终已落盘的 shape id。
    motion = apply_native_presentation_motion(
        presentation,
        slide_roles=tuple(slide.role for slide in plan.slides),
    )
    presentation.save(target_file)
    return motion


def _paint_background(slide: object, palette: _StudioThemePalette) -> None:
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = palette.background
    if palette.design_system == "technology":
        rail = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.28), Inches(7.5))
        rail.fill.solid()
        rail.fill.fore_color.rgb = palette.primary
        rail.line.fill.background()
        for index in range(4):
            marker = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(12.26 - index * 0.31),
                Inches(0.48 + index * 0.31),
                Inches(0.12),
                Inches(0.12),
            )
            marker.fill.solid()
            marker.fill.fore_color.rgb = palette.accent
            marker.line.fill.background()
        return
    if palette.design_system == "editorial":
        upper_rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.42), Inches(11.72), Inches(0.04))
        upper_rule.fill.solid()
        upper_rule.fill.fore_color.rgb = palette.accent
        upper_rule.line.fill.background()
        return
    if palette.design_system == "impact":
        top_band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.42))
        top_band.fill.solid()
        top_band.fill.fore_color.rgb = palette.primary
        top_band.line.fill.background()
        accent_block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.98), Inches(0.42), Inches(1.35), Inches(7.08))
        accent_block.fill.solid()
        accent_block.fill.fore_color.rgb = palette.accent
        accent_block.fill.transparency = 8
        accent_block.line.fill.background()
        return
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.16))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = palette.primary
    top_bar.line.fill.background()
    accent = slide.shapes.add_shape(MSO_SHAPE.ARC, Inches(10.7), Inches(-0.7), Inches(3.3), Inches(3.3))
    accent.fill.solid()
    accent.fill.fore_color.rgb = palette.accent
    accent.fill.transparency = 55
    accent.line.fill.background()


def _render_cover(
    slide: object,
    plan: PresentationStudioPlanResponse,
    palette: _StudioThemePalette,
    *,
    image: _StudioVisualAsset | None,
) -> None:
    if palette.design_system == "technology":
        _render_technology_cover(slide, plan, palette, image=image)
        return
    if palette.design_system == "editorial":
        _render_editorial_cover(slide, plan, palette, image=image)
        return
    if palette.design_system == "impact":
        _render_impact_cover(slide, plan, palette, image=image)
        return
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.78), Inches(1.04), Inches(11.72), Inches(5.18))
    panel.fill.solid()
    panel.fill.fore_color.rgb = palette.surface
    panel.line.color.rgb = palette.border
    ribbon = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.78), Inches(1.04), Inches(0.18), Inches(5.18))
    ribbon.fill.solid()
    ribbon.fill.fore_color.rgb = palette.primary
    ribbon.line.fill.background()
    text_width = Inches(5.55) if image is not None else Inches(9.8)
    _add_text(slide, Inches(1.28), Inches(1.65), text_width, Inches(1.32), plan.brief.title, 30, palette.heading, True)
    _add_text(slide, Inches(1.3), Inches(3.15), text_width, Inches(0.9), plan.brief.core_message, 17, palette.primary, False)
    _add_text(slide, Inches(1.3), Inches(4.47), text_width, Inches(0.66), "AgentFlow · 由用户意图生成的创作计划", 13, palette.muted, False)
    _add_text(slide, Inches(1.3), Inches(5.12), text_width, Inches(0.45), "导出前请核验关键事实、数据与案例。", 11, palette.muted, False)
    if image is not None:
        _add_image(slide, image, left=Inches(7.05), top=Inches(1.62), width=Inches(4.58), height=Inches(2.4))
        _add_text(slide, Inches(7.05), Inches(4.16), Inches(4.58), Inches(0.32), image.credit_text, 8, palette.muted, False)
    _add_footer(slide, "创作依据：用户主题说明", "1 / %s" % len(plan.slides), palette)


def _render_technology_cover(
    slide: object,
    plan: PresentationStudioPlanResponse,
    palette: _StudioThemePalette,
    *,
    image: _StudioVisualAsset | None,
) -> None:
    """技术洞察主题采用网格化分栏与左侧系统轨道，区别于商务卡片封面。"""

    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.82), Inches(1.02), Inches(11.25), Inches(5.3))
    panel.fill.solid()
    panel.fill.fore_color.rgb = palette.surface
    panel.line.color.rgb = palette.border
    _add_text(slide, Inches(1.24), Inches(1.5), Inches(6.1), Inches(1.22), plan.brief.title, 31, palette.heading, True)
    _add_text(slide, Inches(1.26), Inches(3.02), Inches(5.45), Inches(0.9), plan.brief.core_message, 17, palette.primary, False)
    _add_text(slide, Inches(1.26), Inches(4.86), Inches(5.35), Inches(0.38), "结构化方案 · 关键事实导出前核验", 11, palette.muted, False)
    for index, label in enumerate(("目标", "路径", "交付"), start=1):
        top = Inches(4.08 + (index - 1) * 0.44)
        marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.27), top, Inches(0.17), Inches(0.17))
        marker.fill.solid()
        marker.fill.fore_color.rgb = palette.accent
        marker.line.fill.background()
        _add_text(slide, Inches(1.55), top - Inches(0.04), Inches(1.25), Inches(0.26), label, 10, palette.text, True)
    if image is not None:
        _add_image(slide, image, left=Inches(7.24), top=Inches(1.46), width=Inches(4.13), height=Inches(3.26))
        _add_text(slide, Inches(7.24), Inches(4.87), Inches(4.13), Inches(0.22), image.credit_text, 8, palette.muted, False)
    else:
        for row in range(3):
            for column in range(3):
                cell = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(7.28 + column * 1.34),
                    Inches(1.58 + row * 0.92),
                    Inches(1.1),
                    Inches(0.68),
                )
                cell.fill.solid()
                cell.fill.fore_color.rgb = palette.background if (row + column) % 2 else palette.surface
                cell.line.color.rgb = palette.border
    _add_footer(slide, "创作依据：用户主题说明", "1 / %s" % len(plan.slides), palette)


def _render_editorial_cover(
    slide: object,
    plan: PresentationStudioPlanResponse,
    palette: _StudioThemePalette,
    *,
    image: _StudioVisualAsset | None,
) -> None:
    """叙事展示主题使用留白、引言和图片边栏，强调一个核心观点而不是企业仪表盘。"""

    _add_text(slide, Inches(1.02), Inches(1.26), Inches(7.1), Inches(1.45), plan.brief.title, 34, palette.heading, True)
    quote_mark = slide.shapes.add_shape(MSO_SHAPE.ARC, Inches(1.02), Inches(3.0), Inches(0.55), Inches(0.55))
    quote_mark.fill.solid()
    quote_mark.fill.fore_color.rgb = palette.accent
    quote_mark.fill.transparency = 20
    quote_mark.line.fill.background()
    _add_text(slide, Inches(1.75), Inches(3.04), Inches(5.85), Inches(1.02), plan.brief.core_message, 18, palette.primary, False)
    _add_text(slide, Inches(1.04), Inches(5.42), Inches(6.1), Inches(0.36), "一份先讲清主张、再展开论证的演示计划", 12, palette.muted, False)
    if image is not None:
        _add_image(slide, image, left=Inches(8.42), top=Inches(1.2), width=Inches(3.0), height=Inches(4.45))
        _add_text(slide, Inches(8.42), Inches(5.78), Inches(3.0), Inches(0.22), image.credit_text, 8, palette.muted, False)
    else:
        block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.42), Inches(1.2), Inches(3.0), Inches(4.45))
        block.fill.solid()
        block.fill.fore_color.rgb = palette.accent
        block.fill.transparency = 58
        block.line.fill.background()
    _add_footer(slide, "创作依据：用户主题说明", "1 / %s" % len(plan.slides), palette)


def _render_impact_cover(
    slide: object,
    plan: PresentationStudioPlanResponse,
    palette: _StudioThemePalette,
    *,
    image: _StudioVisualAsset | None,
) -> None:
    """强调对比主题以强标题、编号和不对称图像区建立演讲式开场。"""

    number = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.84), Inches(1.0), Inches(1.18), Inches(1.18))
    number.fill.solid()
    number.fill.fore_color.rgb = palette.accent
    number.line.fill.background()
    _add_text(slide, Inches(1.13), Inches(1.28), Inches(0.62), Inches(0.4), "01", 18, palette.surface, True)
    _add_text(slide, Inches(0.86), Inches(2.48), Inches(6.35), Inches(1.54), plan.brief.title, 36, palette.heading, True)
    _add_text(slide, Inches(0.88), Inches(4.34), Inches(5.65), Inches(0.96), plan.brief.core_message, 17, palette.primary, False)
    if image is not None:
        _add_image(slide, image, left=Inches(7.26), top=Inches(1.08), width=Inches(4.12), height=Inches(4.32))
        _add_text(slide, Inches(7.26), Inches(5.55), Inches(4.12), Inches(0.22), image.credit_text, 8, palette.muted, False)
    else:
        visual = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(7.32), Inches(1.3), Inches(3.5), Inches(3.72))
        visual.fill.solid()
        visual.fill.fore_color.rgb = palette.primary
        visual.fill.transparency = 14
        visual.line.fill.background()
    _add_footer(slide, "创作依据：用户主题说明", "1 / %s" % len(plan.slides), palette)


def _render_standard(
    slide: object,
    plan: PresentationStudioSlidePlan,
    *,
    index: int,
    total: int,
    palette: _StudioThemePalette,
    image: _StudioVisualAsset | None,
    asset_sources: tuple[_StudioVisualAsset, ...],
    research_sources: tuple[WikimediaResearchSource, ...],
    structured_data: tuple[_StructuredDataChart, ...],
) -> None:
    title_left = Inches(1.02) if palette.design_system == "technology" else Inches(0.8)
    title_top = Inches(0.74) if palette.design_system == "impact" else Inches(0.66)
    _add_text(slide, title_left, title_top, Inches(10.5), Inches(0.6), plan.title, 24, palette.heading, True)
    underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, title_left, title_top + Inches(0.72), Inches(0.82), Inches(0.07))
    underline.fill.solid()
    underline.fill.fore_color.rgb = palette.primary
    underline.line.fill.background()
    if plan.role == "agenda":
        _render_agenda(slide, plan, palette)
    elif plan.role == "summary":
        _render_summary(slide, plan, palette)
    elif plan.role == "sources":
        _render_sources(
            slide,
            plan,
            palette,
            asset_sources=asset_sources,
            research_sources=research_sources,
            structured_data=structured_data,
        )
    elif chart := next((item for item in structured_data if item.slide_id == plan.slide_id), None):
        _render_structured_data_chart(slide, plan, palette, chart)
    elif image is not None:
        # 受控素材已经绑定到该页。此时优先确保图片实际嵌入且有明确文字关系，避免审计素材
        # 因“选择了无图版式”而丢失；没有外图时才完整体现下面的版式语法。
        _render_image_statement(slide, plan, palette, image=image)
    else:
        _render_content_layout(slide, plan, palette)
    if plan.role not in {"sources", "agenda"} and plan.visual_direction:
        _add_text(slide, Inches(0.98), Inches(6.45), Inches(10.4), Inches(0.28), plan.visual_direction, 9, palette.muted, False)
    footer = "创作依据：用户主题说明"
    if plan.role == "sources":
        footer = (
            "事实边界：公开资料仅作来源参考，关键事实仍需核验"
            if research_sources
            else "事实边界：未联网核验，需确认后使用"
        )
    _add_footer(slide, footer, f"{index} / {total}", palette)


def _render_content_layout(
    slide: object,
    plan: PresentationStudioSlidePlan,
    palette: _StudioThemePalette,
) -> None:
    """按受控 layout 落版；未知值在计划层已归一，这里仍保留卡片兜底。"""

    if plan.layout == "comparison":
        _render_comparison(slide, plan, palette)
    elif plan.layout == "process":
        _render_process(slide, plan, palette)
    elif plan.layout == "timeline":
        _render_timeline(slide, plan, palette)
    elif plan.layout == "metrics":
        _render_metric_tiles(slide, plan, palette)
    elif plan.layout == "quote":
        _render_quote(slide, plan, palette)
    elif plan.layout == "image_statement":
        _render_image_placeholder_statement(slide, plan, palette)
    else:
        _render_insight_cards(slide, plan, palette)


def _render_structured_data_chart(
    slide: object,
    plan: PresentationStudioSlidePlan,
    palette: _StudioThemePalette,
    chart: _StructuredDataChart,
) -> None:
    """用原生表格/Chart 绘制确定性数据视图，确保 PPTX 可编辑且可追溯。"""

    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.98), Inches(1.68), Inches(11.28), Inches(4.62))
    panel.fill.solid()
    panel.fill.fore_color.rgb = palette.surface
    panel.line.color.rgb = palette.border
    _add_text(slide, Inches(1.3), Inches(1.96), Inches(7.1), Inches(0.36), chart.title, 16, palette.heading, True)
    _add_text(slide, Inches(1.3), Inches(2.32), Inches(9.6), Inches(0.27), _chart_subtitle(chart), 9, palette.muted, False)
    if isinstance(chart, ResearchGatewayChartData) and chart.chart_type == "comparison_table":
        _render_research_comparison_table(slide, palette, chart)
    elif isinstance(chart, ResearchGatewayChartData) and chart.chart_type == "trend_table":
        _render_research_trend_table(slide, palette, chart)
    elif isinstance(chart, ResearchGatewayChartData) and chart.chart_type == "grouped_bar":
        _render_research_grouped_bar_chart(slide, palette, chart)
    elif isinstance(chart, ResearchGatewayChartData) and chart.chart_type == "horizontal_bar":
        _render_horizontal_bar_chart(slide, palette, chart)
    elif isinstance(chart, ResearchGatewayChartData) and chart.chart_type in {"share_pie", "share_doughnut"}:
        _render_share_chart(slide, palette, chart)
    elif isinstance(chart, ResearchGatewayChartData) and chart.chart_type == "trend_area":
        _render_trend_area_chart(slide, palette, chart)
    elif chart.chart_type == "comparison_bar":
        _render_comparison_bar_chart(slide, palette, chart)
    else:
        _render_trend_line_chart(slide, palette, chart)
    narrative = " · ".join(_compact_text(item, 50) for item in plan.bullets[:2])
    _add_text(slide, Inches(1.3), Inches(5.68), Inches(10.45), Inches(0.35), narrative, 10, palette.text, False)


def _render_comparison_bar_chart(
    slide: object,
    palette: _StudioThemePalette,
    chart: _StructuredDataChart,
) -> None:
    """同年比较使用 PowerPoint 原生柱状图；年份一致性已在 Provider 层验证。"""

    points = list(chart.points)
    chart_data = CategoryChartData()
    chart_data.categories = [_compact_text(_point_entity(point), 24) for point in points]
    chart_data.add_series(_compact_text(_point_metric(points[0], chart), 40), [point.value for point in points])
    native_chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1.3),
        Inches(2.72),
        Inches(9.85),
        Inches(2.52),
        chart_data,
    ).chart
    _style_native_chart(native_chart, palette, show_legend=False)
    _add_chart_value_ledger(slide, chart, palette, top=5.31)


def _render_trend_line_chart(
    slide: object,
    palette: _StudioThemePalette,
    chart: _StructuredDataChart,
) -> None:
    """趋势图使用 PowerPoint 原生多序列折线图，横轴严格对应已验证期间。"""

    points = list(chart.points)
    periods = list(dict.fromkeys(_point_period(point) for point in points))
    series: dict[str, dict[str, float]] = {}
    for point in points:
        series.setdefault(_point_entity(point), {})[_point_period(point)] = point.value
    chart_data = CategoryChartData()
    chart_data.categories = periods
    for entity, values in series.items():
        chart_data.add_series(_compact_text(entity, 32), [values.get(period) for period in periods])
    native_chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        Inches(1.3),
        Inches(2.72),
        Inches(9.85),
        Inches(2.52),
        chart_data,
    ).chart
    _style_native_chart(native_chart, palette, show_legend=len(series) > 1)
    _add_chart_value_ledger(slide, chart, palette, top=5.31)


def _render_trend_area_chart(
    slide: object,
    palette: _StudioThemePalette,
    chart: ResearchGatewayChartData,
) -> None:
    """连续趋势可使用原生面积图，仍保留可编辑的内嵌工作簿数据。"""

    points = list(chart.points)
    periods = list(dict.fromkeys(_point_period(point) for point in points))
    series: dict[str, dict[str, float]] = {}
    for point in points:
        series.setdefault(_point_entity(point), {})[_point_period(point)] = point.value
    chart_data = CategoryChartData()
    chart_data.categories = periods
    for entity, values in series.items():
        chart_data.add_series(_compact_text(entity, 32), [values.get(period) for period in periods])
    native_chart = slide.shapes.add_chart(
        XL_CHART_TYPE.AREA,
        Inches(1.3),
        Inches(2.72),
        Inches(9.85),
        Inches(2.52),
        chart_data,
    ).chart
    _style_native_chart(native_chart, palette, show_legend=len(series) > 1)
    _add_chart_value_ledger(slide, chart, palette, top=5.31)


def _render_horizontal_bar_chart(
    slide: object,
    palette: _StudioThemePalette,
    chart: ResearchGatewayChartData,
) -> None:
    """横向条形图适合长标签、排名比较和单对象多指标画像。"""

    points = list(chart.points)
    entities = list(dict.fromkeys(point.entity for point in points))
    chart_data = CategoryChartData()
    if len(entities) == 1:
        chart_data.categories = [_compact_text(point.metric, 32) for point in points]
        chart_data.add_series(_compact_text(entities[0], 32), [point.value for point in points])
    else:
        chart_data.categories = [_compact_text(point.entity, 32) for point in points]
        chart_data.add_series(_compact_text(points[0].metric, 40), [point.value for point in points])
    native_chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(1.3),
        Inches(2.72),
        Inches(9.85),
        Inches(2.52),
        chart_data,
    ).chart
    _style_native_chart(native_chart, palette, show_legend=False)
    _add_chart_value_ledger(slide, chart, palette, top=5.31)


def _render_share_chart(
    slide: object,
    palette: _StudioThemePalette,
    chart: ResearchGatewayChartData,
) -> None:
    """饼图与环形图只呈现已经由规划/验证层确认的构成关系。"""

    points = list(chart.points)
    entities = list(dict.fromkeys(point.entity for point in points))
    categories = (
        [_compact_text(point.entity, 28) for point in points]
        if len(entities) > 1
        else [_compact_text(point.metric, 28) for point in points]
    )
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("构成", [point.value for point in points])
    chart_type = XL_CHART_TYPE.PIE if chart.chart_type == "share_pie" else XL_CHART_TYPE.DOUGHNUT
    native_chart = slide.shapes.add_chart(
        chart_type,
        Inches(1.65),
        Inches(2.66),
        Inches(8.9),
        Inches(2.75),
        chart_data,
    ).chart
    native_chart.has_title = False
    native_chart.has_legend = True
    native_chart.legend.position = XL_LEGEND_POSITION.RIGHT
    native_chart.legend.include_in_layout = False
    native_chart.legend.font.name = "Microsoft YaHei UI"
    native_chart.legend.font.size = Pt(9)
    native_chart.chart_style = 10
    plot = native_chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.show_percentage = True
    plot.data_labels.position = XL_LABEL_POSITION.BEST_FIT
    plot.data_labels.font.name = "Microsoft YaHei UI"
    plot.data_labels.font.size = Pt(8)
    _add_chart_value_ledger(slide, chart, palette, top=5.31)


def _render_research_comparison_table(
    slide: object,
    palette: _StudioThemePalette,
    chart: ResearchGatewayChartData,
) -> None:
    """通用研究默认优先用表格承载多指标，避免把不同量纲硬塞进一根纵轴。"""

    entities = list(dict.fromkeys(point.entity for point in chart.points))[:4]
    metrics = list(dict.fromkeys(point.metric for point in chart.points))[:3]
    table = slide.shapes.add_table(
        len(entities) + 1,
        len(metrics) + 1,
        Inches(1.3),
        Inches(2.78),
        Inches(9.9),
        Inches(min(2.55, 0.56 * (len(entities) + 1))),
    ).table
    table.columns[0].width = Inches(2.05)
    for column in range(1, len(metrics) + 1):
        table.columns[column].width = Inches(7.85 / max(1, len(metrics)))
    headers = ["对象"] + [
        _research_metric_header(metric, chart.points)
        for metric in metrics
    ]
    for column, header in enumerate(headers):
        cell = table.cell(0, column)
        cell.text = _compact_text(header, 50)
        cell.fill.solid()
        cell.fill.fore_color.rgb = palette.primary
        _style_table_cell(cell, palette.surface, bold=True, font_size=10)
    values = {(point.entity, point.metric): point for point in chart.points}
    for row, entity in enumerate(entities, start=1):
        # 第一列是对象名称；后续列才与每个指标一一对应。此前从 0 开始枚举 metrics，导致
        # 对象列被重复写入、真实数值列从未赋值，artifact 虽记录图表但客户看到的是空表。
        entity_cell = table.cell(row, 0)
        entity_cell.text = _compact_text(entity, 50)
        entity_cell.fill.solid()
        entity_cell.fill.fore_color.rgb = palette.surface if row % 2 else palette.background
        _style_table_cell(entity_cell, palette.text, bold=True, font_size=10)
        for column, metric in enumerate(metrics, start=1):
            cell = table.cell(row, column)
            point = values.get((entity, metric))
            cell.text = _compact_text(_format_research_table_value(point, chart), 78)
            cell.fill.solid()
            cell.fill.fore_color.rgb = palette.surface if row % 2 else palette.background
            _style_table_cell(cell, palette.text, bold=False, font_size=10)


def _render_research_trend_table(
    slide: object,
    palette: _StudioThemePalette,
    chart: ResearchGatewayChartData,
) -> None:
    """把折线图背后的逐期序列同时交付为可编辑明细表。"""

    entities = list(dict.fromkeys(point.entity for point in chart.points))[:4]
    periods = list(dict.fromkeys(point.period for point in chart.points))[:8]
    values = {(point.period, point.entity): point for point in chart.points}
    table = slide.shapes.add_table(
        len(periods) + 1,
        len(entities) + 1,
        Inches(1.3),
        Inches(2.72),
        Inches(9.9),
        Inches(min(2.72, 0.48 * (len(periods) + 1))),
    ).table
    table.columns[0].width = Inches(2.1)
    for column in range(1, len(entities) + 1):
        table.columns[column].width = Inches(7.8 / max(1, len(entities)))
    headers = ["期间"] + [_compact_text(entity, 36) for entity in entities]
    for column, header in enumerate(headers):
        cell = table.cell(0, column)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = palette.primary
        _style_table_cell(cell, palette.surface, bold=True, font_size=10)
    for row, period in enumerate(periods, start=1):
        period_cell = table.cell(row, 0)
        period_cell.text = _compact_text(period, 40)
        period_cell.fill.solid()
        period_cell.fill.fore_color.rgb = palette.surface if row % 2 else palette.background
        _style_table_cell(period_cell, palette.text, bold=True, font_size=10)
        for column, entity in enumerate(entities, start=1):
            cell = table.cell(row, column)
            point = values.get((period, entity))
            cell.text = _compact_text(_format_research_table_value(point, chart), 78)
            cell.fill.solid()
            cell.fill.fore_color.rgb = palette.surface if row % 2 else palette.background
            _style_table_cell(cell, palette.text, bold=False, font_size=10)


def _render_research_grouped_bar_chart(
    slide: object,
    palette: _StudioThemePalette,
    chart: ResearchGatewayChartData,
) -> None:
    """分组柱状图使用 PowerPoint 原生 Chart，仅服务已验证同量纲的 2-3 个指标。"""

    entities = list(dict.fromkeys(point.entity for point in chart.points))[:4]
    metrics = list(dict.fromkeys(point.metric for point in chart.points))[:3]
    values = {(point.entity, point.metric): point.value for point in chart.points}
    chart_data = CategoryChartData()
    chart_data.categories = entities
    for metric in metrics:
        chart_data.add_series(metric, [values.get((entity, metric)) for entity in entities])
    native_chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1.3),
        Inches(2.72),
        Inches(9.85),
        Inches(2.52),
        chart_data,
    ).chart
    _style_native_chart(native_chart, palette, show_legend=True)
    _add_chart_value_ledger(slide, chart, palette, top=5.31)


def _style_native_chart(chart: object, palette: _StudioThemePalette, *, show_legend: bool) -> None:
    """为原生 Chart 应用克制的项目视觉 token，不依赖易失的外部模板。"""

    chart.has_title = False
    chart.has_legend = show_legend
    if show_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.name = "Microsoft YaHei UI"
        chart.legend.font.size = Pt(9)
    chart.chart_style = 10
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = palette.border
    chart.value_axis.tick_labels.font.name = "Microsoft YaHei UI"
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.category_axis.tick_labels.font.name = "Microsoft YaHei UI"
    chart.category_axis.tick_labels.font.size = Pt(9)
    colors = (palette.primary, palette.accent, palette.heading, palette.muted)
    for index, series in enumerate(chart.series):
        series.format.line.color.rgb = colors[index % len(colors)]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = colors[index % len(colors)]
    for plot in chart.plots:
        plot.has_data_labels = True
        plot.data_labels.show_value = True
        plot.data_labels.font.name = "Microsoft YaHei UI"
        plot.data_labels.font.size = Pt(8)


def _add_chart_value_ledger(
    slide: object,
    chart: _StructuredDataChart,
    palette: _StudioThemePalette,
    *,
    top: float,
) -> None:
    """在图下保留可回读数据标签；原生 Chart 数据仍可由 PowerPoint 直接编辑。"""

    labels = [
        f"{_point_entity(point)} · {_point_period(point)} · {_format_chart_value(point, chart)}"
        for point in chart.points
    ]
    _add_text(
        slide,
        Inches(1.3),
        Inches(top),
        Inches(9.85),
        Inches(0.26),
        _compact_text(" | ".join(labels), 220),
        7,
        palette.muted,
        False,
    )


def _chart_subtitle(chart: _StructuredDataChart) -> str:
    if isinstance(chart, WorldBankChartData):
        return f"World Bank Indicators API · {chart.indicator_code} · 数据仅按已确认计划读取"
    if chart.evidence_level == "ai_knowledge_draft":
        return "AI 智能生成数据 · 已写入可编辑原生图表"
    source_count = len(chart.sources)
    return f"ResearchGateway · {source_count} 条已读取来源 · 每个数值均带证据与来源 ID"


def _point_entity(point: object) -> str:
    return point.country_name if isinstance(point, WorldBankDataPoint) else point.entity


def _point_period(point: object) -> str:
    return str(point.year) if isinstance(point, WorldBankDataPoint) else _compact_text(point.period, 12)


def _point_metric(point: object, chart: _StructuredDataChart) -> str:
    """统一固定 Provider 与通用研究数据点的原生 Chart 系列名。"""

    return chart.indicator_name if isinstance(point, WorldBankDataPoint) else point.metric


def _format_chart_value(point: object | None, chart: _StructuredDataChart) -> str:
    if point is None:
        return "数据不足"
    if isinstance(chart, WorldBankChartData) and isinstance(point, WorldBankDataPoint):
        return _format_indicator_value(point.value, chart.indicator_code)
    assert isinstance(point, ResearchGatewayDataPoint)
    value = f"{point.value:,.2f}".rstrip("0").rstrip(".")
    return _compact_text(f"{value} {point.unit}", 20)


def _research_metric_header(metric: str, points: tuple[ResearchGatewayDataPoint, ...]) -> str:
    metric_points = tuple(point for point in points if point.metric == metric)
    unit = next((point.unit for point in metric_points), "")
    periods = {point.period for point in metric_points}
    period = next(iter(periods)) if len(periods) == 1 else "各项期间见单元格"
    return " · ".join(part for part in (metric, unit, period) if part)


def _format_research_table_value(
    point: ResearchGatewayDataPoint | None,
    chart: ResearchGatewayChartData,
) -> str:
    """不同来源或时点的演示表必须把差异写进单元格，不能假装同口径。"""

    if point is None:
        return "数据不足"
    value = _format_chart_value(point, chart)
    periods = {item.period for item in chart.points if item.metric == point.metric}
    source_sets = {item.source_ids for item in chart.points if item.metric == point.metric}
    if len(periods) == 1 and len(source_sets) == 1:
        return value
    source_label = "/".join(point.source_ids) or "AI草稿"
    return f"{value} · {point.period} · {source_label}"


def _format_indicator_value(value: float, indicator_code: str) -> str:
    """仅格式化已验证数值，不作单位换算或推导，以免改变数据含义。"""

    if indicator_code.startswith("NY.GDP"):
        if abs(value) >= 1_000_000_000_000:
            return f"${value / 1_000_000_000_000:.1f}T"
        if abs(value) >= 1_000_000_000:
            return f"${value / 1_000_000_000:.1f}B"
        return f"${value:,.0f}"
    if indicator_code == "SP.POP.TOTL":
        if abs(value) >= 100_000_000:
            return f"{value / 100_000_000:.1f}亿"
        if abs(value) >= 10_000:
            return f"{value / 10_000:.0f}万"
    return f"{value:,.0f}"


def _content_image_for_slide(
    plan: PresentationStudioSlidePlan,
    *,
    assets_by_slide_id: dict[str, _StudioVisualAsset],
) -> _StudioVisualAsset | None:
    """封面以外只让绑定到该正文页的素材出现，避免同一图片或错误图片充斥整份演示。"""

    if plan.role != "content":
        return None
    return assets_by_slide_id.get(plan.slide_id)


def _render_agenda(slide: object, plan: PresentationStudioSlidePlan, palette: _StudioThemePalette) -> None:
    columns = 2
    for index, bullet in enumerate(plan.bullets[:6]):
        row, column = divmod(index, columns)
        left = Inches(1.0 + column * 5.55)
        top = Inches(1.86 + row * 1.3)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.1), Inches(0.96))
        card.fill.solid()
        card.fill.fore_color.rgb = palette.surface
        card.line.color.rgb = palette.border
        marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.24), top + Inches(0.25), Inches(0.43), Inches(0.43))
        marker.fill.solid()
        marker.fill.fore_color.rgb = palette.primary
        marker.line.fill.background()
        _add_text(slide, left + Inches(0.78), top + Inches(0.2), Inches(4.02), Inches(0.55), _compact_text(bullet, 52), 13, palette.text, True)
        _add_text(slide, left + Inches(0.35), top + Inches(0.26), Inches(0.2), Inches(0.25), str(index + 1), 8, palette.surface, True)


def _render_image_statement(
    slide: object,
    plan: PresentationStudioSlidePlan,
    palette: _StudioThemePalette,
    *,
    image: _StudioVisualAsset,
) -> None:
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.98), Inches(1.72), Inches(5.12), Inches(4.48))
    panel.fill.solid()
    panel.fill.fore_color.rgb = palette.surface
    panel.line.color.rgb = palette.border
    _add_text(slide, Inches(1.34), Inches(2.03), Inches(4.35), Inches(0.65), _compact_text(plan.bullets[0], 94), 17, palette.heading, True)
    for index, bullet in enumerate(plan.bullets[1:3], start=1):
        top = Inches(3.0 + (index - 1) * 0.88)
        marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.36), top + Inches(0.08), Inches(0.18), Inches(0.18))
        marker.fill.solid()
        marker.fill.fore_color.rgb = palette.accent
        marker.line.fill.background()
        _add_text(slide, Inches(1.68), top, Inches(3.95), Inches(0.55), _compact_text(bullet, 95), 12, palette.text, False)
    _add_image(slide, image, left=Inches(6.55), top=Inches(1.75), width=Inches(5.22), height=Inches(2.73))
    _add_text(slide, Inches(6.55), Inches(4.62), Inches(5.22), Inches(0.25), image.credit_text, 8, palette.muted, False)
    if len(plan.bullets) > 3:
        accent_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.55), Inches(5.04), Inches(5.22), Inches(0.78))
        accent_card.fill.solid()
        accent_card.fill.fore_color.rgb = palette.surface
        accent_card.line.color.rgb = palette.border
        _add_text(slide, Inches(6.83), Inches(5.2), Inches(4.68), Inches(0.35), _compact_text(plan.bullets[3], 105), 11, palette.text, True)


def _render_comparison(slide: object, plan: PresentationStudioSlidePlan, palette: _StudioThemePalette) -> None:
    """双栏只表达用户/模型已给出的两个视角，不擅自添加“现状/目标”等事实判断。"""

    first = plan.bullets[0] if plan.bullets else "待确认的第一视角"
    second = plan.bullets[1] if len(plan.bullets) > 1 else "待确认的第二视角"
    for index, (label, content, color) in enumerate(
        (("视角 A", first, palette.primary), ("视角 B", second, palette.accent))
    ):
        left = Inches(1.02 + index * 5.58)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.88), Inches(5.12), Inches(3.55))
        card.fill.solid()
        card.fill.fore_color.rgb = palette.surface
        card.line.color.rgb = palette.border
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.88), Inches(5.12), Inches(0.18))
        band.fill.solid()
        band.fill.fore_color.rgb = color
        band.line.fill.background()
        _add_text(slide, left + Inches(0.38), Inches(2.28), Inches(4.2), Inches(0.3), label, 11, color, True)
        _add_text(slide, left + Inches(0.38), Inches(2.83), Inches(4.25), Inches(1.25), _compact_text(content, 160), 17, palette.heading, True)
        remaining = plan.bullets[2 + index] if len(plan.bullets) > 2 + index else ""
        if remaining:
            _add_text(slide, left + Inches(0.38), Inches(4.48), Inches(4.2), Inches(0.48), _compact_text(remaining, 100), 11, palette.text, False)


def _render_process(slide: object, plan: PresentationStudioSlidePlan, palette: _StudioThemePalette) -> None:
    """流程页把已有要点按顺序表现，不虚构工期、负责人或完成比例。"""

    bullets = plan.bullets[:5]
    width = Inches(10.7 / max(1, len(bullets)))
    for index, bullet in enumerate(bullets, start=1):
        left = Inches(1.06) + (index - 1) * width
        if index > 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left - Inches(0.18), Inches(3.03), Inches(0.32), Inches(0.42))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = palette.accent
            arrow.line.fill.background()
        marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, Inches(2.06), Inches(0.62), Inches(0.62))
        marker.fill.solid()
        marker.fill.fore_color.rgb = palette.primary if index == 1 else palette.accent
        marker.line.fill.background()
        _add_text(slide, left + Inches(0.19), Inches(2.2), Inches(0.24), Inches(0.22), str(index), 10, palette.surface, True)
        _add_text(slide, left, Inches(3.02), width - Inches(0.28), Inches(1.18), _compact_text(bullet, 100), 13, palette.text, index == 1)


def _render_timeline(slide: object, plan: PresentationStudioSlidePlan, palette: _StudioThemePalette) -> None:
    """时间线展示叙事先后，而非声称真实日期或项目进度。"""

    bullets = plan.bullets[:5]
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.18), Inches(3.32), Inches(10.36), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = palette.border
    line.line.fill.background()
    for index, bullet in enumerate(bullets, start=1):
        left = Inches(1.1 + (index - 1) * (10.35 / max(1, len(bullets) - 1))) if len(bullets) > 1 else Inches(6.0)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, Inches(3.08), Inches(0.5), Inches(0.5))
        dot.fill.solid()
        dot.fill.fore_color.rgb = palette.primary if index == 1 else palette.accent
        dot.line.fill.background()
        label_top = Inches(1.93) if index % 2 else Inches(3.92)
        _add_text(slide, left - Inches(0.3), label_top, Inches(1.1), Inches(0.35), f"阶段 {index}", 10, palette.primary, True)
        _add_text(slide, left - Inches(0.47), label_top + Inches(0.38), Inches(1.48), Inches(0.78), _compact_text(bullet, 62), 11, palette.text, False)


def _render_metric_tiles(slide: object, plan: PresentationStudioSlidePlan, palette: _StudioThemePalette) -> None:
    """没有经核验数值时使用“关键点”卡片，避免把文字硬造为数据图表。"""

    bullets = plan.bullets[:4]
    for index, bullet in enumerate(bullets, start=1):
        row, column = divmod(index - 1, 2)
        left = Inches(1.02 + column * 5.55)
        top = Inches(1.85 + row * 1.85)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.05), Inches(1.45))
        card.fill.solid()
        card.fill.fore_color.rgb = palette.surface
        card.line.color.rgb = palette.border
        _add_text(slide, left + Inches(0.34), top + Inches(0.24), Inches(1.2), Inches(0.28), f"关键点 {index:02d}", 10, palette.primary, True)
        _add_text(slide, left + Inches(0.34), top + Inches(0.63), Inches(4.28), Inches(0.58), _compact_text(bullet, 92), 14, palette.heading, index == 1)


def _render_quote(slide: object, plan: PresentationStudioSlidePlan, palette: _StudioThemePalette) -> None:
    """观点页突出一个主张，并把余下要点降为佐证而不是堆成普通列表。"""

    primary = plan.bullets[0] if plan.bullets else "待确认的核心主张"
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.9), Inches(0.16), Inches(3.72))
    bar.fill.solid()
    bar.fill.fore_color.rgb = palette.accent
    bar.line.fill.background()
    _add_text(slide, Inches(1.55), Inches(2.08), Inches(9.6), Inches(1.42), _compact_text(primary, 190), 24, palette.heading, True)
    for index, bullet in enumerate(plan.bullets[1:4], start=1):
        _add_text(slide, Inches(1.58), Inches(4.05 + (index - 1) * 0.54), Inches(8.9), Inches(0.36), f"— {_compact_text(bullet, 120)}", 12, palette.text, False)


def _render_image_placeholder_statement(
    slide: object,
    plan: PresentationStudioSlidePlan,
    palette: _StudioThemePalette,
) -> None:
    """尚未选图或素材降级时保留图文结构，不把一个空白矩形伪装成图片已生成。"""

    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.55), Inches(1.75), Inches(5.22), Inches(3.15))
    panel.fill.solid()
    panel.fill.fore_color.rgb = palette.accent
    panel.fill.transparency = 72
    panel.line.color.rgb = palette.border
    _add_text(slide, Inches(6.95), Inches(3.05), Inches(4.35), Inches(0.36), "此页可在确认导出时加入已批准的视觉素材", 11, palette.muted, False)
    _render_image_statement_text(slide, plan, palette)


def _render_image_statement_text(slide: object, plan: PresentationStudioSlidePlan, palette: _StudioThemePalette) -> None:
    _add_text(slide, Inches(1.18), Inches(2.08), Inches(4.74), Inches(0.82), _compact_text(plan.bullets[0], 105), 18, palette.heading, True)
    for index, bullet in enumerate(plan.bullets[1:3], start=1):
        _add_text(slide, Inches(1.22), Inches(3.28 + (index - 1) * 0.7), Inches(4.5), Inches(0.48), _compact_text(bullet, 100), 12, palette.text, False)


def _render_insight_cards(slide: object, plan: PresentationStudioSlidePlan, palette: _StudioThemePalette) -> None:
    bullets = plan.bullets[:4]
    card_width = Inches(5.15) if len(bullets) > 2 else Inches(10.82)
    for index, bullet in enumerate(bullets):
        row, column = divmod(index, 2) if len(bullets) > 2 else (index, 0)
        left = Inches(1.0 + column * 5.55)
        top = Inches(1.9 + row * 1.7)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_width, Inches(1.35))
        card.fill.solid()
        card.fill.fore_color.rgb = palette.surface
        card.line.color.rgb = palette.border
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.12), Inches(1.35))
        band.fill.solid()
        band.fill.fore_color.rgb = palette.primary if index == 0 else palette.accent
        band.line.fill.background()
        _add_text(slide, left + Inches(0.38), top + Inches(0.25), card_width - Inches(0.66), Inches(0.76), _compact_text(bullet, 135), 14, palette.text, index == 0)


def _render_delivery_table(slide: object, plan: PresentationStudioSlidePlan, palette: _StudioThemePalette) -> None:
    rows = max(2, min(5, len(plan.bullets) + 1))
    table = slide.shapes.add_table(rows, 2, Inches(1.0), Inches(1.92), Inches(11.05), Inches(3.9)).table
    table.columns[0].width = Inches(2.55)
    table.columns[1].width = Inches(8.5)
    headers = ("交付关注点", "本页表达")
    for column, header in enumerate(headers):
        cell = table.cell(0, column)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = palette.primary
        _style_table_cell(cell, palette.surface, bold=True, font_size=12)
    for row in range(1, rows):
        bullet = plan.bullets[row - 1] if row - 1 < len(plan.bullets) else "待确认的补充信息"
        label = f"重点 {row}"
        for column, value in enumerate((label, _compact_text(bullet, 140))):
            cell = table.cell(row, column)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = palette.surface if row % 2 else palette.background
            _style_table_cell(cell, palette.text, bold=column == 0, font_size=11)


def _render_summary(slide: object, plan: PresentationStudioSlidePlan, palette: _StudioThemePalette) -> None:
    highlight = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.82), Inches(11.05), Inches(1.25))
    highlight.fill.solid()
    highlight.fill.fore_color.rgb = palette.primary
    highlight.line.fill.background()
    _add_text(slide, Inches(1.42), Inches(2.12), Inches(10.18), Inches(0.58), _compact_text(plan.bullets[0], 165), 18, palette.surface, True)
    for index, bullet in enumerate(plan.bullets[1:4], start=1):
        left = Inches(1.0 + (index - 1) * 3.7)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(3.68), Inches(3.35), Inches(1.45))
        card.fill.solid()
        card.fill.fore_color.rgb = palette.surface
        card.line.color.rgb = palette.border
        _add_text(slide, left + Inches(0.26), Inches(3.98), Inches(2.84), Inches(0.72), _compact_text(bullet, 86), 12, palette.text, False)


def _render_sources(
    slide: object,
    plan: PresentationStudioSlidePlan,
    palette: _StudioThemePalette,
    *,
    asset_sources: tuple[_StudioVisualAsset, ...],
    research_sources: tuple[WikimediaResearchSource, ...],
    structured_data: tuple[_StructuredDataChart, ...],
) -> None:
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.98), Inches(1.72), Inches(11.25), Inches(4.65))
    panel.fill.solid()
    panel.fill.fore_color.rgb = palette.surface
    panel.line.color.rgb = palette.border
    entries = list(plan.bullets)
    pexels_sources = [image for image in asset_sources if isinstance(image, PexelsImageAsset)]
    seedream_sources = [image for image in asset_sources if isinstance(image, SeedreamImageAsset)]
    # 计划阶段的第二条说明是“未联网”。一旦本次导出已经经过明确确认并使用外部视觉或公开
    # 资料，就必须替换它，避免来源页同时说“未联网”和展示联网结果这两套相互矛盾的事实。
    if pexels_sources or seedream_sources or research_sources or structured_data:
        entries = entries[:1]
        entries.append("本次外部内容均经确认后获取；仅作为视觉素材或公开出处参考，关键事实仍需人工核验。")
    if pexels_sources:
        entries.append("图片来源：Photos provided by Pexels（已在本页保留摄影师与照片链接）。")
        entries.extend(
            f"{image.photographer or 'Pexels 摄影师'} · {image.photo_url or 'Pexels'}" for image in pexels_sources
        )
    if seedream_sources:
        entries.append("生成式视觉：由 Seedream 5.0 生成，未叠加文字水印，已保留模型审计记录。")
        entries.extend(f"AI 生成图片 · 页面意图：{image.query}" for image in seedream_sources)
    if research_sources:
        entries.append("公开资料参考：以下页面仅补充出处，不自动作为统计数据、结论或引用依据。")
        entries.extend(
            f"{source.title} · {source.page_url} · 抓取于 {source.retrieved_at}"
            for source in research_sources
        )
    research_data_sources: dict[str, object] = {}
    ai_draft_chart_count = 0
    for chart in structured_data:
        if isinstance(chart, WorldBankChartData):
            entries.append(
                "可验证数据图表：World Bank Indicators API；图表仅使用同一年度的比较数据或单一国家的年度趋势。"
            )
            entries.append(
                f"{chart.indicator_name} · {chart.indicator_code} · "
                f"{chart.source_url} · 读取于 {chart.retrieved_at}"
            )
        else:
            if chart.evidence_level == "ai_knowledge_draft":
                ai_draft_chart_count += 1
            for source in chart.sources:
                research_data_sources[source.source_url.casefold()] = source
    if research_data_sources:
        entries.append(
            f"可验证数据章节：ResearchGateway 已生成 {sum(isinstance(item, ResearchGatewayChartData) for item in structured_data)} "
            "个视图，并完成来源读取与逐项证据校验。"
        )
        entries.extend(
            # 来源 ID、域名和短标题先出现，确保长 URL 被视觉层截断时客户仍能识别来源；
            # 完整 URL 和抓取时间继续保存在 artifact，并在版面空间允许时显示在后半段。
            f"{source.source_id} · {urlparse(source.source_url).hostname or source.source_url} · "
            f"{_compact_text(source.title, 48)} · {source.source_url} · 读取于 {source.retrieved_at}"
            for source in research_data_sources.values()
        )
    if ai_draft_chart_count:
        entries.append(
            f"AI 数据草稿：{ai_draft_chart_count} 个数据视图由模型知识生成，未经过公开来源逐项核验；"
            "其用途是快速形成可编辑初稿，正式发布前必须复核。"
        )
    _add_bullets(slide, entries, palette, compact=True)


def _style_table_cell(cell: object, color: RGBColor, *, bold: bool, font_size: int) -> None:
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.font.name = "Microsoft YaHei UI"
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    cell.text_frame.word_wrap = True
    cell.margin_left = Inches(0.12)
    cell.margin_right = Inches(0.12)


def _add_image(
    slide: object,
    image: _StudioVisualAsset,
    *,
    left: object,
    top: object,
    width: object,
    height: object,
) -> None:
    """只接收已验证 Provider 的内存字节，不允许渲染器读取任意本地图片路径。"""

    slide.shapes.add_picture(BytesIO(image.image_bytes), left, top, width=width, height=height)


def _add_text(
    slide: object,
    left: object,
    top: object,
    width: object,
    height: object,
    text: str,
    font_size: int,
    color: RGBColor,
    bold: bool,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = 0
    frame.margin_right = 0
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.LEFT
    paragraph.font.name = "Microsoft YaHei UI"
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color


def _add_bullets(
    slide: object,
    bullets: list[str],
    palette: _StudioThemePalette,
    *,
    compact: bool,
) -> None:
    text_box = slide.shapes.add_textbox(Inches(1.34), Inches(2.06), Inches(10.45), Inches(3.95))
    frame = text_box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.vertical_anchor = MSO_ANCHOR.TOP
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"• {_compact_text(bullet, 180 if compact else 150)}"
        paragraph.font.name = "Microsoft YaHei UI"
        paragraph.font.size = Pt(13 if compact else 16)
        paragraph.font.color.rgb = palette.text
        paragraph.space_after = Pt(7 if compact else 13)
        paragraph.line_spacing = 1.12


def _add_footer(slide: object, source: str, page_label: str, palette: _StudioThemePalette) -> None:
    _add_text(slide, Inches(0.82), Inches(6.9), Inches(10.6), Inches(0.24), source, 9, palette.muted, False)
    _add_text(slide, Inches(11.65), Inches(6.88), Inches(0.75), Inches(0.24), page_label, 9, palette.muted, False)


def _verify_studio_presentation(
    path: Path,
    plan: PresentationStudioPlanResponse,
    *,
    assets: _StudioVisualAssetResolution,
    research: WikimediaResearchResolution,
    structured_data: _StructuredDataResolution,
    motion: NativePresentationMotionSummary,
) -> PresentationVerification:
    if not path.exists() or path.stat().st_size <= 0:
        raise PresentationStudioPlanConflictError("PPT 创作文件未生成有效内容。")
    opened = Presentation(path)
    if len(opened.slides) != len(plan.slides):
        raise PresentationStudioPlanConflictError("PPT 创作文件页数与已确认计划不一致。")
    motion_inspection = inspect_native_presentation_motion(path)
    if motion.enabled:
        if motion_inspection.transition_slide_count != len(plan.slides):
            raise PresentationStudioPlanConflictError("PPT 创作文件缺少原生页面转场，已停止交付。")
        if motion_inspection.entrance_effect_count != motion.entrance_effect_count:
            raise PresentationStudioPlanConflictError("PPT 原生入场动效数量与渲染计划不一致，已停止交付。")
        if motion_inspection.invalid_target_count:
            raise PresentationStudioPlanConflictError("PPT 原生入场动效引用了不存在的页面对象，已停止交付。")
    for index, (slide, planned) in enumerate(zip(opened.slides, plan.slides), start=1):
        text = _read_slide_text(slide)
        if planned.title not in text:
            raise PresentationStudioPlanConflictError(f"第 {index} 页缺少已确认标题，已停止交付。")
        if planned.role == "content" and not any(_visible_bullet_prefix(text, bullet) for bullet in planned.bullets):
            raise PresentationStudioPlanConflictError(f"第 {index} 页缺少已确认正文要点，已停止交付。")
    source_slide_count = sum(1 for item in plan.slides if item.role == "sources")
    if source_slide_count != 1:
        raise PresentationStudioPlanConflictError("PPT 创作计划必须包含唯一事实核验页。")
    if assets.images:
        embedded_picture_total = sum(
            1 for slide in opened.slides for shape in slide.shapes if hasattr(shape, "image")
        )
        if embedded_picture_total < len(assets.images):
            raise PresentationStudioPlanConflictError("PPT 创作文件缺少已确认的外部视觉素材，已停止交付。")
    if research.sources:
        source_text = _read_slide_text(opened.slides[-1])
        if not all(source.title in source_text and source.page_url in source_text for source in research.sources):
            raise PresentationStudioPlanConflictError("PPT 创作文件缺少已确认的公开资料来源，已停止交付。")
    structured_data_charts = _structured_data_charts(structured_data)
    if structured_data_charts:
        source_text = _read_slide_text(opened.slides[-1])
        for chart in structured_data_charts:
            chart_slide = next(
                (
                    slide
                    for planned, slide in zip(plan.slides, opened.slides, strict=True)
                    if planned.slide_id == chart.slide_id
                ),
                None,
            )
            if chart_slide is None:
                raise PresentationStudioPlanConflictError(
                    f"PPT 数据视图缺少目标页面 {chart.slide_id}，请重新生成创作计划。"
                )
            chart_text = _read_slide_text(chart_slide)
            if chart.title not in chart_text:
                raise PresentationStudioPlanConflictError(
                    f"PPT 的 {chart.chart_type} 页面缺少已确认的数据标题，已停止交付。"
                )
            if isinstance(chart, WorldBankChartData):
                source_markers = (chart.source_url, chart.indicator_code)
                missing_sources = [marker for marker in source_markers if marker not in source_text]
            else:
                # 来源页每条记录受版面长度限制，长标题 + 完整 URL 可能在 180 字符处被裁切。
                # 完整 URL 仍保存在 artifact；PPTX 回读用来源 ID、标题前缀和域名三者确认，
                # 既能锁定真实来源，又不会把视觉层的合理截断误判为数据交付失败。
                missing_sources = [
                    source.source_id
                    for source in chart.sources
                    if not _research_source_marker_present(source_text, source)
                ]
            if missing_sources:
                raise PresentationStudioPlanConflictError(
                    f"PPT 来源页缺少 {chart.chart_type} 的来源标记：{', '.join(missing_sources[:3])}。"
                )
            if chart.chart_type in {"comparison_table", "trend_table"}:
                if not any(getattr(shape, "has_table", False) for shape in chart_slide.shapes):
                    raise PresentationStudioPlanConflictError("PPT 数据表页面缺少可编辑的原生表格，已停止交付。")
            elif chart.chart_type in {
                "comparison_bar", "grouped_bar", "horizontal_bar", "trend_line", "trend_area",
                "share_pie", "share_doughnut",
            }:
                if not any(getattr(shape, "has_chart", False) for shape in chart_slide.shapes):
                    raise PresentationStudioPlanConflictError("PPT 图表页面缺少可编辑的原生 Chart，已停止交付。")
            # artifact 中有元数据还不够。表格逐项回读可见单元格；原生 Chart 则直接读取
            # 内嵌 workbook 系列值，避免长图表的数据标签被视觉截短后产生假失败。
            if chart.chart_type in {"comparison_table", "trend_table"}:
                expected_values = [_format_chart_value(point, chart) for point in chart.points]
                if not all(value in chart_text for value in expected_values):
                    raise PresentationStudioPlanConflictError("PPT 数据表缺少已验证数值，已停止交付。")
            elif not _native_chart_values_match(chart_slide, chart):
                raise PresentationStudioPlanConflictError("PPT 原生 Chart 的内嵌数据与验证结果不一致，已停止交付。")
    warnings = [*assets.warnings, *research.warnings, *structured_data.warnings, *motion.warnings]
    return PresentationVerification(
        passed=True,
        slide_count=len(opened.slides),
        source_slide_count=source_slide_count,
        warnings=_bounded_verification_warnings(warnings),
    )


def _research_source_marker_present(source_text: str, source: object) -> bool:
    """验证来源页的稳定可见标记；完整 URL 由 artifact 元数据承担审计保存。"""

    source_id = str(getattr(source, "source_id", "")).strip()
    title = str(getattr(source, "title", "")).strip()
    source_url = str(getattr(source, "source_url", "")).strip()
    hostname = (urlparse(source_url).hostname or "").strip()
    title_prefix = _compact_text(title, 48)
    return bool(
        source_id
        and source_id in source_text
        and title_prefix
        and title_prefix in source_text
        and hostname
        and hostname in source_text
    )


def _bounded_verification_warnings(values: list[str]) -> list[str]:
    """回执只展示有限告警；超出的数量压成摘要，不能让成功产物因 Pydantic 上限变成 400。"""

    unique = list(dict.fromkeys(value.strip() for value in values if value.strip()))
    if len(unique) <= 6:
        return unique
    return [*unique[:5], f"另有 {len(unique) - 5} 条降级或来源说明，完整记录已写入任务历史。"]


def _native_chart_values_match(slide: object, expected_chart: _StructuredDataChart) -> bool:
    """以多重集核对原生 Chart 的内嵌数值，重复值也必须保留正确次数。"""

    actual: Counter[str] = Counter()
    for shape in slide.shapes:
        if not getattr(shape, "has_chart", False):
            continue
        for series in shape.chart.series:
            for value in series.values:
                if value is not None:
                    actual[_chart_number_key(float(value))] += 1
    expected = Counter(_chart_number_key(float(point.value)) for point in expected_chart.points)
    return bool(expected) and all(actual[key] >= count for key, count in expected.items())


def _chart_number_key(value: float) -> str:
    return format(value, ".12g")


def _compact_text(value: str, limit: int) -> str:
    text = re.sub(r"\s+", " ", value).strip()
    return text if len(text) <= limit else f"{text[: max(1, limit - 1)].rstrip()}…"


def _visible_bullet_prefix(rendered_text: str, bullet: str) -> bool:
    """版式可按阅读空间截短正文，但回读必须仍能找到足以识别的受控前缀。

    这不是放宽为“任意文本存在”：标题仍逐字验证，正文只接受原计划去除空白后的前 40 字。
    这样既避免长句挤坏表格/图文页，也能发现渲染时把整条要点丢失的错误。
    """

    normalized = re.sub(r"\s+", " ", bullet).strip()
    if not normalized:
        return True
    prefix = normalized[: min(40, len(normalized))].rstrip()
    return bool(prefix) and prefix in rendered_text


def _read_slide_text(slide: object) -> str:
    """读取文本框和表格单元格，供交付回读验证使用。

    `python-pptx` 中表格 Shape 不带 ``text`` 属性；如果只遍历 ``has_text_frame``，新版的
    表格页会被错误判成空白。所有版式都通过这一个入口回读，避免验证规则跟随 UI 分叉。
    """

    chunks: list[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            chunks.append(shape.text)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                chunks.extend(cell.text for cell in row.cells)
    return "\n".join(chunks)
