"""使用本机已配置模型的 PPT 数据链路真实验收。

默认不运行，必须显式传入 ``--live``。脚本不读取、打印或保存 API Key；它只输出脱敏的
供应商名称、已规划的图表类型、数据点数量和 Seedream 图片结果。用于修改模型规划、AI
数据直出或图像 Provider 后的人工授权回归，避免用 mock 误判真实供应商兼容性。
"""

from __future__ import annotations

import argparse
import asyncio
import sys
import tempfile
from pathlib import Path

from pptx import Presentation


BACKEND_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(BACKEND_ROOT))

from app.schemas.presentation_studio import PresentationStudioPlanRequest
from app.services.model_gateway import resolve_model_runtime
from app.services.presentation_research_gateway import fetch_ai_knowledge_draft_chart_data
from app.services.presentation_studio import build_presentation_studio_plan
from app.services.presentation_studio_delivery import _render_studio_presentation
from app.services.seedream_assets import generate_seedream_images


def main() -> None:
    """执行一次真实模型规划、数据生成和单张图片生成探针。"""

    parser = argparse.ArgumentParser(description="运行真实 PPT 数据与图片回归。")
    parser.add_argument("--live", action="store_true", help="确认使用本机已配置的真实模型与图像额度。")
    parser.add_argument("--render", action="store_true", help="将真实数据写入临时 PPTX 并回读原生对象。")
    args = parser.parse_args()
    if not args.live:
        raise SystemExit("此脚本会消耗真实 Provider 额度；请显式传入 --live。")

    runtime = resolve_model_runtime()
    request = PresentationStudioPlanRequest(
        # Unicode escape 避免 Windows 控制台管道把中文验收语句转成问号。
        intent=(
            "\u5e2e\u6211\u751f\u6210\u6885\u897f\u751f\u6daf\u6570\u636ePPT\uff0c"
            "\u8981\u5305\u542b\u591a\u79cd\u6570\u636e\u3002"
        ),
        target_slide_count=8,
        visual_asset_provider="seedream",
        structured_data_enabled=True,
    )
    plan = asyncio.run(build_presentation_studio_plan(request=request))
    expected_visuals = list(plan.data_plan.requested_visuals)
    print(
        "live plan: "
        f"provider={runtime.provider} model={runtime.model} "
        f"state={plan.data_plan.state} entities={plan.data_plan.entities} visuals={plan.data_plan.requested_visuals} "
        f"contract={plan.data_plan.required_visual_count}"
    )
    if (
        plan.data_plan.state != "research_planned"
        or plan.data_plan.entities != ["梅西"]
        or len(expected_visuals) < 4
        or not any(visual.endswith("table") for visual in expected_visuals)
    ):
        raise RuntimeError(f"简短单对象主题没有形成受限的多视图数据计划：{plan.data_plan.notice}")

    resolution = fetch_ai_knowledge_draft_chart_data(plan.data_plan)
    delivered_visuals = [chart.chart_type for chart in resolution.charts]
    point_counts = [len(chart.points) for chart in resolution.charts]
    print(f"live data: visuals={delivered_visuals} points={point_counts} warnings={len(resolution.warnings)}")
    if resolution.warnings:
        print("live data warning: " + " | ".join(resolution.warnings))
    if delivered_visuals != expected_visuals:
        raise RuntimeError("真实模型没有完成所有已请求的数据视图。")

    if args.render:
        # 只把真实模型数据写入系统临时目录。该验证不创建客户 artifact，也不消耗额外模型或图片额度。
        with tempfile.TemporaryDirectory(prefix="agentflow_live_presentation_") as temporary_directory:
            target = Path(temporary_directory) / "live_data_contract.pptx"
            with target.open("wb") as target_file:
                _render_studio_presentation(
                    target_file,
                    plan,
                    assets=(),
                    assets_by_slide_id={},
                    research_sources=(),
                    structured_data=resolution.charts,
                )
            presentation = Presentation(target)
            table_count = sum(
                1
                for slide in presentation.slides
                for shape in slide.shapes
                if getattr(shape, "has_table", False)
            )
            chart_count = sum(
                1
                for slide in presentation.slides
                for shape in slide.shapes
                if getattr(shape, "has_chart", False)
            )
            print(f"live render: slides={len(presentation.slides)} tables={table_count} charts={chart_count}")
            expected_table_count = sum(visual.endswith("table") for visual in expected_visuals)
            expected_chart_count = len(expected_visuals) - expected_table_count
            if table_count < expected_table_count or chart_count < expected_chart_count:
                raise RuntimeError("真实数据未被写入足量的 PowerPoint 原生表格或图表。")

    # 图片提示词用英文是 Provider 的最佳实践，与客户中文意图无关，也便于确认返回的是图像。
    images = generate_seedream_images(
        ["Lionel Messi football career data story, editorial stadium portrait, blue gold presentation visual"],
        limit=1,
    )
    print(f"live Seedream: images={len(images.images)} warnings={len(images.warnings)}")
    if images.warnings:
        print("live Seedream warning: " + " | ".join(images.warnings))


if __name__ == "__main__":
    main()
