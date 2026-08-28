"""离线验证 PPT 制作 V3 的“自然语言意图 -> 创作计划 -> 多版式交付 -> 任务历史”闭环。

脚本强制 mock 和临时数据目录：它不读取真实模型 Key、不联网，也不向项目 output/ 写入
PPTX 或素材。实际文件交付与授权素材会在后续确认阶段单独接入并扩大验证覆盖。
"""

from __future__ import annotations

import atexit
import asyncio
import json
import os
import shutil
import sys
import tempfile
import time
import warnings
from dataclasses import replace
from pathlib import Path

import httpx

warnings.filterwarnings(
    "ignore",
    message="Using `httpx` with starlette.testclient is deprecated.*",
)

BACKEND_ROOT = Path(__file__).resolve().parents[1]
_VERIFY_ROOT = Path(tempfile.mkdtemp(prefix="agentflow_verify_presentation_studio_"))
os.environ["AGENTFLOW_DATA_DIR"] = str(_VERIFY_ROOT / "data")
os.environ["AGENTFLOW_DOCUMENT_PRESENTATION_OUTPUT_DIR"] = str(_VERIFY_ROOT / "presentations")
os.environ["AGENTFLOW_CHAT_MODE"] = "mock"
# 离线回归必须稳定降级，不能因开发者本机配置的图库 Key 触发真实联网请求。
os.environ["AGENTFLOW_PEXELS_API_KEY"] = ""
os.environ["AGENTFLOW_SEEDREAM_API_KEY"] = ""
# 离线回归不使用开发机的稳定搜索 Key；Tavily adapter 本身另以 MockTransport 覆盖，避免
# PPT 全量验证因本机 .env、代理或外部额度发生网络副作用。
os.environ["AGENTFLOW_TAVILY_API_KEY"] = ""
os.environ["AGENTFLOW_PRESENTATION_RESEARCH_SEARCH_PROVIDER"] = "deepseek_native"
atexit.register(lambda: shutil.rmtree(_VERIFY_ROOT, ignore_errors=True))
sys.path.insert(0, str(BACKEND_ROOT))

from fastapi.testclient import TestClient
from pptx import Presentation

from main import app
from app.schemas.presentation_studio import (
    PresentationStudioBrief,
    PresentationStudioDataPlan,
    PresentationStudioPlanRequest,
    PresentationStudioPlanResponse,
    PresentationStudioSlidePlan,
)
from app.services.presentation_studio_delivery import (
    _bounded_verification_warnings,
    _effective_asset_slots,
    _research_source_marker_present,
    _structured_data_contract_gap,
)
from app.services import presentation_studio_delivery as delivery_module
from app.services.presentation_native_animations import inspect_native_presentation_motion
from app.services.presentation_studio import (
    _StudioContentSlide,
    _StudioModelOutput,
    _StudioResearchBlueprint,
    _align_research_time_scope,
    _data_plan,
    _data_visual_intent,
    _ensure_pair_comparison_query,
    _has_explicit_data_research_intent,
    _has_world_bank_shortcut,
    _focus_broad_comparison_blueprint,
    _ensure_data_visual_queries,
    _infer_conservative_research_blueprint,
    _parse_model_output,
    _parse_research_blueprint,
    _request_research_blueprint,
    _requested_entity_scope,
    _requested_data_visual_count_hint,
    _strip_unverified_numeric_claims,
    _trend_query_seeds,
)
from app.services.seedream_assets import SeedreamImageAsset, _generate_one_image
from app.services.model_gateway import (
    ModelGatewayError,
    ModelToolTurn,
    NativeWebSearchResult,
    NativeWebSearchSource,
)
from app.services.presentation_research_gateway import (
    ResearchGatewayChartData,
    ResearchGatewayDataPoint,
    ResearchGatewayResolution,
    ResearchGatewaySource,
    ResearchGatewayValidationError,
    _ResearchSearchCandidate,
    _extract_from_sources,
    _from_tavily_candidate,
    _locate_evidence_quote,
    _read_sources,
    _prioritize_candidates,
    _search_research_sources,
    _split_research_queries,
    _trend_period_label,
    _validate_extraction,
    _validate_extraction_views,
    complete_research_resolution_with_ai_draft,
    fetch_ai_knowledge_draft_chart_data,
    fetch_research_gateway_chart_data,
)
from app.services.presentation_research_network import research_httpx_options
from app.services.tavily_research import fetch_tavily_research_sources
from app.services.tavily_research import TavilyResearchCandidate
from app.services.wikimedia_research import fetch_wikimedia_references, is_safe_wikimedia_url
from app.services.world_bank_data import fetch_world_bank_chart_data


def main() -> None:
    """覆盖同步计划、异步受理、历史恢复与“计划阶段零副作用”。"""

    client = TestClient(app)

    # Seedream 单页素材请求必须显式关闭组图，避免每个 PPT 页面槽位被 Provider 当作组图
    # 任务排队。这里用 MockTransport 验证请求契约，不读取本机 Key 或发起真实图片生成。
    def seedream_handler(request: httpx.Request) -> httpx.Response:
        assert request.method == "POST"
        assert request.url.path.endswith("/images/generations")
        payload = json.loads(request.content)
        assert payload["sequential_image_generation"] == "disabled"
        assert payload["response_format"] == "b64_json"
        assert payload["watermark"] is False
        return httpx.Response(200, json={"data": [{"id": "seedream-fixture", "b64_json": "/9j/"}]})

    with httpx.Client(transport=httpx.MockTransport(seedream_handler)) as mock_client:
        seedream_asset = _generate_one_image(mock_client, query="single image request fixture")
    assert seedream_asset.asset_id == "seedream-fixture"
    assert seedream_asset.image_bytes.startswith(b"\xff\xd8\xff")

    assert _requested_data_visual_count_hint(
        PresentationStudioPlanRequest(
            intent="多生成几张数据表格，并包含柱状图和折线图。",
            structured_data_enabled=True,
        )
    ) == 4
    assert _requested_data_visual_count_hint(
        PresentationStudioPlanRequest(
            intent="帮我生成梅西生涯数据 PPT。",
            structured_data_enabled=True,
        )
    ) == 4

    # 稳定搜索 Adapter 只接收已固化的查询并请求清洗正文；不能消费 Tavily 的自然语言 answer。
    # Advanced Search 的 content 是来源页按查询重排的原文片段，应置于整页正文前避免关键数值被截断。
    def tavily_handler(request: httpx.Request) -> httpx.Response:
        assert request.method == "POST"
        assert request.url == httpx.URL("https://api.tavily.com/search")
        assert request.headers["authorization"] == "Bearer tavily-test-key"
        payload = json.loads(request.content)
        assert payload["query"] == "甲队乙队 2024 采用率 官方统计"
        assert payload["include_answer"] is False
        assert payload["include_raw_content"] == "text"
        return httpx.Response(
            200,
            json={
                "answer": "这段供应商答案不能进入数据图表。",
                "results": [
                    {
                        "title": "甲队官方统计",
                        "url": "https://stats.example.com/team-a",
                        "content": "甲队相关片段：甲队采用率为 42%。",
                        "raw_content": "甲队采用率为 42%，统计时间为 2024 年。官方页面同时公开了统计口径。",
                    },
                    {
                        "title": "乙队官方统计",
                        "url": "https://stats.example.com/team-b",
                        "content": "乙队相关片段：乙队采用率为 38%，统计时间为 2024 年。",
                        "raw_content": None,
                    },
                ],
            },
        )

    with httpx.Client(transport=httpx.MockTransport(tavily_handler)) as mock_client:
        tavily_sample = fetch_tavily_research_sources(
            ["甲队乙队 2024 采用率 官方统计"],
            api_key="tavily-test-key",
            client=mock_client,
        )
    assert tavily_sample.query_count == 1
    assert len(tavily_sample.candidates) == 2
    assert tavily_sample.candidates[0].raw_content.startswith("甲队相关片段")
    assert "PAGE TEXT:" in tavily_sample.candidates[0].raw_content
    assert "甲队采用率为 42%" in tavily_sample.candidates[0].raw_content
    assert tavily_sample.candidates[1].raw_content.startswith("乙队相关片段")
    assert all("供应商答案" not in item.raw_content for item in tavily_sample.candidates)

    # 多对象研究会并发执行三条聚焦查询，但聚合顺序仍须与蓝图一致，不能因线程完成顺序让
    # 候选来源和对象映射随机变化。每条查询只返回一条，便于直接验证 query -> candidate 关系。
    concurrent_queries = ("双方对比", "甲队单独", "乙队单独")

    def concurrent_tavily_handler(request: httpx.Request) -> httpx.Response:
        payload = json.loads(request.content)
        query = payload["query"]
        return httpx.Response(
            200,
            json={
                "results": [
                    {
                        "title": f"{query}来源",
                        "url": f"https://stats.example.com/{concurrent_queries.index(query)}",
                        "content": f"{query}的查询相关原文片段包含数值 42。",
                        "raw_content": f"{query}的完整正文。",
                    }
                ]
            },
        )

    with httpx.Client(transport=httpx.MockTransport(concurrent_tavily_handler)) as mock_client:
        concurrent_sample = fetch_tavily_research_sources(
            concurrent_queries,
            api_key="tavily-test-key",
            client=mock_client,
        )
    assert [item.source_query for item in concurrent_sample.candidates] == list(concurrent_queries)
    assert [item.title for item in concurrent_sample.candidates] == [f"{query}来源" for query in concurrent_queries]

    # 同一趋势查询返回多个页面时，应优先读取有多个期间和密集数字的统计表，不能让视频或
    # 社交页面占用每组唯一的首轮读取名额。排序只决定读取顺序，不替代后续逐点证据校验。
    trend_query = "甲队 season appearances goals every season statistics table"
    prioritized_trend = _prioritize_candidates(
        (
            _ResearchSearchCandidate(
                title="甲队单赛季视频",
                url="https://www.youtube.com/watch?v=example",
                prefetched_excerpt="甲队 2024 年进球回顾，共 20 球。",
                source_query=trend_query,
            ),
            _ResearchSearchCandidate(
                title="甲队历年统计表",
                url="https://stats.example.com/team-a-seasons",
                prefetched_excerpt=(
                    "Season Apps Goals 2021/22 30 18 2022/23 32 21 "
                    "2023/24 35 24 2024/25 36 27"
                ),
                source_query=trend_query,
            ),
        )
    )
    assert prioritized_trend[0].title == "甲队历年统计表"

    # `auto` 只在配置稳定搜索 Key 后选择 Tavily；这能保证现有只配置 DeepSeek 的用户继续走
    # 原生搜索。这里替换网关模块的依赖，验证路由判断本身，不触发任何真实网络或 Key 读取。
    import app.services.presentation_research_gateway as research_gateway_module

    original_gateway_settings = research_gateway_module.settings
    original_tavily_fetch = research_gateway_module.fetch_tavily_research_sources
    tavily_route_queries: list[tuple[str, ...]] = []

    def fake_tavily_fetch(queries: tuple[str, ...]) -> object:
        tavily_route_queries.append(queries)
        return type(
            "FakeTavilyResult",
            (),
            {
                "candidates": (
                    TavilyResearchCandidate(
                        title="甲队官方统计",
                        url="https://stats.example.com/team-a",
                        raw_content="甲队采用率为 42%，统计时间为 2024 年。官方口径已公开。",
                    ),
                ),
                "query_count": len(queries),
                "warnings": (),
            },
        )()

    research_gateway_module.settings = replace(
        original_gateway_settings,
        presentation_research_search_provider="auto",
        tavily_api_key="tavily-test-key",
    )
    research_gateway_module.fetch_tavily_research_sources = fake_tavily_fetch
    try:
        routed_search = asyncio.run(
            _search_research_sources(
                # Tavily 路径不应访问聊天模型 Runtime；传入无方法对象可防止测试误把该分支
                # 退回到 DeepSeek 原生搜索。
                runtime=object(),  # type: ignore[arg-type]
                queries=("甲队乙队 2024 采用率 官方统计",),
            )
        )
    finally:
        research_gateway_module.settings = original_gateway_settings
        research_gateway_module.fetch_tavily_research_sources = original_tavily_fetch
    assert tavily_route_queries == [("甲队乙队 2024 采用率 官方统计",)]
    assert routed_search.candidates[0].retrieval_method == "tavily_raw_content"

    # 自动模式中稳定搜索若临时返回零候选，才允许使用研究蓝图尚未消费的查询回退到 DeepSeek。
    # 验证它不会重复第一组查询，也不会因为回退而把剩余额度再次留给后续补查。
    native_fallback_queries: list[tuple[str, ...]] = []

    class FallbackRuntime:
        async def native_web_search_sources(
            self,
            *,
            queries: tuple[str, ...],
            max_uses: int,
        ) -> NativeWebSearchResult:
            assert max_uses == len(queries)
            native_fallback_queries.append(queries)
            return NativeWebSearchResult(
                sources=(
                    NativeWebSearchSource(
                        title="乙队官方统计",
                        url="https://stats.example.com/team-b",
                    ),
                ),
                query_count=len(queries),
                retrieved_at="2026-08-12T00:00:00+00:00",
            )

    def fake_empty_tavily_fetch(queries: tuple[str, ...]) -> object:
        tavily_route_queries.append(queries)
        return type("FakeEmptyTavilyResult", (), {"candidates": (), "query_count": len(queries), "warnings": ()})()

    research_gateway_module.settings = replace(
        original_gateway_settings,
        presentation_research_search_provider="auto",
        tavily_api_key="tavily-test-key",
    )
    research_gateway_module.fetch_tavily_research_sources = fake_empty_tavily_fetch
    try:
        fallback_search = asyncio.run(
            _search_research_sources(
                runtime=FallbackRuntime(),  # type: ignore[arg-type]
                queries=("第一组查询",),
                fallback_queries=("第二组查询", "第三组查询"),
            )
        )
    finally:
        research_gateway_module.settings = original_gateway_settings
        research_gateway_module.fetch_tavily_research_sources = original_tavily_fetch
    assert native_fallback_queries == [("第二组查询", "第三组查询")]
    assert fallback_search.provider == "deepseek_native_fallback"
    assert fallback_search.query_count == 3
    assert fallback_search.fallback_used is True
    assert fallback_search.candidates[0].url == "https://stats.example.com/team-b"

    # 浏览器代理不一定传给 Python。专用代理或直连都必须由客户显式选择，并且不能被写入
    # artifact；这里只验证 httpx 参数，不发起网络请求。
    import app.services.presentation_research_network as research_network_module

    original_network_settings = research_network_module.settings
    research_network_module.settings = replace(
        original_network_settings,
        presentation_research_network_mode="environment",
        presentation_research_proxy_url="http://127.0.0.1:7890",
    )
    try:
        assert research_httpx_options() == {
            "trust_env": False,
            "proxy": "http://127.0.0.1:7890",
        }
        research_network_module.settings = replace(
            original_network_settings,
            presentation_research_network_mode="direct",
            presentation_research_proxy_url="http://127.0.0.1:7890",
        )
        assert research_httpx_options() == {"trust_env": False}
    finally:
        research_network_module.settings = original_network_settings
    # 客户仅提供一句主题时，默认计划应保有完整叙事而非早期的八页简版；内置主题本身不是
    # 外部素材来源，因此不选择配图时仍应明确说明其会继续生效。
    default_plan = client.post(
        "/api/agents/document_agent/presentation-studio/run",
        json={"intent": "为客户介绍一项可确认的产品升级方案。"},
    )
    assert default_plan.status_code == 200, default_plan.text
    default_payload = default_plan.json()
    assert len(default_payload["slides"]) == 10
    assert default_payload["asset_plan"]["state"] == "not_requested"
    assert "内置主题" in default_payload["asset_plan"]["notice"]
    assert default_payload["research_plan"]["state"] == "not_requested"

    # 公开资料 Provider 只允许固定的中文 Wikimedia API。这里使用 MockTransport 覆盖解析、
    # HTML 清理、链接边界和去重，不让离线回归因为开发机网络状态而变得不稳定。
    def wikimedia_handler(request: httpx.Request) -> httpx.Response:
        assert request.url.host == "zh.wikipedia.org"
        assert request.url.path == "/w/api.php"
        assert request.url.params["action"] == "query"
        return httpx.Response(
            200,
            json={
                "query": {
                    "search": [
                        {"pageid": 42, "title": "公开资料样例", "snippet": "<span>可追溯</span> 的资料参考"},
                        {"pageid": 43, "title": "公开资料样例", "snippet": "重复结果"},
                    ]
                }
            },
        )

    with httpx.Client(transport=httpx.MockTransport(wikimedia_handler)) as mock_client:
        research_sample = fetch_wikimedia_references(
            ["创作主题", "另一个主题"],
            client=mock_client,
        )
    assert len(research_sample.sources) == 1
    assert research_sample.sources[0].snippet == "可追溯 的资料参考"
    assert is_safe_wikimedia_url(research_sample.sources[0].page_url)

    # 结构化数据 Provider 只接受固定 World Bank 指标计划。比较图必须选择所有国家都存在的
    # 最新共同年份，而不是各取最新值后伪装为可比数据。
    def world_bank_handler(request: httpx.Request) -> httpx.Response:
        assert request.url.host == "api.worldbank.org"
        assert request.url.path.endswith("/indicator/NY.GDP.MKTP.CD")
        assert request.url.params["format"] == "json"
        return httpx.Response(
            200,
            json=[
                {"page": 1, "pages": 1, "per_page": "100", "total": 4},
                [
                    {"countryiso3code": "CHN", "country": {"value": "中国"}, "indicator": {"id": "NY.GDP.MKTP.CD"}, "date": "2024", "value": 18_000_000_000_000},
                    {"countryiso3code": "USA", "country": {"value": "美国"}, "indicator": {"id": "NY.GDP.MKTP.CD"}, "date": "2024", "value": 28_000_000_000_000},
                    {"countryiso3code": "CHN", "country": {"value": "中国"}, "indicator": {"id": "NY.GDP.MKTP.CD"}, "date": "2023", "value": 17_000_000_000_000},
                    {"countryiso3code": "USA", "country": {"value": "美国"}, "indicator": {"id": "NY.GDP.MKTP.CD"}, "date": "2023", "value": 27_000_000_000_000},
                ],
            ],
        )

    data_sample_plan = PresentationStudioDataPlan(
        state="provider_planned",
        provider="world_bank",
        chart_type="comparison_bar",
        slide_id="content_1",
        indicator_code="NY.GDP.MKTP.CD",
        indicator_name="GDP（现价美元）",
        country_codes=["CHN", "USA"],
        country_names=["中国", "美国"],
        max_points=2,
        notice="验证数据图表计划。",
    )
    with httpx.Client(transport=httpx.MockTransport(world_bank_handler)) as mock_client:
        data_sample = fetch_world_bank_chart_data(data_sample_plan, client=mock_client)
    assert data_sample.chart is not None
    assert data_sample.chart.chart_type == "comparison_bar"
    assert {point.year for point in data_sample.chart.points} == {2024}
    assert len(_bounded_verification_warnings([f"告警 {index}" for index in range(8)])) == 6
    assert "另有 3 条" in _bounded_verification_warnings([f"告警 {index}" for index in range(8)])[-1]

    # 通用主题不需要为每个领域增加专用 MCP。首次模型只产出二次研究蓝图，计划阶段保持
    # 零联网、零事实数值；后续 ResearchGateway 再在客户确认后执行查询和验证。
    generic_plan = _data_plan(
        PresentationStudioPlanRequest(
            intent="制作梅西与 C 罗职业数据对比 PPT。",
            structured_data_enabled=True,
        ),
        slides=[
            PresentationStudioSlidePlan(
                slide_id="content_1",
                role="content",
                title="职业数据对比",
                bullets=["比较核心数据。", "统一赛事与时间口径。"],
                visual_direction="用可验证数据表呈现差异。",
            )
        ],
        brief=PresentationStudioBrief(
            title="梅西与 C 罗：职业数据对比",
            purpose="用统一口径比较两位球员。",
            audience="足球爱好者",
            core_message="数据只有在来源、时间和赛事口径一致时才可比较。",
            theme="impact_contrast",
            theme_reason="对比主题需要清晰区分对象。",
            fact_check_notice="所有数据必须在联网确认后核验。",
        ),
        blueprint=_StudioResearchBlueprint(
            needed=True,
            research_question="梅西与 C 罗在统一赛事口径下的职业进球与助攻如何比较？",
            entities=["梅西", "C 罗"],
            metrics=["职业进球", "职业助攻"],
            time_scope="截至可获得的最新完整赛季",
            comparison_scope="同一数据来源、赛事范围和统计时间",
            chart_type="grouped_bar",
            target_slide_index=1,
            required_data_points=4,
            search_queries=[
                "Messi Ronaldo career goals assists official statistics",
                "Lionel Messi career statistics official profile",
                "Cristiano Ronaldo career statistics official profile",
            ],
            preferred_source_types=["official_profile", "recognized_database"],
        ),
    )
    assert generic_plan.state == "research_planned"
    assert generic_plan.provider == "research_gateway"
    assert generic_plan.entities == ["梅西", "C 罗"]
    assert generic_plan.chart_type == "comparison_table"
    # 该单元夹具只有一张正文页，因此第二个视图不能与表格重叠；完整计划会按正文页扩展。
    assert generic_plan.requested_visuals == ["comparison_table"]
    assert 4 <= len(generic_plan.search_queries) <= 6
    assert _has_explicit_data_research_intent("制作梅西与 C 罗的数据对比 PPT")
    assert not _has_explicit_data_research_intent("制作一份温暖的团队文化介绍 PPT")
    assert _has_world_bank_shortcut("制作中国与美国 GDP 对比 PPT")
    assert not _has_world_bank_shortcut("制作梅西与 C 罗的数据对比 PPT")

    # 简短的数据主题由规划模型决定指标与视图组合，Harness 只补足基础层级和页面上限；
    # 不应继续因为客户没有逐字说“柱状图/折线图”就退化成单表。
    adaptive_slides = [
        PresentationStudioSlidePlan(
            slide_id=f"content_{index}",
            role="content",
            title=f"生涯数据视图 {index}",
            bullets=["用合适的数据形态回答不同问题。"],
            visual_direction="按数据用途选择图表。",
        )
        for index in range(1, 5)
    ]
    adaptive_plan = _data_plan(
        PresentationStudioPlanRequest(
            intent="帮我生成梅西生涯数据 PPT。",
            structured_data_enabled=True,
        ),
        slides=adaptive_slides,
        brief=PresentationStudioBrief(
            title="梅西生涯数据全景",
            purpose="用多种数据形态呈现职业生涯。",
            audience="足球爱好者",
            core_message="从总量、结构和趋势理解职业表现。",
            theme="impact_contrast",
            theme_reason="数据主题需要清晰层次。",
            fact_check_notice="AI 草稿需在正式发布前复核。",
        ),
        blueprint=_StudioResearchBlueprint(
            needed=True,
            research_question="梅西职业生涯的总量指标和逐赛季趋势如何？",
            entities=["梅西"],
            entity_search_names=["Lionel Messi 梅西"],
            metrics=["职业生涯出场次数", "职业生涯总进球数", "职业生涯助攻数"],
            trend_metric="逐赛季联赛进球数",
            time_scope="职业生涯",
            comparison_scope="统一赛事范围和统计口径",
            chart_type="comparison_table",
            recommended_visuals=["comparison_table", "horizontal_bar", "trend_line"],
            visual_metrics=[
                ["职业生涯出场次数", "职业生涯总进球数", "职业生涯助攻数"],
                ["职业生涯出场次数", "职业生涯总进球数", "职业生涯助攻数"],
                ["逐赛季联赛进球数"],
            ],
            target_slide_index=1,
            required_data_points=12,
            search_queries=[
                "Lionel Messi career appearances goals assists statistics",
                "Lionel Messi season league goals table",
                "梅西 生涯 出场 进球 助攻 数据",
            ],
        ),
    )
    assert adaptive_plan.requested_visuals == [
        "comparison_table", "horizontal_bar", "trend_line", "trend_area"
    ]
    assert len(adaptive_plan.visual_slide_ids) == 4
    assert adaptive_plan.visual_metrics[1] == [
        "职业生涯出场次数", "职业生涯总进球数", "职业生涯助攻数"
    ]

    # 客户只点名“梅西生涯数据”时，规划模型即使错误联想出 C 罗，也不能把第二个人带进
    # 数据模型、图表或内部查询。该回归直接模拟真实客户反馈的跑偏蓝图。
    assert _requested_entity_scope("帮我生成梅西生涯数据 PPT，要包含多种数据。") == ["梅西"]
    assert _requested_entity_scope("制作梅西与 C 罗的各种数据对比 PPT。") == ["梅西", "C 罗"]
    entity_drift_plan = _data_plan(
        PresentationStudioPlanRequest(
            intent="帮我生成梅西生涯数据 PPT，要包含多种数据。",
            structured_data_enabled=True,
        ),
        slides=adaptive_slides,
        brief=PresentationStudioBrief(
            title="梅西生涯数据全景",
            purpose="用多种数据形态呈现职业生涯。",
            audience="足球爱好者",
            core_message="从总量、结构和趋势理解职业表现。",
            theme="impact_contrast",
            theme_reason="数据主题需要清晰层次。",
            fact_check_notice="AI 草稿需在正式发布前复核。",
        ),
        blueprint=_StudioResearchBlueprint(
            needed=True,
            research_question="比较梅西与 C 罗的职业生涯数据。",
            entities=["梅西", "C 罗"],
            entity_search_names=["Lionel Messi 梅西", "Cristiano Ronaldo C 罗"],
            metrics=["职业生涯出场次数", "职业生涯总进球数", "职业生涯助攻数"],
            trend_metric="逐赛季联赛进球数",
            time_scope="职业生涯",
            comparison_scope="梅西与 C 罗的对比范围",
            chart_type="comparison_table",
            recommended_visuals=["comparison_table", "grouped_bar", "trend_line"],
            target_slide_index=1,
            required_data_points=12,
            search_queries=["梅西 C 罗数据", "梅西数据", "C 罗数据"],
        ),
    )
    assert entity_drift_plan.entities == ["梅西"]
    assert entity_drift_plan.entity_search_names == ["梅西"]
    assert "C 罗" not in entity_drift_plan.research_question
    assert "C 罗" not in entity_drift_plan.comparison_scope
    assert all("C 罗" not in query for query in entity_drift_plan.search_queries)
    assert entity_drift_plan.requested_visuals == [
        "comparison_table", "horizontal_bar", "trend_line", "trend_area"
    ]

    class AiDraftRuntime:
        """固定返回单对象总量与趋势，验证模型知识回退不依赖联网来源。"""

        async def tool_turn(self, **_: object) -> ModelToolTurn:
            return ModelToolTurn(
                content=json.dumps(
                    {
                        "status": "complete",
                        "title": "梅西生涯数据草稿",
                        "points": [
                            {"entity": "梅西", "metric": "职业生涯出场次数", "value": 1100, "unit": "场", "period": "职业生涯"},
                            {"entity": "梅西", "metric": "职业生涯总进球数", "value": 850, "unit": "个", "period": "职业生涯"},
                            {"entity": "梅西", "metric": "职业生涯助攻数", "value": 380, "unit": "次", "period": "职业生涯"},
                            {"entity": "梅西", "metric": "逐赛季联赛进球数", "value": 25, "unit": "球", "period": "2019/20"},
                            {"entity": "梅西", "metric": "逐赛季联赛进球数", "value": 30, "unit": "球", "period": "2020/21"},
                            {"entity": "梅西", "metric": "逐赛季联赛进球数", "value": 11, "unit": "球", "period": "2021/22"},
                            {"entity": "梅西", "metric": "逐赛季联赛进球数", "value": 16, "unit": "球", "period": "2022/23"},
                        ],
                        "notes": ["示例数值仅用于离线验证。"],
                    },
                    ensure_ascii=False,
                )
            )

    ai_completed = complete_research_resolution_with_ai_draft(
        adaptive_plan,
        ResearchGatewayResolution(chart=None, warnings=("公开来源不足。",)),
        runtime=AiDraftRuntime(),  # type: ignore[arg-type]
    )
    assert len(ai_completed.charts) == 4
    assert [chart.chart_type for chart in ai_completed.charts] == adaptive_plan.requested_visuals
    assert all(chart.evidence_level == "ai_knowledge_draft" for chart in ai_completed.charts)
    assert all(not chart.sources for chart in ai_completed.charts)
    assert any("正式使用前复核" in warning for warning in ai_completed.warnings)

    # 客户点名“4 张图表 + 表格/折线/饼图/柱状图”时，类型和总数都是硬合同。单对象饼图
    # 必须拿到三项组成指标，不能因为旧逻辑只给一个主指标而在渲染阶段消失。
    exact_visual_plan = _data_plan(
        PresentationStudioPlanRequest(
            intent="帮我制作一份梅西生涯数据 PPT，要求 4 张图表，表格、折线图、饼状图、柱状图。",
            structured_data_enabled=True,
        ),
        slides=adaptive_slides,
        brief=PresentationStudioBrief(
            title="梅西生涯数据",
            purpose="用多种数据视图展示职业生涯。",
            audience="足球爱好者",
            core_message="以总览、趋势和组成视图呈现数据。",
            theme="impact_contrast",
            theme_reason="适合数据对比表达。",
            fact_check_notice="数据由模型生成，供创作编辑使用。",
        ),
        blueprint=_StudioResearchBlueprint(
            needed=True,
            research_question="展示梅西职业生涯出场、进球、助攻和赛季进球趋势。",
            entities=["梅西"],
            entity_search_names=["Lionel Messi 梅西"],
            metrics=["职业生涯出场次数", "职业生涯总进球数", "职业生涯助攻数"],
            trend_metric="逐赛季联赛进球数",
            time_scope="职业生涯",
            comparison_scope="同一对象的生涯统计视图",
            chart_type="comparison_table",
            recommended_visuals=["comparison_table", "horizontal_bar", "trend_line", "share_pie"],
            visual_metrics=[
                ["职业生涯出场次数", "职业生涯总进球数", "职业生涯助攻数"],
                ["职业生涯出场次数", "职业生涯总进球数", "职业生涯助攻数"],
                ["逐赛季联赛进球数"],
                ["职业生涯出场次数", "职业生涯总进球数", "职业生涯助攻数"],
            ],
            target_slide_index=1,
            required_data_points=12,
            search_queries=["梅西职业生涯数据", "梅西赛季进球", "Lionel Messi career statistics"],
            preferred_source_types=["recognized_database"],
        ),
    )
    assert exact_visual_plan.requested_visuals == [
        "trend_table", "comparison_bar", "trend_line", "share_pie"
    ]
    assert exact_visual_plan.required_visual_count == 4
    assert exact_visual_plan.visual_metrics[-1] == [
        "职业生涯出场次数", "职业生涯总进球数", "职业生涯助攻数"
    ]
    ai_direct = fetch_ai_knowledge_draft_chart_data(
        exact_visual_plan,
        runtime=AiDraftRuntime(),  # type: ignore[arg-type]
    )
    assert [chart.chart_type for chart in ai_direct.charts] == exact_visual_plan.requested_visuals
    assert len(ai_direct.charts) == 4

    # 用户明确要多表格、柱状图和折线图时，趋势必须拥有独立的逐期指标和查询，不得把
    # “职业生涯总进球”这种累计快照强行画成折线。客户明确说“多张表”时还要把
    # 趋势明细交付为第二张原生表；四类视图仍只复用 4 + 6 = 10 个数据点。
    multi_visual_slides = [
        PresentationStudioSlidePlan(
            slide_id=f"content_{index}",
            role="content",
            title=f"数据视图 {index}",
            bullets=["使用可验证数据。"],
            visual_direction="按数据用途选择版式。",
        )
        for index in range(1, 5)
    ]
    multi_visual_plan = _data_plan(
        PresentationStudioPlanRequest(
            intent="制作梅西与 C 罗数据对比 PPT，多生成几张数据表格，并包含柱状图和折线图。",
            structured_data_enabled=True,
        ),
        slides=multi_visual_slides,
        brief=generic_plan.brief if hasattr(generic_plan, "brief") else PresentationStudioBrief(
            title="梅西与 C 罗数据对比",
            purpose="展示多维职业数据。",
            audience="足球爱好者",
            core_message="用多种数据视图比较两位球员。",
            theme="impact_contrast",
            theme_reason="适合对比表达。",
            fact_check_notice="所有数字均需核验来源。",
        ),
        blueprint=_StudioResearchBlueprint(
            needed=True,
            research_question="比较梅西与 C 罗的职业生涯总进球数和职业生涯助攻数。",
            entities=["梅西", "C罗"],
            entity_search_names=["Lionel Messi 梅西", "Cristiano Ronaldo C罗"],
            metrics=["职业生涯总进球数", "职业生涯助攻数"],
            time_scope="同一公开统计期间",
            comparison_scope="同一范围、单位和统计期间",
            chart_type="comparison_table",
            target_slide_index=1,
            required_data_points=4,
            search_queries=[
                "梅西 C罗 职业生涯总进球数 职业生涯助攻数 statistics",
                "梅西 职业生涯统计",
                "C罗 职业生涯统计",
            ],
            preferred_source_types=["official_profile"],
        ),
    )
    assert multi_visual_plan.requested_visuals == [
        "comparison_table", "trend_table", "comparison_bar", "trend_line"
    ]
    assert multi_visual_plan.visual_slide_ids == ["content_1", "content_2", "content_3", "content_4"]
    assert multi_visual_plan.trend_metric == "逐赛季联赛进球数"
    assert multi_visual_plan.required_data_points == 10
    assert multi_visual_plan.max_points == 10
    assert multi_visual_plan.entity_search_names == ["梅西", "C 罗"]
    assert any(
        "梅西 C 罗" in query and "comparison table" in query
        for query in multi_visual_plan.search_queries
    )
    assert any("梅西" in query and "every season" in query for query in multi_visual_plan.search_queries)
    assert any("C 罗" in query and "every season" in query for query in multi_visual_plan.search_queries)
    assert all(
        "season appearances goals" in query
        for query in multi_visual_plan.search_queries[:3]
    )

    # 用户可以只说一句自然语言。数量相同的重复表达取最大值而不是相加；三张表会拆成
    # 两张不同横向指标表和一张折线底稿，另有原生柱图与折线图，共五个独立页面。
    explicit_visual_request = PresentationStudioPlanRequest(
        intent=(
            "帮我审查梅西与C罗生涯数据对比ppt，要有多种数据对比，要有至少三个表格，"
            "至少一个表格，一个折线图，一个柱状图"
        ),
        structured_data_enabled=True,
    )
    explicit_intent = _data_visual_intent(explicit_visual_request)
    assert (explicit_intent.table_count, explicit_intent.bar_count, explicit_intent.line_count) == (3, 1, 1)
    assert explicit_intent.total == 5
    assert _requested_data_visual_count_hint(explicit_visual_request) == 5
    explicit_plan = _data_plan(
        explicit_visual_request,
        slides=[
            PresentationStudioSlidePlan(
                slide_id=f"content_{index}",
                role="content",
                title=f"数据视图 {index}",
                bullets=["使用可验证数据。"],
                visual_direction="按数据用途选择版式。",
            )
            for index in range(1, 6)
        ],
        brief=PresentationStudioBrief(
            title="梅西与 C 罗生涯数据对比",
            purpose="用多个独立数据视图完成职业生涯比较。",
            audience="足球爱好者",
            core_message="数据表和图表各自回答不同问题。",
            theme="impact_contrast",
            theme_reason="适合对比表达。",
            fact_check_notice="所有数字均需核验来源。",
        ),
        blueprint=_StudioResearchBlueprint(
            needed=True,
            research_question="比较梅西与 C 罗的六项职业数据及逐赛季联赛进球。",
            entities=["梅西", "C罗"],
            entity_search_names=["Lionel Messi 梅西", "Cristiano Ronaldo C罗"],
            metrics=["总进球", "总助攻", "总出场", "冠军数", "金球奖数", "帽子戏法数"],
            trend_metric="逐赛季联赛进球数",
            time_scope="同一公开统计期间",
            comparison_scope="逐项展示来源、单位和期间",
            chart_type="comparison_table",
            target_slide_index=1,
            required_data_points=18,
            search_queries=[
                "Lionel Messi Cristiano Ronaldo career goals assists appearances trophies statistics",
                "Lionel Messi career statistics",
                "Cristiano Ronaldo career statistics",
            ],
            preferred_source_types=["official_profile", "recognized_database"],
        ),
    )
    assert explicit_plan.requested_visuals == [
        "comparison_table", "comparison_table", "trend_table", "comparison_bar", "trend_line"
    ]
    assert explicit_plan.visual_slide_ids == [f"content_{index}" for index in range(1, 6)]
    assert explicit_plan.visual_metrics == [
        ["总进球", "总助攻", "总出场"],
        ["冠军数", "金球奖数", "帽子戏法数"],
        ["逐赛季联赛进球数"],
        ["总进球"],
        ["逐赛季联赛进球数"],
    ]
    assert explicit_plan.visual_contract_explicit is True
    assert explicit_plan.required_table_count == 3
    assert explicit_plan.required_bar_chart_count == 1
    assert explicit_plan.required_line_chart_count == 1
    pair_seed, entity_seeds = _trend_query_seeds(
        _StudioResearchBlueprint(
            needed=True,
            entities=["梅西", "C罗"],
            entity_search_names=["Lionel Messi 梅西", "Cristiano Ronaldo C罗"],
            metrics=["进球数"],
            search_queries=[
                "梅西 C罗 进球数",
                "梅西 C罗 进球数 statistics",
                "Lionel Messi 梅西 career statistics",
                "Cristiano Ronaldo C罗 career statistics",
            ],
        )
    )
    assert pair_seed == "Lionel Messi 梅西 Cristiano Ronaldo C罗"
    assert entity_seeds == ["Lionel Messi 梅西", "Cristiano Ronaldo C罗"]

    # 宽泛比较最多保留六个指标，既控制研究预算，也能给多个表格分配不同主题。
    broad_blueprint = _StudioResearchBlueprint(
        needed=True,
        research_question="比较甲公司与乙公司的收入、利润、市场份额和用户增长。",
        entities=["甲公司", "乙公司"],
        metrics=["营业收入", "净利润", "市场份额", "用户增长"],
        time_scope="同一公开统计期间",
        comparison_scope="同一来源、单位和统计期间",
        chart_type="comparison_table",
        target_slide_index=1,
        required_data_points=8,
        search_queries=["甲公司乙公司营业收入净利润市场份额用户增长统计"],
        preferred_source_types=["official_report"],
    )
    focused_blueprint = _focus_broad_comparison_blueprint(broad_blueprint)
    assert focused_blueprint.entities == ["甲公司", "乙公司"]
    assert focused_blueprint.metrics == ["营业收入", "净利润", "市场份额", "用户增长"]
    assert focused_blueprint.required_data_points == 8
    assert focused_blueprint.search_queries == ["甲公司乙公司营业收入净利润市场份额用户增长统计"]
    # 计划阶段不能凭空把“截至某年”加进客户没有限定时间的主题；否则动态公开统计即使来自
    # 同一页，也会被后续 Verifier 错误地当成与虚构年份冲突。客户明确限定年份时则保留原条件。
    unbounded_blueprint = _align_research_time_scope(
        PresentationStudioPlanRequest(intent="制作甲公司与乙公司营业收入对比 PPT。"),
        focused_blueprint,
    )
    assert "读取快照" in unbounded_blueprint.time_scope
    assert "读取快照" in unbounded_blueprint.comparison_scope
    assert "读取快照" in unbounded_blueprint.research_question
    bounded_blueprint = _align_research_time_scope(
        PresentationStudioPlanRequest(intent="制作甲公司与乙公司 2024 年营业收入对比 PPT。"),
        focused_blueprint,
    )
    assert bounded_blueprint.time_scope == focused_blueprint.time_scope
    pair_query_blueprint = _ensure_pair_comparison_query(
        focused_blueprint.model_copy(
            update={"chart_type": "comparison_table", "search_queries": ["甲公司乙公司营收 statistics"]}
        )
    )
    assert pair_query_blueprint.search_queries[0] == "甲公司 乙公司 营业收入 净利润 市场份额"
    assert len(pair_query_blueprint.search_queries) == 2

    # 专用研究规划器是图表链路的入口。不同模型常把字段改成同义名、整数改成文字或包一层
    # research_plan；这些属于格式差异，必须先在零联网阶段归一化，不能因此退化成普通 PPT。
    aliased_blueprint_payload = {
        "research_plan": {
            "research_goal": "比较梅西与 C 罗的职业进球和助攻。",
            "comparison_entities": "梅西，C 罗",
            "entity_aliases": ["Lionel Messi 梅西", "Cristiano Ronaldo C 罗"],
            "indicators": "职业进球；职业助攻",
            "timeframe": "截至同一公开统计截止日期",
            "comparison_rule": "同一赛事范围、单位和截止日期",
            "chart": "分组柱状图",
            "target_slide": "第 2 页",
            "data_points": "4 项",
            "queries": [
                "Messi Ronaldo career goals assists official statistics",
                "Lionel Messi career statistics official profile",
                "Cristiano Ronaldo career statistics official profile",
            ],
            "source_types": "official_profile,recognized_database",
        }
    }
    aliased_blueprint = _parse_research_blueprint(json.dumps(aliased_blueprint_payload, ensure_ascii=False))
    assert aliased_blueprint.entities == ["梅西", "C 罗"]
    assert aliased_blueprint.entity_search_names == ["Lionel Messi 梅西", "Cristiano Ronaldo C 罗"]
    assert aliased_blueprint.metrics == ["职业进球", "职业助攻"]
    assert aliased_blueprint.chart_type == "grouped_bar"
    assert aliased_blueprint.target_slide_index == 2
    assert aliased_blueprint.required_data_points == 4

    research_planner_output = _StudioModelOutput(
        title="梅西与 C 罗：数据视角下的双骄对比",
        purpose="用可验证的统一口径比较两位球员。",
        audience="足球爱好者",
        core_message="图表必须基于同一来源和同一截止日期。",
        theme="impact_contrast",
        theme_reason="对比主题需要清晰区分对象。",
        content_slides=[
            _StudioContentSlide(
                title="进球与助攻对比",
                bullets=["先统一统计范围。", "只采用可追溯数据。"],
                layout="comparison",
                visual_direction="使用数据表突出可比性。",
            )
        ],
        asset_queries=["football data comparison"],
    )
    numeric_output = research_planner_output.model_copy(
        update={
            "content_slides": [
                _StudioContentSlide(
                    title="未核验数字处理",
                    bullets=["模型声称进球数为 850+，但客户没有提供来源。", "先统一统计范围。"],
                    layout="comparison",
                    visual_direction="保留数据位置。",
                )
            ]
        }
    )
    sanitized_output, stripped_numeric_claims = _strip_unverified_numeric_claims(
        numeric_output,
        request=PresentationStudioPlanRequest(intent="制作梅西与 C 罗数据对比 PPT。"),
    )
    assert stripped_numeric_claims is True
    assert "850" not in sanitized_output.content_slides[0].bullets[0]
    assert "可追溯来源" in sanitized_output.content_slides[0].bullets[0]

    class ResearchBlueprintRuntime:
        """用固定模型回包验证研究规划修复，不触发网络或真实模型。"""

        def __init__(self, responses: list[str]) -> None:
            self._responses = iter(responses)
            self.call_count = 0

        async def tool_turn(self, **_: object) -> ModelToolTurn:
            self.call_count += 1
            return ModelToolTurn(content=next(self._responses))

    repair_runtime = ResearchBlueprintRuntime(
        ["首次回复没有 JSON", json.dumps(aliased_blueprint_payload, ensure_ascii=False)]
    )
    repaired_blueprint, repair_used, fallback_used = asyncio.run(
        _request_research_blueprint(
            runtime=repair_runtime,  # type: ignore[arg-type]
            request=PresentationStudioPlanRequest(
                intent="制作梅西与 C 罗的职业数据对比 PPT。",
                structured_data_enabled=True,
            ),
            output=research_planner_output,
        )
    )
    assert repair_runtime.call_count == 2
    assert repair_used is True
    assert fallback_used is False
    assert repaired_blueprint.chart_type == "grouped_bar"

    # 两次 JSON 都无法解析时，仅“明确双对象 + 已出现量化指标”允许建立查询蓝图。它不含
    # 数字、来源或 URL，确认导出后仍必须经过 ResearchGateway 的联网证据和 Verifier。
    fallback_request = PresentationStudioPlanRequest(
        intent="帮我做一个梅西与C罗的进球数据对比 PPT。",
        structured_data_enabled=True,
    )
    fallback_runtime = ResearchBlueprintRuntime(["无结构化输出", "仍然不是 JSON"])
    fallback_blueprint, fallback_repair_used, fallback_used = asyncio.run(
        _request_research_blueprint(
            runtime=fallback_runtime,  # type: ignore[arg-type]
            request=fallback_request,
            output=research_planner_output,
        )
    )
    assert fallback_runtime.call_count == 2
    assert fallback_repair_used is True
    assert fallback_used is True
    assert fallback_blueprint.entities == ["梅西", "C罗"]
    assert fallback_blueprint.metrics == ["职业生涯总进球数", "职业生涯助攻数"]
    assert all(not any(character.isdigit() for character in query) for query in fallback_blueprint.search_queries)
    single_fallback_output = research_planner_output.model_copy(
        update={
            "title": "梅西生涯数据全景",
            "content_slides": [
                _StudioContentSlide(
                    title="进球、助攻与出场",
                    bullets=["比较生涯总量。", "观察逐赛季进球趋势。"],
                    layout="metrics",
                    visual_direction="使用数据视图。",
                )
            ],
        }
    )
    single_fallback = _infer_conservative_research_blueprint(
        request=PresentationStudioPlanRequest(
            intent="帮我生成梅西生涯数据 PPT。",
            structured_data_enabled=True,
        ),
        output=single_fallback_output,
    )
    assert single_fallback is not None
    assert single_fallback.entities == ["梅西"]
    assert single_fallback.recommended_visuals == [
        "comparison_table", "horizontal_bar", "trend_line"
    ]
    no_metric_output = research_planner_output.model_copy(
        update={
            "title": "甲与乙：主题介绍",
            "content_slides": [
                _StudioContentSlide(
                    title="主题背景",
                    bullets=["说明背景信息。", "不声明量化指标。"],
                    layout="insight_cards",
                    visual_direction="使用概念视觉。",
                )
            ],
        }
    )
    assert _infer_conservative_research_blueprint(
        request=PresentationStudioPlanRequest(
            intent="制作甲与乙的数据对比 PPT。",
            structured_data_enabled=True,
        ),
        output=no_metric_output,
    ) is None

    # 通用 ResearchGateway 不读取模型记忆：原生搜索只给出候选 URL，随后受控读取少量页面，
    # 第二次模型仅从这些片段抽取数据。这里用 MockTransport 与假 Runtime 验证完整链路不联网。
    research_runtime_calls: list[int] = []

    class ResearchGatewayRuntime:
        async def native_web_search_sources(self, *, queries: tuple[str, ...], max_uses: int) -> NativeWebSearchResult:
            assert queries == (
                "甲队乙队 2024 采用率 官方统计",
                "甲队 2024 采用率 统计",
                "乙队 2024 采用率 统计",
            )
            assert max_uses == 3
            return NativeWebSearchResult(
                sources=(
                    NativeWebSearchSource(title="甲队官方统计", url="https://stats.example.com/team-a"),
                    NativeWebSearchSource(title="乙队官方统计", url="https://stats.example.com/team-b"),
                ),
                query_count=3,
                retrieved_at="2026-08-12T00:00:00+00:00",
            )

        async def tool_turn(self, **_: object) -> ModelToolTurn:
            research_runtime_calls.append(1)
            return ModelToolTurn(
                content=json.dumps(
                    {
                        "status": "complete",
                        "title": "甲队与乙队采用率对比（2024）",
                        "points": [
                            {
                                "entity": "甲队",
                                "metric": "采用率",
                                "value": 42,
                                "unit": "%",
                                "period": "2024年",
                                "source_ids": ["S1"],
                                "evidence_quote": "甲队采用率为 42%，统计时间为 2024 年。",
                            },
                            {
                                "entity": "乙队",
                                "metric": "采用率",
                                "value": 38,
                                "unit": "%",
                                "period": "2024年",
                                "source_ids": ["S2"],
                                "evidence_quote": "乙队采用率为 38%，统计时间为 2024 年。",
                            },
                        ],
                        "notes": [],
                    },
                    ensure_ascii=False,
                )
            )

    research_plan = PresentationStudioDataPlan(
        state="research_planned",
        provider="research_gateway",
        chart_type="comparison_table",
        slide_id="content_1",
        research_question="甲队与乙队 2024 年采用率如何比较？",
        entities=["甲队", "乙队"],
        metrics=["采用率"],
        time_scope="2024年",
        comparison_scope="同一统计口径、单位和时间",
        required_data_points=2,
        search_queries=[
            "甲队乙队 2024 采用率 官方统计",
            "甲队 2024 采用率 统计",
            "乙队 2024 采用率 统计",
        ],
        preferred_source_types=["official_statistics"],
        notice="验证 ResearchGateway 计划。",
    )

    def research_page_handler(request: httpx.Request) -> httpx.Response:
        assert request.url.host == "stats.example.com"
        if request.url.path == "/unreadable":
            # 模拟来源可访问但不是网关允许读取的正文类型；这会走真实的 ValueError -> warning
            # 分支，而不是让测试夹具自身抛 KeyError 干扰“来源回退”行为。
            return httpx.Response(200, headers={"content-type": "application/json"}, text="{}")
        pages = {
            "/team-a": "<html><body>甲队采用率为 42%，统计时间为 2024 年。"
            "该口径只统计已完成的年度业务样本，并由官方统计部门公开发布。"
            "本页用于展示口径说明、统计范围和可复核的年度结果，不包含预测值。</body></html>",
            "/team-b": "<html><body>乙队采用率为 38%，统计时间为 2024 年。"
            "该口径只统计已完成的年度业务样本，并由官方统计部门公开发布。"
            "本页用于展示口径说明、统计范围和可复核的年度结果，不包含预测值。</body></html>",
        }
        return httpx.Response(200, headers={"content-type": "text/html"}, text=pages[request.url.path])

    with httpx.Client(transport=httpx.MockTransport(research_page_handler)) as mock_client:
        research_data_sample = fetch_research_gateway_chart_data(
            research_plan,
            runtime=ResearchGatewayRuntime(),  # type: ignore[arg-type]
            page_client=mock_client,
        )
    assert research_data_sample.chart is not None
    assert research_data_sample.chart.chart_type == "comparison_table"
    assert research_data_sample.chart.search_provider == "deepseek_native"
    assert len(research_data_sample.chart.sources) == 2
    assert research_runtime_calls == [1]

    # 一次抽取可复用为表格、柱状图和双序列折线图。各视图独立选择完整数据子集，不能再把
    # 多指标全集塞给单指标柱图，也不能把两个对象首尾相连成一条错误折线。
    multi_view_source = ResearchGatewaySource(
        source_id="S20",
        title="甲队与乙队多期统计",
        source_url="https://stats.example.com/multi-view",
        excerpt=(
            "甲队总量为42，乙队总量为38。甲队2022年得分30，2023年得分35，2024年得分42；"
            "乙队2022年得分28，2023年得分33，2024年得分38。"
        ),
        retrieved_at="2026-08-13T08:30:00+00:00",
        retrieval_method="tavily_raw_content",
    )
    multi_view_points = [
        {"entity": "甲队", "metric": "总量", "value": 42, "unit": "分", "period": "2024年", "source_ids": ["S20"], "evidence_quote": "甲队总量为42，乙队总量为38"},
        {"entity": "乙队", "metric": "总量", "value": 38, "unit": "分", "period": "2024年", "source_ids": ["S20"], "evidence_quote": "乙队总量为38"},
    ]
    for entity, values in (("甲队", (30, 35, 42)), ("乙队", (28, 33, 38))):
        for year, value in zip((2022, 2023, 2024), values, strict=True):
            multi_view_points.append(
                {
                    "entity": entity,
                    "metric": "年度得分",
                    "value": value,
                    "unit": "分",
                    "period": f"{year}年",
                    "source_ids": ["S20"],
                    "evidence_quote": (
                        f"{entity}{year}年得分{value}"
                        if year == 2022
                        else f"{year}年得分{value}"
                    ),
                }
            )
    multi_view_plan = research_plan.model_copy(
        update={
            "metrics": ["总量", "年度得分"],
            "trend_metric": "年度得分",
            "required_data_points": 8,
            "requested_visuals": ["comparison_table", "trend_table", "comparison_bar", "trend_line"],
            "visual_slide_ids": ["content_1", "content_2", "content_3", "content_4"],
        }
    )
    multi_view_charts, multi_view_warnings = _validate_extraction_views(
        json.dumps({"status": "complete", "title": "甲乙数据", "points": multi_view_points, "notes": []}, ensure_ascii=False),
        plan=multi_view_plan,
        sources=(multi_view_source,),
        search_provider="tavily",
        query_count=3,
        extraction_attempts=1,
    )
    assert [chart.chart_type for chart in multi_view_charts] == [
        "comparison_table", "trend_table", "comparison_bar", "trend_line"
    ]
    assert [chart.slide_id for chart in multi_view_charts] == [
        "content_1", "content_2", "content_3", "content_4"
    ]
    assert len(multi_view_charts[0].points) == 4
    assert len(multi_view_charts[1].points) == 6
    assert len(multi_view_charts[2].points) == 2
    assert len(multi_view_charts[3].points) == 6
    assert not multi_view_warnings

    # 五视图合同使用同一次抽取结果：两张横向表各取不同的三项指标，趋势表与折线图复用
    # 三期序列，柱图只取第一项同口径指标。最终数量按真实 Chart/Table 类型计数。
    contract_metrics = ["总进球", "总助攻", "总出场", "冠军数", "金球奖数", "帽子戏法数"]
    contract_sentences: list[str] = []
    contract_points: list[dict[str, object]] = []
    for entity, offset in (("梅西", 0), ("C罗", 10)):
        for metric_index, metric in enumerate(contract_metrics, start=1):
            value = 100 + offset + metric_index
            sentence = f"{entity}在2025年统计中的{metric}为{value}项。"
            contract_sentences.append(sentence)
            contract_points.append(
                {
                    "entity": entity,
                    "metric": metric,
                    "value": value,
                    "unit": "项",
                    "period": "2025年",
                    "source_ids": ["S40"],
                    "evidence_quote": sentence,
                }
            )
        for season, value in zip(("2021/22", "2022/23", "2023/24"), (30 + offset, 32 + offset, 35 + offset), strict=True):
            sentence = f"{entity}{season}赛季的联赛进球数为{value}球。"
            contract_sentences.append(sentence)
            contract_points.append(
                {
                    "entity": entity,
                    "metric": "逐赛季联赛进球数",
                    "value": value,
                    "unit": "球",
                    "period": season,
                    "source_ids": ["S40"],
                    "evidence_quote": sentence,
                }
            )
    contract_source = ResearchGatewaySource(
        source_id="S40",
        title="梅西与C罗职业数据测试来源",
        source_url="https://stats.example.com/contract",
        excerpt=" ".join(contract_sentences),
        retrieved_at="2026-08-14T00:00:00+00:00",
        retrieval_method="test_fixture",
    )
    contract_charts, contract_warnings = _validate_extraction_views(
        json.dumps(
            {"status": "complete", "title": "梅西与C罗生涯数据", "points": contract_points, "notes": []},
            ensure_ascii=False,
        ),
        plan=explicit_plan,
        sources=(contract_source,),
        search_provider="fixture",
        query_count=3,
        extraction_attempts=1,
    )
    assert [chart.chart_type for chart in contract_charts] == [
        "comparison_table", "comparison_table", "trend_table", "comparison_bar", "trend_line"
    ]
    assert {point.metric for point in contract_charts[0].points} == set(contract_metrics[:3])
    assert {point.metric for point in contract_charts[1].points} == set(contract_metrics[3:])
    assert not contract_warnings
    assert _structured_data_contract_gap(explicit_plan, contract_charts) == ""
    assert _structured_data_contract_gap(explicit_plan, contract_charts[:3]) == "柱状图 0/1、折线图 0/1、数据视图 3/5"
    assert _research_source_marker_present(
        "S20 · 甲队与乙队多期统计 · https://stats.example.com/multi-vi",
        multi_view_source,
    )
    assert not _research_source_marker_present(
        "S20 · 甲队与乙队多期统计",
        multi_view_source,
    )
    assert _trend_period_label("2010–11") == "2010/11"
    assert _trend_period_label("2010-11") == "2010/11"
    assert _trend_period_label("2010/2011") == "2010/11"
    assert _trend_period_label("2024年") == "2024年"

    # 第二次已有抽取预算若只缺趋势视图，必须明确只抽趋势点，不能再次把职业总量抢占输出。
    captured_completion_prompt: dict[str, str] = {}

    class TrendCompletionRuntime:
        async def tool_turn(self, **kwargs: object) -> ModelToolTurn:
            captured_completion_prompt["system"] = str(kwargs["system_prompt"])
            messages = kwargs["messages"]
            captured_completion_prompt["user"] = messages[0].content  # type: ignore[index,union-attr]
            return ModelToolTurn(content='{"status":"insufficient","title":"","points":[],"notes":["no trend"]}')

    asyncio.run(
        _extract_from_sources(
            runtime=TrendCompletionRuntime(),  # type: ignore[arg-type]
            plan=multi_view_plan,
            sources=(multi_view_source,),
            target_visuals=("trend_table", "trend_line"),
            repair_reason="aggregate views exist; trend views are missing",
        )
    )
    assert "ONLY for trend_table/trend_line" in captured_completion_prompt["user"]
    assert "exactly three common period labels" in captured_completion_prompt["user"]
    assert 'Candidate metrics for this pass: ["年度得分"]' in captured_completion_prompt["user"]
    assert "Candidate metrics for this pass: [\"总量\"" not in captured_completion_prompt["user"]
    assert "Maximum points: 6" in captured_completion_prompt["user"]
    assert 'Requested views for this pass: ["trend_table", "trend_line"]' in captured_completion_prompt["user"]

    # Tavily 偶尔能给出 URL 却没有可读取正文。自动模式只可把蓝图剩余查询交给 DeepSeek
    # 原生搜索，随后继续走相同的来源和口径校验；不能复用第一组查询或继续扩大查询预算。
    fallback_native_calls: list[tuple[str, ...]] = []

    class SourceFallbackRuntime(ResearchGatewayRuntime):
        async def native_web_search_sources(
            self,
            *,
            queries: tuple[str, ...],
            max_uses: int,
        ) -> NativeWebSearchResult:
            assert queries == ("甲队与乙队采用率对比 官方统计", "甲队乙队年度业务样本统计")
            assert max_uses == 2
            fallback_native_calls.append(queries)
            return NativeWebSearchResult(
                sources=(
                    NativeWebSearchSource(title="甲队官方统计", url="https://stats.example.com/team-a"),
                    NativeWebSearchSource(title="乙队官方统计", url="https://stats.example.com/team-b"),
                ),
                query_count=2,
                retrieved_at="2026-08-12T00:00:00+00:00",
            )

    def unreadable_tavily_fetch(_: tuple[str, ...]) -> object:
        return type(
            "UnreadableTavilyResult",
            (),
            {
                "candidates": (
                    TavilyResearchCandidate(
                        title="无法核验的候选来源",
                        url="https://stats.example.com/unreadable",
                    ),
                ),
                "query_count": 3,
                "warnings": (),
            },
        )()

    original_gateway_settings = research_gateway_module.settings
    original_tavily_fetch = research_gateway_module.fetch_tavily_research_sources
    source_fallback_plan = research_plan.model_copy(
        update={
            "search_queries": [
                "甲队乙队 2024 采用率 官方统计",
                "甲队与乙队采用率对比 官方统计",
                "甲队乙队年度业务样本统计",
                "甲队 2024 采用率 统计",
                "乙队 2024 采用率 统计",
            ]
        }
    )
    research_gateway_module.settings = replace(
        original_gateway_settings,
        presentation_research_search_provider="auto",
        tavily_api_key="tavily-test-key",
    )
    research_gateway_module.fetch_tavily_research_sources = unreadable_tavily_fetch
    try:
        with httpx.Client(transport=httpx.MockTransport(research_page_handler)) as mock_client:
            source_fallback = fetch_research_gateway_chart_data(
                source_fallback_plan,
                runtime=SourceFallbackRuntime(),  # type: ignore[arg-type]
                page_client=mock_client,
            )
    finally:
        research_gateway_module.settings = original_gateway_settings
        research_gateway_module.fetch_tavily_research_sources = original_tavily_fetch
    assert fallback_native_calls == [
        ("甲队与乙队采用率对比 官方统计", "甲队乙队年度业务样本统计")
    ], fallback_native_calls
    assert source_fallback.chart is not None
    assert source_fallback.chart.search_provider == "deepseek_native_fallback"
    assert source_fallback.chart.query_count == 5
    assert any("切换至 DeepSeek" in warning for warning in source_fallback.warnings)

    # 显式模型路由启用后，Kimi 可能因公开网页片段触发 content_filter，但 Runtime 不得读取
    # DeepSeek 的 Key 并悄悄换模型。既有修复预算仍可在同一 Provider 内尝试一次；两次都
    # 失败时，本次只跳过图表并返回可操作原因。
    class ContentFilterRuntime(ResearchGatewayRuntime):
        provider = "kimi"

        def __init__(self) -> None:
            super().__init__()
            self.calls = 0

        async def tool_turn(self, **_: object) -> ModelToolTurn:
            self.calls += 1
            raise ModelGatewayError(
                "模型接口返回 HTTP 400 (content_filter · The request was rejected because it was considered high risk)"
            )

    content_filter_runtime = ContentFilterRuntime()
    with httpx.Client(transport=httpx.MockTransport(research_page_handler)) as mock_client:
        content_filter_result = fetch_research_gateway_chart_data(
            research_plan,
            runtime=content_filter_runtime,  # type: ignore[arg-type]
            page_client=mock_client,
        )
    assert content_filter_runtime.calls == 2
    assert content_filter_result.chart is None
    assert any("首次抽取未通过" in warning for warning in content_filter_result.warnings)
    assert not any("DeepSeek" in warning for warning in content_filter_result.warnings)

    # Tavily 返回的清洗正文会跳过本机逐页抓取，但仍与直连页面使用相同的来源契约。通过
    # 只提供两个有逐字证据的数据点，验证“稳定搜索服务”不会绕开对象、时间或单位核验。
    prefetched_candidates = (
        {
            "title": "甲队官方统计",
            "url": "https://stats.example.com/team-a",
            "raw_content": "甲队采用率为 42%，统计时间为 2024 年。官方页面公开了统计口径和年度结果。",
        },
        {
            "title": "乙队官方统计",
            "url": "https://stats.example.com/team-b",
            "raw_content": "乙队采用率为 38%，统计时间为 2024 年。官方页面公开了统计口径和年度结果。",
        },
    )

    # 直接调用内部读取/Verifier，避免在进程内修改全局 Settings；外部 Tavily API 的请求参数
    # 已在前面的 MockTransport 验证，下面只验证其清洗正文不能绕过图表数据契约。

    prefetched_sources, prefetched_warnings = _read_sources(
        tuple(
            _from_tavily_candidate(
                TavilyResearchCandidate(
                    title=item["title"],
                    url=item["url"],
                    raw_content=item["raw_content"],
                )
            )
            for item in prefetched_candidates
        ),
        page_client=None,
    )
    assert not prefetched_warnings
    assert all(source.retrieval_method == "tavily_raw_content" for source in prefetched_sources)
    prefetched_chart = _validate_extraction(
        json.dumps(
            {
                "status": "complete",
                "title": "甲队与乙队采用率对比（2024）",
                "points": [
                    {
                        "entity": "甲队",
                        "metric": "采用率",
                        "value": 42,
                        "unit": "%",
                        "period": "2024年",
                        "source_ids": ["S1"],
                        "evidence_quote": "甲队采用率为 42%，统计时间为 2024 年。",
                    },
                    {
                        "entity": "乙队",
                        "metric": "采用率",
                        "value": 38,
                        "unit": "%",
                        "period": "2024年",
                        "source_ids": ["S2"],
                        "evidence_quote": "乙队采用率为 38%，统计时间为 2024 年。",
                    },
                ],
                "notes": [],
            },
            ensure_ascii=False,
        ),
        plan=research_plan,
        sources=prefetched_sources,
        search_provider="tavily",
        query_count=1,
        extraction_attempts=1,
    )
    assert prefetched_chart.audit_metadata()["sources"][0]["retrieval_method"] == "tavily_raw_content"
    assert prefetched_chart.audit_metadata()["search_provider"] == "tavily"

    # 少数模型会把可选 notes 误写成单个字符串。它不属于图表事实字段，允许归一化为数组；
    # 来源、数值、时间和单位仍须完整通过相同的 Verifier，不能借此放松数据准入标准。
    note_string_payload = {
        "status": "complete",
        "title": "甲队与乙队采用率对比（2024）",
        "points": [
            {
                "entity": "甲队",
                "metric": "采用率",
                "value": 42,
                "unit": "%",
                "period": "2024年",
                "source_ids": ["S1"],
                "evidence_quote": "甲队采用率为 42%，统计时间为 2024 年。",
            },
            {
                "entity": "乙队",
                "metric": "采用率",
                "value": 38,
                "unit": "%",
                "period": "2024年",
                "source_ids": ["S2"],
                "evidence_quote": "乙队采用率为 38%，统计时间为 2024 年。",
            },
        ],
        "notes": "来源已按同一年度口径核对。",
    }
    note_string_chart = _validate_extraction(
        json.dumps(note_string_payload, ensure_ascii=False),
        plan=research_plan,
        sources=prefetched_sources,
        search_provider="tavily",
        query_count=1,
        extraction_attempts=1,
    )
    assert len(note_string_chart.points) == 2

    # 多视图抽取中一条坏引用不能撤回其它已核验数据。Verifier 应逐点丢弃它并保留警告；
    # 但若所有点都坏，后面的完整对象/指标校验仍会拒绝整项交付。
    partially_invalid_payload = json.loads(json.dumps(note_string_payload, ensure_ascii=False))
    partially_invalid_payload["points"].append(
        {
            "entity": "甲队",
            "metric": "不存在的指标",
            "value": 999,
            "unit": "分",
            "period": "2024年",
            "source_ids": ["S1"],
            "evidence_quote": "来源正文中没有这条伪造引文",
        }
    )
    partial_charts, partial_warnings = _validate_extraction_views(
        json.dumps(partially_invalid_payload, ensure_ascii=False),
        plan=research_plan,
        sources=prefetched_sources,
        search_provider="tavily",
        query_count=1,
        extraction_attempts=1,
    )
    assert len(partial_charts) == 1
    assert len(partial_charts[0].points) == 2
    assert any("丢弃" in warning and "1 条" in warning for warning in partial_warnings)

    # 国际公开统计页常使用英文表头，而研究蓝图和最终页面使用中文指标。恢复器只能在
    # 对象、数值、赛季和指标别名同时命中时返回来源中的连续原文，不能信任模型翻译的引文。
    bilingual_trend_source = ResearchGatewaySource(
        source_id="S30",
        title="Season goals",
        source_url="https://stats.example.com/season-goals",
        excerpt=(
            "31 Messi League Goals 2012/13 46 Ronaldo League Goals 2012/13 "
            "50 Messi League Goals 2011/12 46 Ronaldo League Goals 2011/12"
        ),
        retrieved_at="2026-08-13T08:30:00+00:00",
        retrieval_method="tavily_raw_content",
    )
    recovered_quote = _locate_evidence_quote(
        "Messi scored 31 goals in 2012/13",
        entity="梅西",
        metric="逐赛季进球数",
        value=31,
        period="2012/13",
        source_ids=("S30",),
        source_map={"S30": bilingual_trend_source},
    )
    assert "31" in recovered_quote and "Messi" in recovered_quote and "2012/13" in recovered_quote
    assert _locate_evidence_quote(
        "Messi scored 99 goals in 2012/13",
        entity="梅西",
        metric="逐赛季进球数",
        value=99,
        period="2012/13",
        source_ids=("S30",),
        source_map={"S30": bilingual_trend_source},
    ) == ""

    # 动态对比页常同时列出两方“当前值”，却没有写明截止日期。只有双对象、单指标、同一
    # 已读取页面的两个数值可以使用来源快照；最终会写入真实读取日期，不能伪装为统计年份。
    snapshot_source = ResearchGatewaySource(
        source_id="S9",
        title="甲队与乙队当前采用率",
        source_url="https://stats.example.com/current-comparison",
        excerpt="当前页面显示甲队采用率为 42%，乙队采用率为 38%。",
        retrieved_at="2026-08-13T08:30:00+00:00",
        retrieval_method="tavily_raw_content",
    )
    snapshot_payload = {
        "status": "complete",
        "title": "甲队与乙队采用率对比",
        "points": [
            {
                "entity": "甲队",
                "metric": "采用率",
                "value": 42,
                "unit": "%",
                "period": "source_snapshot",
                "source_ids": ["S9"],
                "evidence_quote": "甲队采用率为 42%",
            },
            {
                "entity": "乙队",
                "metric": "采用率",
                "value": 38,
                "unit": "%",
                "period": "source_snapshot",
                "source_ids": ["S9"],
                "evidence_quote": "乙队采用率为 38%",
            },
        ],
        "notes": [],
    }
    snapshot_chart = _validate_extraction(
        json.dumps(snapshot_payload, ensure_ascii=False),
        plan=research_plan,
        sources=(snapshot_source,),
        search_provider="tavily",
        query_count=1,
        extraction_attempts=1,
    )
    assert {point.period for point in snapshot_chart.points} == {"网页读取快照（2026-08-13）"}
    # 演示型对比表允许分别读取两个公开页面，但会在每个单元格保留来源和读取时间；这比因为
    # 达不到审计级同源标准而把已查到的整张表静默删除更符合客户预期。
    mixed_source_payload = json.loads(json.dumps(snapshot_payload, ensure_ascii=False))
    mixed_source_payload["points"][1]["source_ids"] = ["S10"]
    second_snapshot_source = replace(snapshot_source, source_id="S10")
    mixed_source_chart = _validate_extraction(
        json.dumps(mixed_source_payload, ensure_ascii=False),
        plan=research_plan,
        sources=(snapshot_source, second_snapshot_source),
        search_provider="tavily",
        query_count=1,
        extraction_attempts=1,
    )
    assert mixed_source_chart.chart_type == "comparison_table"
    assert len(mixed_source_chart.sources) == 2
    mixed_source_bar = _validate_extraction(
        json.dumps(mixed_source_payload, ensure_ascii=False),
        plan=research_plan.model_copy(update={"chart_type": "comparison_bar"}),
        sources=(snapshot_source, second_snapshot_source),
        search_provider="tavily",
        query_count=1,
        extraction_attempts=1,
    )
    assert mixed_source_bar.chart_type == "comparison_table"

    mixed_period_payload = json.loads(json.dumps(snapshot_payload, ensure_ascii=False))
    mixed_period_payload["points"][1]["period"] = "2024年"
    try:
        _validate_extraction(
            json.dumps(mixed_period_payload, ensure_ascii=False),
            plan=research_plan,
            sources=(snapshot_source,),
            search_provider="tavily",
            query_count=1,
            extraction_attempts=1,
        )
    except ResearchGatewayValidationError:
        pass
    else:
        raise AssertionError("来源没有 2024 年期间时，不能把读取快照伪造成 2024 年数据")

    # 首轮查询不能再机械截前三条。成对查询、甲方单独查询和乙方单独查询应一次并发完成，
    # 其余语句才留作失败后的受限补查。
    query_plan = research_plan.model_copy(
        update={
            "search_queries": [
                "甲队 乙队 采用率对比",
                "甲队 采用率 官方数据",
                "甲队 乙队 数据统计",
                "乙队 采用率 官方数据",
                "双方采用率 历史资料",
            ]
        }
    )
    primary_queries, remaining_queries = _split_research_queries(query_plan)
    assert primary_queries == (
        "甲队 乙队 采用率对比",
        "甲队 采用率 官方数据",
        "乙队 采用率 官方数据",
    )
    assert len(remaining_queries) == 2

    # 来源 ID、时间和单位是数据 Verifier 的硬约束；任一失真都只能降级，不能绘制“看似合理”的图。
    class InvalidResearchGatewayRuntime(ResearchGatewayRuntime):
        async def tool_turn(self, **_: object) -> ModelToolTurn:
            return ModelToolTurn(
                content=json.dumps(
                    {
                        "status": "complete",
                        "points": [
                            {
                                "entity": "甲队",
                                "metric": "采用率",
                                "value": 42,
                                "unit": "%",
                                "period": "2024年",
                                "source_ids": ["S1"],
                                "evidence_quote": "甲队采用率为 42%，统计时间为 2024 年。",
                            },
                            {
                                "entity": "乙队",
                                "metric": "采用率",
                                "value": 38,
                                "unit": "人",
                                "period": "2023年",
                                "source_ids": ["S2"],
                                "evidence_quote": "乙队采用率为 38%，统计时间为 2024 年。",
                            },
                        ],
                        "notes": [],
                    },
                    ensure_ascii=False,
                )
            )

    with httpx.Client(transport=httpx.MockTransport(research_page_handler)) as mock_client:
        invalid_research = fetch_research_gateway_chart_data(
            research_plan,
            runtime=InvalidResearchGatewayRuntime(),  # type: ignore[arg-type]
            page_client=mock_client,
        )
    assert invalid_research.chart is None
    assert any("数据" in warning or "单位" in warning for warning in invalid_research.warnings)

    # 主创作模型偶尔会多返回一个旧版/半成品 research_blueprint；它必须被忽略，正式研究
    # 蓝图只能来自专用规划器，不能让附带字段破坏已经合法的主创作计划。
    extra_blueprint_payload = {
        "title": "职责隔离验证",
        "purpose": "验证主创作计划不消费研究半成品。",
        "audience": "测试人员",
        "core_message": "创作与研究规划应解耦。",
        "theme": "executive_blue",
        "theme_reason": "采用稳健主题。",
        "content_slides": [
            {
                "title": "验证页",
                "bullets": ["主内容保持合法。", "附带研究字段被忽略。"],
                "layout": "comparison",
                "visual_direction": "清楚呈现职责边界。",
            }
        ],
        "asset_queries": ["presentation contract verification"],
        "research_blueprint": {"needed": True, "required_data_points": ["错误类型"]},
    }
    parsed_main_plan = _parse_model_output(
        json.dumps(extra_blueprint_payload, ensure_ascii=False),
        expected_content_slide_count=1,
    )
    assert parsed_main_plan.research_blueprint.needed is False

    # 即梦素材的可追溯性来自 artifact 元数据，不依赖把水印烧进客户交付图片。
    seedream_audit_sample = SeedreamImageAsset(
        asset_id="verify-seedream",
        query="presentation visual",
        image_bytes=b"\\xff\\xd8\\xff",
        model="verify-model",
        prompt_digest="verify-digest",
    )
    assert seedream_audit_sample.audit_metadata()["watermark"] is False

    request = {
        "intent": "为一家智能制造团队制作一份面向客户的数字化升级方案 PPT。",
        "target_slide_count": 7,
        "allow_licensed_assets": True,
        "public_research_enabled": True,
    }
    planned = client.post("/api/agents/document_agent/presentation-studio/run", json=request)
    assert planned.status_code == 200, planned.text
    payload = planned.json()
    assert payload["mode"] == "mock"
    assert len(payload["slides"]) == 7
    assert payload["slides"][0]["role"] == "cover"
    assert payload["slides"][-1]["role"] == "sources"
    assert payload["brief"]["fact_check_notice"]
    assert payload["asset_plan"]["state"] == "planned"
    assert payload["asset_plan"]["provider"] == "pexels"
    assert payload["asset_plan"]["queries"]
    assert payload["asset_plan"]["slots"]
    assert len(payload["asset_plan"]["slots"]) == len(payload["asset_plan"]["queries"])
    assert all(item["slide_id"] in {slide["slide_id"] for slide in payload["slides"]} for item in payload["asset_plan"]["slots"])
    assert "没有联网" in payload["asset_plan"]["notice"]
    assert payload["research_plan"]["state"] == "planned"
    assert payload["research_plan"]["provider"] == "wikimedia"
    assert payload["research_plan"]["max_sources"] == 3
    assert "没有联网" in payload["research_plan"]["notice"]
    # 旧 V2.1 计划只有有序查询词；升级后仍要能在不改写历史快照的前提下恢复页面映射。
    legacy_payload = dict(payload)
    legacy_asset_plan = dict(payload["asset_plan"])
    legacy_asset_plan.pop("slots")
    legacy_payload["asset_plan"] = legacy_asset_plan
    legacy_plan = PresentationStudioPlanResponse.model_validate(legacy_payload)
    legacy_slots = _effective_asset_slots(legacy_plan)
    assert [slot.query for slot in legacy_slots] == payload["asset_plan"]["queries"]
    assert legacy_slots[0].slide_id == payload["slides"][0]["slide_id"]
    task_id = payload["task_id"]

    restored = client.get(f"/api/agents/document_agent/presentation-studio/{task_id}/result")
    assert restored.status_code == 200, restored.text
    assert restored.json()["status"] == "completed"
    assert restored.json()["result"]["plan_id"] == payload["plan_id"]

    history = client.get(f"/api/tasks/{task_id}")
    assert history.status_code == 200, history.text
    # /api/tasks/{task_id} 直接返回 WorkflowRun，不额外包一层 workflow_run。
    steps = history.json()["steps"]
    assert steps[-1]["step_id"] == "presentation_studio_plan"
    assert steps[-1]["output"]["external_assets_fetched"] is False
    assert steps[-1]["output"]["public_research_fetched"] is False
    assert steps[-1]["output"]["structured_data_fetched"] is False

    artifacts = client.get(f"/api/tasks/{task_id}/artifacts")
    assert artifacts.status_code == 200, artifacts.text
    assert artifacts.json()["artifacts"] == []
    assert not (_VERIFY_ROOT / "presentations").exists()

    missing_confirmation = client.post(
        f"/api/agents/document_agent/presentation-studio/{task_id}/export",
        json={
            "plan_id": payload["plan_id"],
            "filename": "未确认创作方案.pptx",
            "confirmed": False,
        },
    )
    assert missing_confirmation.status_code == 409, missing_confirmation.text

    exported = client.post(
        f"/api/agents/document_agent/presentation-studio/{task_id}/export",
        json={
            "plan_id": payload["plan_id"],
            "filename": "智能制造升级方案.pptx",
            "confirmed": True,
            "fetch_licensed_assets": True,
            "network_confirmed": True,
        },
    )
    assert exported.status_code == 200, exported.text
    exported_payload = exported.json()
    assert exported_payload["verification"]["passed"] is True
    assert exported_payload["verification"]["warnings"]
    output_path = _VERIFY_ROOT / "presentations" / "智能制造升级方案.pptx"
    assert output_path.exists() and output_path.stat().st_size > 0
    presentation = Presentation(output_path)
    assert len(presentation.slides) == len(payload["slides"])
    # 动效必须在最终 PPTX 的 slide XML 中真实存在。这里不依赖 XML 前缀字符串，而由
    # 专用回读器按 PresentationML namespace 核对转场、入场数量与每个 shape 目标。
    motion_inspection = inspect_native_presentation_motion(output_path)
    assert motion_inspection.transition_slide_count == len(payload["slides"])
    assert motion_inspection.entrance_effect_count > 0
    assert motion_inspection.invalid_target_count == 0
    # 未配置 Key 时必须成功降级，正文按受控布局语法而不是按页码硬塞同一种表格。
    assert exported_payload["verification"]["warnings"]
    content_layouts = [slide["layout"] for slide in payload["slides"] if slide["role"] == "content"]
    assert content_layouts == ["comparison", "process", "timeline"]

    persisted_artifacts = client.get(f"/api/tasks/{task_id}/artifacts")
    assert persisted_artifacts.status_code == 200, persisted_artifacts.text
    artifact = next(
        item
        for item in persisted_artifacts.json()["artifacts"]
        if item["artifact_id"] == exported_payload["artifact_id"]
    )
    assert artifact["metadata"]["presentation_mode"] == "studio_v2"
    assert artifact["metadata"]["external_assets_fetched"] is False
    assert artifact["metadata"]["asset_count"] == 0
    assert artifact["metadata"]["public_research_count"] == 0
    motion = artifact["metadata"]["native_presentation_motion"]
    assert motion["enabled"] is True
    assert motion["transition"] == "fade"
    assert motion["entrance"] == "fade_on_click"
    assert motion["transition_slide_count"] == len(payload["slides"])
    assert motion["entrance_effect_count"] == motion_inspection.entrance_effect_count

    # Qt 在导出前会先建立新的事件通道，避免已经结束的“计划阶段”流关闭后丢失联网、
    # 渲染与回读进度。导出 API 也要把这些阶段写入 SQLite，供历史页补拉。
    export_events = client.get(f"/api/tasks/{task_id}/logs")
    assert export_events.status_code == 200, export_events.text
    export_event_names = [event["event"] for event in export_events.json()["events"]]
    assert "presentation_export_queued" in export_event_names
    assert "presentation_export_started" in export_event_names
    assert "presentation_render_started" in export_event_names
    assert "presentation_render_verified" in export_event_names
    assert "presentation_export_completed" in export_event_names

    prepared_export = client.post(
        f"/api/agents/document_agent/presentation-studio/{task_id}/export/prepare"
    )
    assert prepared_export.status_code == 202, prepared_export.text
    assert prepared_export.json()["task_id"] == task_id
    prepared_events = client.get(f"/api/tasks/{task_id}/logs")
    assert prepared_events.status_code == 200, prepared_events.text
    assert any(event["event"] == "presentation_export_channel_ready" for event in prepared_events.json()["events"])

    # 已规划公开资料时，导出前仍必须给出本次联网确认；不允许仅靠计划开关绕过确认。
    unconfirmed_research = client.post(
        f"/api/agents/document_agent/presentation-studio/{task_id}/export",
        json={
            "plan_id": payload["plan_id"],
            "filename": "未确认公开资料.pptx",
            "confirmed": True,
            "fetch_public_research": True,
            "network_confirmed": False,
        },
    )
    assert unconfirmed_research.status_code == 409, unconfirmed_research.text

    # API 联调不访问互联网。把受控 Provider 换成固定结果，验证资料只进入来源页与 artifact，
    # 不会改变正文计划或伪装成模型生成的事实。
    original_research_resolver = delivery_module._resolve_public_research
    delivery_module._resolve_public_research = lambda **_: research_sample
    try:
        research_export = client.post(
            f"/api/agents/document_agent/presentation-studio/{task_id}/export",
            json={
                "plan_id": payload["plan_id"],
                "filename": "智能制造公开资料验证.pptx",
                "confirmed": True,
                "fetch_public_research": True,
                "network_confirmed": True,
            },
        )
    finally:
        delivery_module._resolve_public_research = original_research_resolver
    assert research_export.status_code == 200, research_export.text
    research_export_payload = research_export.json()
    assert research_export_payload["verification"]["passed"] is True
    research_presentation = Presentation(_VERIFY_ROOT / "presentations" / "智能制造公开资料验证.pptx")
    source_text = "\n".join(shape.text for shape in research_presentation.slides[-1].shapes if getattr(shape, "has_text_frame", False))
    assert research_sample.sources[0].title in source_text
    assert research_sample.sources[0].page_url in source_text
    research_artifact = next(
        item
        for item in client.get(f"/api/tasks/{task_id}/artifacts").json()["artifacts"]
        if item["artifact_id"] == research_export_payload["artifact_id"]
    )
    assert research_artifact["metadata"]["public_research_count"] == 1
    assert research_artifact["metadata"]["public_research_sources"][0]["scope"] == "public_reference_only"

    # 数据图表只会在计划已明确匹配、用户确认联网后执行。实际 HTTP 调用在此用固定结果替代，
    # 验证图表、来源页和 artifact 三处都带上同一个受控数据身份。
    data_plan_request = client.post(
        "/api/agents/document_agent/presentation-studio/run",
        json={
            "intent": "制作一份中国和美国 GDP 对比的客户演示 PPT。",
            "target_slide_count": 7,
            "visual_asset_provider": "none",
            "structured_data_enabled": True,
        },
    )
    assert data_plan_request.status_code == 200, data_plan_request.text
    data_plan_payload = data_plan_request.json()
    assert data_plan_payload["data_plan"]["state"] == "provider_planned"
    assert data_plan_payload["data_plan"]["chart_type"] == "comparison_bar"
    assert data_plan_payload["data_plan"]["country_codes"] == ["CHN", "USA"]
    assert data_plan_payload["data_plan"]["indicator_code"] == "NY.GDP.MKTP.CD"
    assert data_plan_payload["asset_plan"]["slots"] == []
    unconfirmed_data = client.post(
        f"/api/agents/document_agent/presentation-studio/{data_plan_payload['task_id']}/export",
        json={
            "plan_id": data_plan_payload["plan_id"],
            "filename": "未确认数据图表.pptx",
            "confirmed": True,
            "fetch_structured_data": True,
            "network_confirmed": False,
        },
    )
    assert unconfirmed_data.status_code == 409, unconfirmed_data.text
    original_data_resolver = delivery_module._resolve_structured_data
    delivery_module._resolve_structured_data = lambda **_: data_sample
    try:
        data_export = client.post(
            f"/api/agents/document_agent/presentation-studio/{data_plan_payload['task_id']}/export",
            json={
                "plan_id": data_plan_payload["plan_id"],
                "filename": "中国美国GDP数据图表验证.pptx",
                "confirmed": True,
                "fetch_structured_data": True,
                "network_confirmed": True,
            },
        )
    finally:
        delivery_module._resolve_structured_data = original_data_resolver
    assert data_export.status_code == 200, data_export.text
    data_presentation = Presentation(_VERIFY_ROOT / "presentations" / "中国美国GDP数据图表验证.pptx")
    data_source_text = "\n".join(
        shape.text for shape in data_presentation.slides[-1].shapes if getattr(shape, "has_text_frame", False)
    )
    assert data_sample.chart.title in "\n".join(
        shape.text for shape in data_presentation.slides[2].shapes if getattr(shape, "has_text_frame", False)
    )
    assert any(getattr(shape, "has_chart", False) for shape in data_presentation.slides[2].shapes)
    assert data_sample.chart.source_url in data_source_text
    data_artifact = next(
        item
        for item in client.get(f"/api/tasks/{data_plan_payload['task_id']}/artifacts").json()["artifacts"]
        if item["artifact_id"] == data_export.json()["artifact_id"]
    )
    assert data_artifact["metadata"]["structured_data_chart"]["indicator_code"] == "NY.GDP.MKTP.CD"

    # 通用研究同样必须穿过导出、来源页和 artifact；这里复用已验证的离线数据，确保业务层
    # 不会把 ResearchGateway 特有字段硬编码成 World Bank 字段。
    generic_delivery_plan = client.post(
        "/api/agents/document_agent/presentation-studio/run",
        json={
            "intent": (
                "帮我审查梅西与C罗生涯数据对比ppt，要有多种数据对比，要有至少三个表格，"
                "一个折线图，一个柱状图。"
            ),
            "target_slide_count": 9,
            "structured_data_enabled": True,
        },
    )
    assert generic_delivery_plan.status_code == 200, generic_delivery_plan.text
    generic_delivery_payload = generic_delivery_plan.json()
    generic_plan_snapshot = PresentationStudioPlanResponse.model_validate(generic_delivery_payload)
    generic_plan_snapshot = generic_plan_snapshot.model_copy(
        update={"data_plan": explicit_plan}
    )
    multi_view_resolution = ResearchGatewayResolution(
        chart=contract_charts[0],
        charts=contract_charts,
        warnings=(),
    )
    original_plan_loader = delivery_module._load_plan
    original_data_resolver = delivery_module._resolve_structured_data
    delivery_module._load_plan = lambda _: generic_plan_snapshot
    delivery_module._resolve_structured_data = lambda **_: multi_view_resolution
    try:
        generic_data_export = client.post(
            f"/api/agents/document_agent/presentation-studio/{generic_delivery_payload['task_id']}/export",
            json={
                "plan_id": generic_plan_snapshot.plan_id,
                "filename": "通用研究数据验证.pptx",
                "confirmed": True,
                "fetch_structured_data": True,
                "network_confirmed": True,
            },
        )
    finally:
        delivery_module._load_plan = original_plan_loader
        delivery_module._resolve_structured_data = original_data_resolver
    assert generic_data_export.status_code == 200, generic_data_export.text
    generic_presentation = Presentation(_VERIFY_ROOT / "presentations" / "通用研究数据验证.pptx")
    generic_source_text = "\n".join(
        shape.text for shape in generic_presentation.slides[-1].shapes if getattr(shape, "has_text_frame", False)
    )
    generic_chart_text = "\n".join(
        cell.text
        for shape in generic_presentation.slides[2].shapes
        if getattr(shape, "has_table", False)
        for row in shape.table.rows
        for cell in row.cells
    )
    all_slide_text = [
        "\n".join(
            [
                shape.text
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False)
            ]
            + [
                cell.text
                for shape in slide.shapes
                if getattr(shape, "has_table", False)
                for row in shape.table.rows
                for cell in row.cells
            ]
        )
        for slide in generic_presentation.slides
    ]
    assert contract_source.source_url in generic_source_text
    assert any(getattr(shape, "has_table", False) for shape in generic_presentation.slides[2].shapes)
    assert any(getattr(shape, "has_table", False) for shape in generic_presentation.slides[3].shapes)
    assert any(getattr(shape, "has_table", False) for shape in generic_presentation.slides[4].shapes)
    assert any(getattr(shape, "has_chart", False) for shape in generic_presentation.slides[5].shapes)
    assert any(getattr(shape, "has_chart", False) for shape in generic_presentation.slides[6].shapes)
    assert "101 项" in generic_chart_text and "111 项" in generic_chart_text
    assert contract_charts[1].title in all_slide_text[3]
    assert contract_charts[2].title in all_slide_text[4]
    assert contract_charts[3].title in all_slide_text[5]
    assert contract_charts[4].title in all_slide_text[6]
    assert "2021/22" in all_slide_text[4] and "2023/24" in all_slide_text[6]
    generic_artifact = next(
        item
        for item in client.get(f"/api/tasks/{generic_delivery_payload['task_id']}/artifacts").json()["artifacts"]
        if item["artifact_id"] == generic_data_export.json()["artifact_id"]
    )
    assert generic_artifact["metadata"]["structured_data_contract_complete"] is True
    assert generic_artifact["metadata"]["structured_data_contract_gap"] == ""
    assert generic_artifact["metadata"]["structured_data_provider"] == "research_gateway"
    assert generic_artifact["metadata"]["structured_data_chart"]["search_provider"] == "fixture"
    assert len(generic_artifact["metadata"]["structured_data_charts"]) == 5

    # Seedream 是与聊天模型分离的生成式视觉 Provider。离线验证不配置 Key，必须保留计划、
    # 明确产生降级 warning、仍能导出并绝不向外发请求。
    seedream_plan = client.post(
        "/api/agents/document_agent/presentation-studio/run",
        json={
            "intent": "为一家智能制造团队制作一份面向客户的数字化升级方案 PPT。",
            "target_slide_count": 7,
            "visual_asset_provider": "seedream",
        },
    )
    assert seedream_plan.status_code == 200, seedream_plan.text
    seedream_payload = seedream_plan.json()
    assert seedream_payload["asset_plan"]["state"] == "planned"
    assert seedream_payload["asset_plan"]["provider"] == "seedream"
    assert "没有调用图像模型" in seedream_payload["asset_plan"]["notice"]
    seedream_export = client.post(
        f"/api/agents/document_agent/presentation-studio/{seedream_payload['task_id']}/export",
        json={
            "plan_id": seedream_payload["plan_id"],
            "filename": "即梦离线降级验证.pptx",
            "confirmed": True,
            "fetch_external_assets": True,
            "network_confirmed": True,
        },
    )
    assert seedream_export.status_code == 200, seedream_export.text
    assert any("未配置 Seedream" in item for item in seedream_export.json()["verification"]["warnings"])
    seedream_artifacts = client.get(f"/api/tasks/{seedream_payload['task_id']}/artifacts")
    assert seedream_artifacts.status_code == 200, seedream_artifacts.text
    seedream_artifact = next(
        item for item in seedream_artifacts.json()["artifacts"]
        if item["artifact_id"] == seedream_export.json()["artifact_id"]
    )
    assert seedream_artifact["metadata"]["external_assets_fetched"] is False
    assert seedream_artifact["metadata"]["asset_count"] == 0

    duplicate = client.post(
        f"/api/agents/document_agent/presentation-studio/{task_id}/export",
        json={
            "plan_id": payload["plan_id"],
            "filename": "智能制造升级方案.pptx",
            "confirmed": True,
            "fetch_licensed_assets": True,
            "network_confirmed": True,
        },
    )
    assert duplicate.status_code == 409, duplicate.text

    invalid = client.post(
        "/api/agents/document_agent/presentation-studio/run",
        json={"intent": "太短", "target_slide_count": 7},
    )
    assert invalid.status_code == 422, invalid.text

    # 四套主题都必须实际完成渲染，不允许只在计划 JSON 中换一个颜色名称。
    for theme in ("executive_blue", "technology_emerald", "narrative_warm", "impact_contrast"):
        themed = client.post(
            "/api/agents/document_agent/presentation-studio/run",
            json={
                "intent": "为客户说明一个可确认的服务升级方案。",
                "target_slide_count": 5,
                "theme_preference": theme,
                "allow_licensed_assets": False,
            },
        )
        assert themed.status_code == 200, themed.text
        themed_payload = themed.json()
        assert themed_payload["brief"]["theme"] == theme
        assert themed_payload["slides"][2]["layout"] == "comparison"
        themed_export = client.post(
            f"/api/agents/document_agent/presentation-studio/{themed_payload['task_id']}/export",
            json={
                "plan_id": themed_payload["plan_id"],
                "filename": f"主题-{theme}.pptx",
                "confirmed": True,
            },
        )
        assert themed_export.status_code == 200, themed_export.text
        rendered = Presentation(_VERIFY_ROOT / "presentations" / f"主题-{theme}.pptx")
        assert len(rendered.slides) == len(themed_payload["slides"])

    started = client.post("/api/agents/document_agent/presentation-studio/start", json=request)
    assert started.status_code == 202, started.text
    background_task_id = started.json()["task_id"]
    result = None
    for _ in range(30):
        candidate = client.get(
            f"/api/agents/document_agent/presentation-studio/{background_task_id}/result"
        )
        assert candidate.status_code == 200, candidate.text
        result = candidate.json()
        if result["status"] == "completed":
            break
        time.sleep(0.03)
    assert result is not None and result["status"] == "completed", result
    assert result["result"]["task_id"] == background_task_id
    print("Presentation studio verification passed.")


if __name__ == "__main__":
    main()
